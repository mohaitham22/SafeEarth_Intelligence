"""
build_pptx.py — generate an editable PowerPoint from the SafeEarth backend pitch deck.

Run:  py -3.12 docs/build_pptx.py
Out:  docs/SafeEarth_Backend_Presentation.pptx

Native python-pptx (editable text boxes + tables), no browser needed.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)   # deep green (titles/headers)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)   # accent green
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # accent amber
SLATE  = RGBColor(0x33, 0x41, 0x55)   # body text
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # light panel
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)

W, H = Inches(13.333), Inches(7.5)     # 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _set(frame, size, color, bold=False, align=PP_ALIGN.LEFT):
    for p in frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def title_bar(s, text, kicker=None):
    band(s, DEEP, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(28); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"
    tf.word_wrap = True


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    tf.word_wrap = True
    first = True
    for it in items:
        # tuple (text, level, color, bold)
        if isinstance(it, tuple):
            text, level, color, bold = (it + (0, SLATE, False))[:4]
        else:
            text, level, color, bold = it, 0, SLATE, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0))
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DEEP
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def code(s, lines, left, top, width, height, fsize=13):
    band(s, RGBColor(0x0F, 0x17, 0x2A), left, top, width, height)
    tf = box(s, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = ln
        r.font.name = "Consolas"; r.font.size = Pt(fsize)
        r.font.color.rgb = RGBColor(0x7E, 0xE7, 0xB0)
    return tf


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.7), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.25), Inches(11.3), Inches(1.4))
for i, (txt, sz, col, bold) in enumerate([
    ("Turning an ML Model Into a Real Product", 26, GREEN, True),
    ("Backend & API Architecture", 20, RGBColor(0xCB,0xD5,0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "FastAPI · PostgreSQL · XGBoost · RAG (Groq)"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94,0xA3,0xB8); r.font.name = "Calibri"
notes(s, "Opening line: I trained a disaster-prediction model. This talk is about the "
         "90% of engineering that happens AFTER the model — turning a notebook into "
         "something thousands of people can actually use.")

# ── SLIDE 2 — Problem ─────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Problem")
bullets(s, [
    ("Natural disasters are predictable from history — but predictions are useless locked in a notebook.", -1, DEEP, True),
    ("16,126 historical disaster events (EM-DAT, 1900–2021)", 0, SLATE, False),
    ("A model that estimates type, probability, severity, and human / economic impact for any location", 0, SLATE, False),
    ("But: a model on my laptop helps exactly one person — me.", 0, AMBER, True),
], top=Inches(1.6), size=20, gap=16)
band(s, LIGHT, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.4))
tf = box(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Goal: "
r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
r = p.add_run(); r.text = "make ML predictions usable by anyone — a website, a phone, a scheduled alert — safely, reliably, 24/7."
r.font.size = Pt(19); r.font.color.rgb = SLATE; r.font.name = "Calibri"

# ── SLIDE 3 — What is SafeEarth ───────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What Is SafeEarth?")
bullets(s, [
    ("Predicts natural disasters for any region (XGBoost + SHAP explanations)", 0, SLATE, False),
    ("Estimates impact: deaths, injuries, affected, economic damage", 0, SLATE, False),
    ("Recommends safety actions via a RAG + LLM pipeline on official guidelines", 0, SLATE, False),
    ("Alerts subscribers by email (automated, scheduled)", 0, SLATE, False),
    ("Forecasts day-by-day risk for the next 30 days", 0, SLATE, False),
    ("Visualizes global trends & analytics", 0, SLATE, False),
], top=Inches(1.6), size=20, gap=14)
band(s, GREEN, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.8))
tf = box(s, Inches(0.9), Inches(6.22), Inches(11.5), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Live:  safeearth.tech (frontend)   ·   api.safeearth.tech (backend)"
r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

# ── SLIDE 4 — Notebook-to-Product Gap ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Core Challenge: the Notebook-to-Product Gap")
table(s, [
    ["In a notebook", "Behind an API"],
    ["model.predict(X) when I like", "A stranger triggers it over the internet"],
    ["I trust my own inputs", "Inputs must be validated"],
    ["Load .pkl once, in my session", "Load once at startup, reuse forever"],
    ["I eyeball the output", "Output is a fixed JSON contract"],
    ["An error crashes my cell (fine)", "An error must be a clean HTTP response"],
    ["One call at a time", "Many concurrent users"],
], Inches(0.6), Inches(1.5), Inches(12.1),
   col_widths=[Inches(5.6), Inches(6.5)], fsize=16)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "The backend is everything required to close this gap."
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
notes(s, "This slide is the thesis of the whole talk. Land it slowly. Professors care "
         "that you understand WHY a backend exists, not just that you wrote one.")

# ── SLIDE 5 — Restaurant analogy ──────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Restaurant Analogy")
tf = box(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.5))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "A backend is a restaurant, not a home kitchen:"
r.font.size = Pt(18); r.font.color.rgb = SLATE; r.font.name = "Calibri"
table(s, [
    ["Restaurant", "Software", "In the repo"],
    ["Dining room & menu", "Frontend", "frontend/ (Next.js)"],
    ["The waiter", "API layer", "backend/routers/ (FastAPI)"],
    ["The kitchen", "Business logic", "backend/services/"],
    ["The pantry", "Database", "PostgreSQL via backend/models/"],
    ["The recipe binder", "ML models (loaded once)", "backend/saved_models/*.pkl"],
], Inches(0.6), Inches(2.0), Inches(12.1),
   col_widths=[Inches(3.6), Inches(3.8), Inches(4.7)], fsize=15)
tf = box(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "A diner never walks into the kitchen — they place an order and get a dish."
r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = GREY; r.font.name = "Calibri"

# ── SLIDE 6 — Architecture ────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "System Architecture")
# simple boxed flow rendered as shapes
def node(s, l, t, w, h, text, fill, fg=WHITE, fs=13):
    sh = band(s, fill, l, t, w, h)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = fg; r.font.name = "Calibri"
    return sh
node(s, Inches(0.6), Inches(2.6), Inches(2.2), Inches(0.9), "Browser /\nPhone\n(Next.js)", SLATE)
node(s, Inches(0.6), Inches(4.0), Inches(2.2), Inches(0.9), "n8n\n(scheduled\nrobot)", AMBER, DEEP)
# FastAPI cluster
band(s, LIGHT, Inches(3.3), Inches(1.6), Inches(5.0), Inches(4.6))
tf = box(s, Inches(3.4), Inches(1.65), Inches(4.8), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "FastAPI backend"; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name="Calibri"
node(s, Inches(3.5), Inches(2.15), Inches(4.6), Inches(0.7), "Routers  →  Auth & Permissions", GREEN)
node(s, Inches(3.5), Inches(3.05), Inches(4.6), Inches(0.7), "Services (business logic)", DEEP)
node(s, Inches(3.5), Inches(3.95), Inches(2.2), Inches(0.7), "ML models\n(loaded once)", SLATE, fs=11)
node(s, Inches(5.9), Inches(3.95), Inches(2.2), Inches(0.7), "RAG\nrecommender", SLATE, fs=11)
node(s, Inches(3.5), Inches(5.15), Inches(4.6), Inches(0.8), "validate → auth → run model → save → reply", GREEN, fs=12)
# right side externals
node(s, Inches(8.8), Inches(2.15), Inches(3.9), Inches(0.7), "PostgreSQL / Neon (database)", DEEP)
node(s, Inches(8.8), Inches(3.05), Inches(3.9), Inches(0.7), "Groq LLM (recommendations)", AMBER, DEEP)
node(s, Inches(8.8), Inches(3.95), Inches(3.9), Inches(0.7), "Resend / Gmail (email alerts)", SLATE)
node(s, Inches(8.8), Inches(4.85), Inches(3.9), Inches(0.7), "HuggingFace (model files)", SLATE)
tf = box(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Every request: validate → authenticate → authorize → run model → save → augment → reply"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = SLATE; r.font.name = "Calibri"

# ── SLIDE 7 — Tech stack ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Tech Stack at a Glance")
table(s, [
    ["Layer", "Tool", "Why"],
    ["API framework", "FastAPI", "Type hints → validation + docs for free"],
    ["Server", "Uvicorn (ASGI)", "Async-native"],
    ["Validation", "Pydantic v2", "Enforced request / response contracts"],
    ["Database", "PostgreSQL (Neon)", "Relational integrity, JSONB, free tier"],
    ["ORM", "SQLAlchemy 2.0 async", "Python objects ↔ rows, no raw SQL"],
    ["Migrations", "Alembic", "Version-controlled schema"],
    ["Auth", "JWT + bcrypt", "Stateless login, hashed passwords"],
    ["ML", "XGBoost + CatBoost + SHAP", "Tabular accuracy + explainability"],
    ["RAG", "Groq + PyMuPDF", "Grounded LLM recommendations"],
    ["Hosting", "Render · Vercel · Neon", "Zero-cost, infra-as-code"],
], Inches(0.6), Inches(1.45), Inches(12.1),
   col_widths=[Inches(2.7), Inches(3.6), Inches(5.8)], fsize=13)

# ── Component slide helper ────────────────────────────────────────────────────
def component(num, title, what, body_items, code_lines=None, chose=None, note=None):
    s = slide(); bg(s, WHITE)
    title_bar(s, title, kicker=f"Component {num}")
    top = Inches(1.45)
    if what:
        tf = box(s, Inches(0.6), top, Inches(12.1), Inches(0.7))
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = what
        r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
        top = Inches(2.2)
    if code_lines:
        code(s, code_lines, Inches(0.6), top, Inches(12.1), Inches(0.42 * len(code_lines) + 0.3), fsize=13)
        top = Inches(top.inches + 0.42 * len(code_lines) + 0.55)
    bullets(s, body_items, top=top, size=17, gap=9)
    if chose:
        y = Inches(6.55)
        band(s, LIGHT, Inches(0.6), y, Inches(12.1), Inches(0.7))
        tf = box(s, Inches(0.8), Inches(y.inches + 0.12), Inches(11.7), Inches(0.5))
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = "Chose over: "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
        r = p.add_run(); r.text = chose
        r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    if note:
        notes(s, note)
    return s

# ── SLIDE 8 — FastAPI ─────────────────────────────────────────────────────────
component(
    1, "The API Layer (FastAPI)",
    "Turns ordinary typed Python functions into internet-callable endpoints.",
    [
        ("Validation is free: latitude outside [-90, 90] → auto 422, before my code runs", 0, SLATE, False),
        ("Docs are free: interactive Swagger UI at /docs", 0, SLATE, False),
        ("~30 endpoints across 10 routers (auth, predictions, regions, alerts, …)", 0, SLATE, False),
    ],
    code_lines=[
        '@router.post("/predict", response_model=PredictionResponse)',
        '@limiter.limit("60/minute")',
        'async def predict(body: PredictRequest, user = Depends(require_subscriber)):',
        '    return await predictor_service.run_prediction_for_request(...)',
    ],
    chose="Flask (needs bolt-ons for validation + docs) and Django (too heavy for a JSON-only ML API).",
)

# ── SLIDE 9 — Async ───────────────────────────────────────────────────────────
component(
    2, "Async / Concurrency",
    "Each prediction waits on the database, the LLM, and email servers.",
    [
        ("Analogy: a good waiter serves tables 2 and 3 while table 1's food cooks", 0, SLATE, False),
        ("One cheap server handles many concurrent users by overlapping the waits", 0, SLATE, False),
        ("Critical on a 512 MB free tier — can't afford many worker processes", 0, AMBER, True),
    ],
    code_lines=['result = await db.execute(select(User).where(User.email == email))'],
    chose="Celery / a full task queue — overkill; built-in BackgroundTasks + one asyncio loop suffice.",
)

# ── SLIDE 10 — Persistence ────────────────────────────────────────────────────
component(
    3, "Persistence — Database + ORM + Migrations",
    "The data layer = a DataFrame that survives restarts and serves many users at once.",
    [
        ("PostgreSQL (Neon): 9 tables — users, predictions, payments… with FK integrity", 0, SLATE, False),
        ("SQLAlchemy ORM: rows as Python objects → no SQL injection, type-safe", 0, SLATE, False),
        ("Alembic migrations = git commits for the schema (reproducible everywhere)", 0, SLATE, False),
        ("Real example: ads table + unsubscribe token added AFTER launch — safely", 0, GREY, False),
    ],
    code_lines=[
        'row = Prediction(user_id=user_id, disaster_type=..., risk_score=...)',
        'db.add(row); await db.commit()',
    ],
    chose="SQLite (weak concurrency), MongoDB (data is highly relational), raw SQL (injection-prone).",
)

# ── SLIDE 11 — ML serving ─────────────────────────────────────────────────────
component(
    4, "ML Model Serving — the load-once rule",
    "Cardinal rule: load the .pkl ONCE at startup, never per request.",
    [
        ("v4.2 XGBoost + CatBoost soft ensemble + per-type impact regressors", 0, SLATE, False),
        ("SHAP top-3 explanation returned with every prediction", 0, SLATE, False),
        ("Memory trick: rebuild SHAP from in-memory model → saves ~300 MB", 0, AMBER, True),
        ("Models too big for git → downloaded from HuggingFace at startup", 0, SLATE, False),
    ],
    code_lines=[
        'async def lifespan(app):',
        '    ml_predictor.load_models(...)   # once, into module-level vars',
        '    yield                           # serve requests, models in RAM',
    ],
    chose="A dedicated model server (Triton / SageMaker) — pricier & more complex; in-process is fast enough.",
    note="Model quality if asked: Macro F1 ~ 0.71, Weighted F1 ~ 0.76 on a held-out "
         "1970–2021 test set. The point of THIS slide is serving, not training.",
)

# ── SLIDE 12 — RAG ────────────────────────────────────────────────────────────
component(
    5, "The RAG Pipeline (+ an honest pivot)",
    "RAG = give the LLM the relevant reference text, then ask it to answer FROM that text.",
    [
        ("Generates 6 grounded safety recommendations per prediction", 0, SLATE, False),
        ("Original design: ChromaDB + sentence-transformer embeddings (vector RAG)", 0, SLATE, False),
        ("Reality: that stack pulls in PyTorch (~2 GB) → out-of-memory on the free tier", 0, AMBER, True),
        ("Pivot: chapter-based retrieval → send the right PDF chapter to Groq", 0, SLATE, False),
        ("Resilience: Groq down? → fallback to a seeded DB table. Predictions never break.", 0, GREEN, True),
    ],
    chose="Full vector RAG (OOM on host) and blind LLM prompting (would hallucinate unsafe advice).",
    note="This is the slide professors will probe. Lesson: match the solution to real "
         "constraints, and degrade gracefully. Own the trade-off.",
)

# ── SLIDE 13 — Auth ───────────────────────────────────────────────────────────
component(
    6, "Authentication & Authorization",
    "Authn = who are you (JWT). Authz = what may you do (central permissions).",
    [
        ("JWT: signed, tamper-proof token (access 30 min, refresh 7 days)", 0, SLATE, False),
        ("Passwords hashed with bcrypt — never stored in plaintext", 0, SLATE, False),
        ("Golden rule: the server is the security boundary; frontend check is UX-only", 0, AMBER, True),
        ("Dual-credential machine guard: shared secret OR admin JWT (constant-time compare)", 0, SLATE, False),
    ],
    code_lines=[
        'ROLE_RANK = {"guest": 0, "subscriber": 1, "premium": 2, "admin": 3}',
        'async def predict(..., user = Depends(require_subscriber)):  # 401/403 here',
    ],
    chose="Server-side sessions — awkward across split frontend/backend; managed auth — adds cost & lock-in.",
)

# ── SLIDE 14 — Background jobs ─────────────────────────────────────────────────
component(
    7, "Background Jobs & Automation",
    "Work that shouldn't make the user wait, or that runs on a timer.",
    [
        ("BackgroundTasks: send verification email AFTER replying → instant registration", 0, SLATE, False),
        ("asyncio loop: daily premium-expiry checker, started in lifespan", 0, SLATE, False),
        ("n8n (external cron): weekly alert dispatch + monthly Egypt CSV export", 0, SLATE, False),
        ("Rule: n8n is only a timer that calls the API — FastAPI does all the real work", 0, AMBER, True),
    ],
    code_lines=[
        'background_tasks.add_task(email_service.send_verification_email, email, token)',
    ],
    chose="Celery + Redis broker — real infrastructure & cost; unnecessary at this scale.",
)

# ── SLIDE 15 — Service abstraction ────────────────────────────────────────────
component(
    8, "Service Abstraction (design pattern)",
    "Code to an interface — like scikit-learn's .fit()/.predict(): any estimator is swappable.",
    [
        ("Full checkout → webhook → 'you're now premium' flow works with zero real money", 0, SLATE, False),
        ("Real Stripe / Paymob later = one new file, no route changes", 0, GREEN, True),
    ],
    code_lines=[
        'class PaymentService(ABC):          # the contract',
        '    async def create_checkout_session(...): ...',
        'class MockPaymentService(PaymentService): ...   # v1, no real money',
        'def get_payment_service():          # factory — the only place that picks one',
    ],
    chose="Calling Stripe's SDK directly in routes — hard-wires a vendor & breaks testability.",
)

# ── SLIDE 16 — Deployment ─────────────────────────────────────────────────────
component(
    9, "Deployment & Infrastructure",
    "Local cooking vs opening the restaurant to the public.",
    [
        ("Backend → Render, Frontend → Vercel, Database → Neon — all free tier", 0, SLATE, False),
        ("Infrastructure-as-code (render.yaml) — reproducible, reviewable deploys", 0, SLATE, False),
        ("Secrets in environment variables, never in git", 0, SLATE, False),
        ("Build script preps the kitchen: install → stats → RAG chapters → migrate → seed", 0, SLATE, False),
        ("Trade-off: free tier sleeps when idle → UptimeRobot ping every 14 min", 0, AMBER, True),
    ],
    chose="A raw VM (manual ops), Kubernetes (over-engineered), serverless (bad at resident ML models).",
)

# ── SLIDE 17 — End-to-end request (money slide) ───────────────────────────────
s = slide(); bg(s, WHITE)
title_bar(s, "One Request, End to End", kicker="The Money Slide")
tf = box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "POST /api/v1/predictions/predict"
r.font.name = "Consolas"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DEEP
steps = [
    ("1", "Request arrives", "Router predict() in routers/predictions.py"),
    ("2", "Rate limit", "@limiter.limit('60/minute') — 429 if exceeded"),
    ("3", "Validate body", "Pydantic PredictRequest → 422 if bad, before my code runs"),
    ("4", "Authenticate + authorize", "require_subscriber: decode JWT, check role → 401/403"),
    ("5", "Into the service layer", "predictor_service.run_prediction_for_request()"),
    ("6", "Run the model", "predictor.predict() — model already in RAM + SHAP top-3"),
    ("7", "Persist to database", "db.add(Prediction) → await db.commit()"),
    ("8", "Augment with RAG", "6 recommendations via Groq (DB fallback if down)"),
    ("9", "Serialize & reply", "200 OK + PredictionResponse JSON"),
    ("10", "(If Critical) alerts", "fan out via BackgroundTasks — AFTER the reply"),
]
y = 1.95
for num, head, detail in steps:
    band(s, GREEN, Inches(0.6), Inches(y), Inches(0.5), Inches(0.42))
    tf = box(s, Inches(0.6), Inches(y - 0.02), Inches(0.5), Inches(0.46)); p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name="Calibri"
    tf = box(s, Inches(1.25), Inches(y - 0.04), Inches(11.4), Inches(0.5)); p = tf.paragraphs[0]
    r = p.add_run(); r.text = head + "  —  "
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
    r = p.add_run(); r.text = detail
    r.font.size = Pt(13); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    y += 0.5
notes(s, "Walk this slowly, step by step, naming the real file at each hop. This is "
         "where it all clicks. If short on time, cut a component slide, NOT this one.")

# ── SLIDE 18 — Trade-offs ─────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Engineering Decisions & Trade-offs")
table(s, [
    ["Decision", "Chosen", "Over", "Why"],
    ["API framework", "FastAPI", "Flask / Django", "Free validation + docs, async"],
    ["Concurrency", "async, 1 process", "Celery / workers", "Fits 512 MB free tier"],
    ["RAG retrieval", "Chapter + Groq", "ChromaDB vectors", "PyTorch OOM'd the host"],
    ["Auth", "Stateless JWT", "Server sessions", "Scales across split apps"],
    ["Model serving", "In-process", "Model server", "Cheaper, simpler, fast enough"],
    ["Payments", "Mock behind ABC", "Real Stripe now", "v1 needs no money; easy swap"],
], Inches(0.5), Inches(1.45), Inches(12.3),
   col_widths=[Inches(2.3), Inches(2.6), Inches(2.7), Inches(4.7)], fsize=13)
tf = box(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Every choice optimized for a single-developer, zero-cost, memory-constrained system."
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
notes(s, "This slide shows engineering maturity: not 'I used the trendiest tool' but "
         "'I matched tools to constraints, and I know what I traded away.'")

# ── SLIDE 19 — Challenges ─────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Challenges I Solved")
bullets(s, [
    ("Memory: SHAP pkl (109 MB → ~300 MB in RAM) → rebuilt from in-memory model", 0, SLATE, False),
    ("Memory: vector-RAG stack (~2 GB) → chapter-based Groq pipeline", 0, SLATE, False),
    ("Cloud DB quirks: Neon SSL params rejected by async driver → stripped + ssl=True", 0, SLATE, False),
    ("Security discipline: verify webhook signature BEFORE any DB write", 0, SLATE, False),
    ("Resilience: a RAG failure can never 500 a prediction (wrapped + DB fallback)", 0, SLATE, False),
    ("Big models + git: download .pkl from HuggingFace at startup", 0, SLATE, False),
], top=Inches(1.7), size=19, gap=16)

# ── SLIDE 20 — Honest engineering ─────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Honest Engineering — Portfolio vs Production")
tf = box(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Production-grade here:"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
bullets(s, [
    ("Centralized permissions · migrations · graceful degradation · infra-as-code · ~190 automated tests", 0, SLATE, False),
], top=Inches(1.95), size=17, gap=6)
tf = box(s, Inches(0.6), Inches(2.7), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Conscious shortcuts (with an upgrade path):"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = AMBER; r.font.name = "Calibri"
bullets(s, [
    ("JWT can't be instantly revoked → short 30-min expiry mitigates", 0, SLATE, False),
    ("BackgroundTasks lost on restart → fine at this scale; Celery later", 0, SLATE, False),
    ("In-process model serving → split out if traffic grows", 0, SLATE, False),
    ("Mock payments → swap-in real provider via the existing interface", 0, SLATE, False),
], top=Inches(3.25), size=17, gap=10)
band(s, LIGHT, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.8))
tf = box(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Knowing what you traded away is the engineering — not a weakness."
r.font.size = Pt(17); r.font.bold = True; r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"

# ── SLIDE 21 — Results ────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Results")
bullets(s, [
    ("Live in production: safeearth.tech + api.safeearth.tech", 0, GREEN, True),
    ("9-table relational schema, fully migrated", 0, SLATE, False),
    ("~30 endpoints, all validated & role-gated", 0, SLATE, False),
    ("ML serving within a 2–5 s latency budget on a 512 MB host", 0, SLATE, False),
    ("RAG + LLM recommendations with DB fallback", 0, SLATE, False),
    ("Automated alerts via scheduled n8n → API", 0, SLATE, False),
    ("~190 automated backend tests (pytest)", 0, SLATE, False),
], top=Inches(1.7), size=19, gap=14)

# ── SLIDE 22 — Future work ────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Future Work")
bullets(s, [
    ("Replace mock payments with real Stripe / Paymob (1-file swap, interface in place)", 0, SLATE, False),
    ("Re-introduce vector RAG on a larger-memory tier for finer retrieval", 0, SLATE, False),
    ("Add a task queue (Celery) for retries + at-least-once email delivery", 0, SLATE, False),
    ("Token revocation / refresh-token blacklist on logout", 0, SLATE, False),
    ("Computer-vision disaster classification (deferred v2)", 0, SLATE, False),
], top=Inches(1.8), size=20, gap=18)

# ── SLIDE 23 — Takeaways ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("A backend is model.predict() made safe, repeatable, and always-on for strangers.", 0, DEEP, True),
    ("Validation, auth, persistence, and serving close the notebook-to-product gap.", 0, SLATE, False),
    ("Good engineering = matching tools to real constraints — and degrading gracefully.", 0, SLATE, False),
    ("Every trade-off was deliberate, documented, and reversible.", 0, SLATE, False),
], top=Inches(2.0), size=21, gap=22)

# ── SLIDE 24 — Thank you ──────────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.0), W, Inches(0.08))
tf = box(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(4.6), Inches(11.3), Inches(1.4)); p = tf.paragraphs[0]
for i, txt in enumerate([
    "SafeEarth Intelligence — FastAPI · PostgreSQL · XGBoost · RAG (Groq)",
    "Full technical deep-dive: docs/BACKEND_EXPLAINED.md",
    "safeearth.tech",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94,0xA3,0xB8); r.font.name = "Calibri"
notes(s, "Likely Qs: Why not microservices? (one small app, coupling = simplicity). "
         "Model drift? (retrain scripts + versioned bundles, HF hosting). JWT secure? "
         "(signed HS256, short expiry; revocation is the known gap). Testing the LLM "
         "path? (mock Groq, assert 6 items + fallback). Median vs mean? (disaster data "
         "is right-skewed: flood mean 1,735 vs median 16).")

out = Path(__file__).resolve().parent / "SafeEarth_Backend_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
