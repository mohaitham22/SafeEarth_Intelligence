"""
build_frontend_pptx.py — generate an editable PowerPoint from the SafeEarth FRONTEND pitch deck.

Run:  py -3.12 docs/build_frontend_pptx.py
Out:  docs/SafeEarth_Frontend_Presentation.pptx

Native python-pptx (editable text boxes + tables), no browser needed.
Companion to build_pptx.py (the backend deck) — same palette + helpers, plus a
blue accent to signal "frontend chapter". Content sourced from FRONTEND_EXPLAINED.md
and FRONTEND_PITCH_DECK.md.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette (matched to the backend deck, + BLUE for frontend flavor) ───────────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)   # deep green (brand — title bars)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)   # accent green
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # accent amber (honest trade-offs)
BLUE   = RGBColor(0x25, 0x63, 0xEB)   # accent blue (frontend / DS-translation)
SLATE  = RGBColor(0x33, 0x41, 0x55)   # body text
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # light panel
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)

W, H = Inches(13.333), Inches(7.5)     # 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def title_bar(s, text, kicker=None):
    band(s, DEEP, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(28); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    tf.word_wrap = True
    first = True
    for it in items:
        # tuple (text, level, color, bold)
        if isinstance(it, tuple):
            text, level, color, bold = (it + (0, SLATE, False))[:4]
        else:
            text, level, color, bold = it, 0, SLATE, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0))
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DEEP
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def code(s, lines, left, top, width, height, fsize=13):
    band(s, RGBColor(0x0F, 0x17, 0x2A), left, top, width, height)
    tf = box(s, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = ln
        r.font.name = "Consolas"; r.font.size = Pt(fsize)
        r.font.color.rgb = RGBColor(0x7E, 0xE7, 0xB0)
    return tf


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def node(s, l, t, w, h, text, fill, fg=WHITE, fs=13):
    sh = band(s, fill, l, t, w, h)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = fg; r.font.name = "Calibri"
    return sh


# ── Component slide helper (mirrors backend deck) ──────────────────────────────
def component(num, title, what, body_items, code_lines=None, chose=None, note=None):
    s = slide(); bg(s, WHITE)
    title_bar(s, title, kicker=f"Component {num}")
    top = Inches(1.45)
    if what:
        tf = box(s, Inches(0.6), top, Inches(12.1), Inches(0.7))
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = what
        r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
        top = Inches(2.2)
    if code_lines:
        code(s, code_lines, Inches(0.6), top, Inches(12.1), Inches(0.42 * len(code_lines) + 0.3), fsize=13)
        top = Inches(top.inches + 0.42 * len(code_lines) + 0.55)
    bullets(s, body_items, top=top, size=17, gap=9)
    if chose:
        y = Inches(6.55)
        band(s, LIGHT, Inches(0.6), y, Inches(12.1), Inches(0.7))
        tf = box(s, Inches(0.8), Inches(y.inches + 0.12), Inches(11.7), Inches(0.5))
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = "Chose over: "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = "Calibri"
        r = p.add_run(); r.text = chose
        r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    if note:
        notes(s, note)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, BLUE, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.7), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(48); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.25), Inches(11.3), Inches(1.4))
for i, (txt, sz, col, bold) in enumerate([
    ("Turning JSON Into a Product People Can Use", 26, BLUE, True),
    ("The Frontend — How the Model Reaches Real People", 20, RGBColor(0xCB, 0xD5, 0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "Next.js 14 · React · TypeScript · Tailwind · Leaflet · Recharts"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "The backend gives me a model behind an API that returns JSON. This talk is the "
         "other half: how a flood-zone resident who will NEVER open a notebook actually "
         "uses that model. The frontend is the difference between a chart in a notebook "
         "and an app anyone can click.")

# ── SLIDE 2 — The Problem ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Problem")
bullets(s, [
    ("A trained model + API returns JSON — and nobody outside the team reads JSON.", -1, DEEP, True),
    ('{ "disaster_type": "Flood", "probability_score": 0.44, ... }', 0, SLATE, False),
    ("A prediction locked behind a curl command helps exactly one person.", 0, SLATE, False),
    ("A notebook chart helps me; a deployed app helps everyone.", 0, AMBER, True),
], top=Inches(1.6), size=20, gap=16)
band(s, LIGHT, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.4))
tf = box(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Goal: "
r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = "Calibri"
r = p.add_run(); r.text = ("take JSON from the backend and paint it as a map, a severity badge, a chart, "
                           "a PDF — and let a stranger ask for more without seeing a single brace.")
r.font.size = Pt(19); r.font.color.rgb = SLATE; r.font.name = "Calibri"
notes(s, "The frontend's entire job is translation: machine output to something a human "
         "acts on. This is a product problem, not a decoration problem.")

# ── SLIDE 3 — What the Frontend Delivers ───────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What the Frontend Delivers")
bullets(s, [
    ("Risk map — click anywhere on Earth → a live disaster prediction (/map)", 0, SLATE, False),
    ("Global analytics — trends, continent comparison, insurance gap, time series", 0, SLATE, False),
    ("Prediction dashboard — location → severity, probability, impact + SHAP", 0, SLATE, False),
    ("30-day forecast — day-by-day risk calendar + trend line", 0, SLATE, False),
    ("Accounts & alerts — register, verify, log in, subscribe to regions", 0, SLATE, False),
    ("Premium — checkout flow + downloadable PDF report", 0, SLATE, False),
], top=Inches(1.6), size=20, gap=13)
band(s, BLUE, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.8))
tf = box(s, Inches(0.9), Inches(6.22), Inches(11.5), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "17 routes, one Next.js app. Each feature = fetch JSON from the backend, render it for a human."
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

# ── SLIDE 4 — The Mental Shift ─────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Mental Shift: Notebook → Web App")
table(s, [
    ["In a notebook (you)", "In a web app (the user)"],
    ["You run cells top to bottom", "The USER decides what runs, by clicking"],
    ["Output prints once", "UI must re-render whenever data changes"],
    ["You trust your own inputs", "Every input must be validated"],
    ["One person — you", "Many strangers, concurrently, 24/7"],
    ["plt.show()", "A live, interactive, navigable screen"],
], Inches(0.6), Inches(1.55), Inches(12.1),
   col_widths=[Inches(5.6), Inches(6.5)], fsize=16)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "The hard part isn't syntax — control inverts. Your code must be ready to react at any moment."
r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
notes(s, "Land this: in a notebook you are in control; in a frontend the user is. The code "
         "is event-driven, not start-to-finish.")

# ── SLIDE 5 — Restaurant Analogy ───────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Restaurant Analogy")
tf = box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Backend = the kitchen (does the work, unseen).  Frontend = the dining room (menu, tables, waiter)."
r.font.size = Pt(17); r.font.color.rgb = SLATE; r.font.name = "Calibri"
table(s, [
    ["Web term", "Plain meaning", "Data-science analogy"],
    ["HTML", "Structure — headings, buttons", "A figure's axes & labels"],
    ["CSS", "Styling — color, spacing, layout", "plt.style / seaborn theme"],
    ["JavaScript", "Behavior — clicks, typing", "Code that re-runs a cell on input"],
    ["The DOM", "Live in-memory tree of the page", "A mutable figure you edit piece-by-piece"],
    ["Rendering", "Turning code/data into pixels", "plt.show() — but it fires many times"],
], Inches(0.6), Inches(1.95), Inches(12.1),
   col_widths=[Inches(2.5), Inches(4.8), Inches(4.8)], fsize=14)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "A great kitchen with no dining room is just a notebook — only the cook can use it."
r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = GREY; r.font.name = "Calibri"

# ── SLIDE 6 — System Architecture ──────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "System Architecture")
node(s, Inches(0.6), Inches(3.3), Inches(2.3), Inches(1.0), "User's\nbrowser", SLATE)
node(s, Inches(3.15), Inches(1.5), Inches(2.3), Inches(0.7), "middleware.ts\n(auth gate)", AMBER, DEEP, fs=12)
# Next.js cluster
band(s, LIGHT, Inches(3.05), Inches(2.45), Inches(5.4), Inches(3.8))
tf = box(s, Inches(3.2), Inches(2.5), Inches(5.1), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Next.js app  (Vercel)"; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
node(s, Inches(3.25), Inches(3.0), Inches(5.0), Inches(0.95), "Server Components\nbuilt on server → finished HTML\n(home, analytics)", GREEN, fs=11)
node(s, Inches(3.25), Inches(4.15), Inches(5.0), Inches(0.95), "Client Components\nrun in browser, interactive\n(map, dashboard, nav)", BLUE, fs=11)
node(s, Inches(3.25), Inches(5.3), Inches(5.0), Inches(0.7), "axios  (lib/api.ts → endpoints.ts)\nJWT auto-attached", DEEP, fs=11)
# right side
node(s, Inches(8.95), Inches(2.6), Inches(3.8), Inches(0.8), "FastAPI backend (Render)", DEEP)
node(s, Inches(8.95), Inches(3.6), Inches(3.8), Inches(0.8), "XGBoost · RAG · EM-DAT JSON", SLATE)
node(s, Inches(8.95), Inches(4.6), Inches(3.8), Inches(0.8), "PostgreSQL (Neon)", SLATE)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Server-rendered pages arrive instantly; interactive parts boot in the browser and fetch JSON on demand."
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = SLATE; r.font.name = "Calibri"
notes(s, "Two big ideas: (1) some pages are pre-built on a server — fast, no spinner; (2) "
         "interactive parts come alive only in the browser. Analytics is server-rendered; "
         "the Leaflet map is browser-only because Leaflet needs window.")

# ── SLIDE 7 — Tech Stack ───────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Tech Stack at a Glance")
table(s, [
    ["Tool", "One-line job", "Where in the repo"],
    ["Next.js 14", "Framework: routing + server rendering", "whole frontend/"],
    ["React 18", "Build UI from reusable components", "every .tsx file"],
    ["TypeScript (strict)", "JS + enforced types; backend contract", "frontend/types/"],
    ["Tailwind v3", "Styling via utility classes", "tailwind.config.ts"],
    ["NextAuth v5", "Login / session in a secure cookie", "auth.ts, middleware.ts"],
    ["Axios", "HTTP client to the backend", "lib/api.ts, lib/endpoints.ts"],
    ["Leaflet + react-leaflet", "Interactive maps", "components/RiskMap.tsx"],
    ["Recharts", "Charts (line / bar / composed)", "components/analytics/*"],
    ["next-pwa", "Installable + offline cache (5 GETs)", "next.config.js"],
], Inches(0.6), Inches(1.4), Inches(12.1),
   col_widths=[Inches(3.1), Inches(5.0), Inches(4.0)], fsize=13)
tf = box(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Deliberately NOT here: no TanStack Query/SWR, no react-hook-form/zod, no shadcn/Radix — hand-built (covered later)."
r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = AMBER; r.font.name = "Calibri"

# ── SLIDE 8 — Component 1: React + Next.js ─────────────────────────────────────
component(
    1, "React + Next.js (the framework)",
    "A React component is a reusable plotting function that returns UI instead of a figure.",
    [
        ("Props = function arguments (read-only). State = a variable that, when changed, re-renders.", 0, SLATE, False),
        ("JSX = the HTML-looking syntax describing what to draw.", 0, SLATE, False),
        ("React draws components; Next.js adds routing, servers, and deployment on top.", 0, SLATE, False),
    ],
    code_lines=[
        'export function SeverityBadge({ level }: { level: SeverityLevel }) {',
        '  const cls = SEVERITY_CLASSES[level] ?? SEVERITY_CLASSES.Low',
        '  return <span className={`... ${cls}`}>{S(`severity.${level}`)}</span>',
        '}',
    ],
    chose="plain React + Vite (hand-wire routing/SSR), Remix, Vue/Svelte — Next.js is the lowest-friction path to a real, deployable site.",
)

# ── SLIDE 9 — Component 2: Routing ─────────────────────────────────────────────
component(
    2, "Routing — the App Router",
    "The folder structure literally IS the routing table. A folder map/ with page.tsx = the /map route.",
    [
        ("Route groups in (parentheses) organize without changing the URL:", 0, SLATE, False),
        ("(public)/ → home, map, analytics, pricing, forecast", 1, SLATE, False),
        ("(auth)/ → login, register, verify-email", 1, SLATE, False),
        ("(protected)/ → dashboard, admin  (gated by middleware)", 1, SLATE, False),
        ("layout.tsx = shared shell · loading.tsx = automatic skeleton.", 0, SLATE, False),
    ],
    chose="the older Pages Router (legacy) and React Router (means leaving Next.js). Only the App Router supports Server Components.",
)

# ── SLIDE 10 — Component 3: Server vs Client + Rendering ────────────────────────
s = slide(); bg(s, WHITE)
title_bar(s, "Server vs Client Components & Rendering", kicker="Component 3")
tf = box(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.7)); p = tf.paragraphs[0]
r = p.add_run(); r.text = ("Server Component = a cell run once on a server, you get the output.  "
                           'Client Component ("use client") = an interactive widget in the browser.')
r.font.size = Pt(15); r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
bullets(s, [
    ("Rule of thumb: display-only → Server; anything the user touches → Client.", 0, SLATE, True),
], top=Inches(2.15), size=17, gap=6)
table(s, [
    ["Strategy", "Plain meaning", "Where we use it"],
    ["ISR", "Static page, auto-refreshed every N sec", "analytics (revalidate=86400), home (3600)"],
    ["CSR", "Built live in the browser", "the Leaflet map — dynamic(..., {ssr:false})"],
    ["SSG / SSR", "Prebuilt at deploy / per request", "most of ~17 routes prerendered"],
], Inches(0.6), Inches(2.85), Inches(12.1),
   col_widths=[Inches(2.0), Inches(5.0), Inches(5.1)], fsize=14)
band(s, LIGHT, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.2))
tf = box(s, Inches(0.8), Inches(5.75), Inches(11.7), Inches(1.0)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Why it matters:  "
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = "Calibri"
r = p.add_run(); r.text = ("analytics renders on the server from precomputed JSON — zero runtime DB queries. "
                           "The map is forced browser-only because Leaflet needs `window`, which a server doesn't have.")
r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"

# ── SLIDE 11 — Component 4: TypeScript ─────────────────────────────────────────
component(
    4, "TypeScript — the contract with the backend",
    "TypeScript = JavaScript with enforced types — like mypy turned on and blocking.",
    [
        ("frontend/types/ is a 1:1 mirror of the backend's Pydantic schemas.", 0, SLATE, False),
        ("A renamed/missing backend field → a red squiggle at edit time, not a runtime crash for a user.", 0, AMBER, True),
    ],
    code_lines=[
        'interface PredictionResult {',
        '  disaster_type: DisasterType            // one of 8 literals',
        '  probability_score: number',
        '  severity_level: "Low"|"Medium"|"High"|"Critical"',
        '  risk_score: number',
        '}',
    ],
    chose="plain JavaScript — a backend rename would silently break the UI with no warning.",
)

# ── SLIDE 12 — Component 5: Tailwind ───────────────────────────────────────────
component(
    5, "Styling with Tailwind",
    'CSS = how things look. Tailwind = "utility-first": compose tiny classes on the element itself.',
    [
        ("p-5 padding · rounded-xl corners · text-slate-900 color · border.", 0, SLATE, False),
        ("DS analogy: like passing color=, linewidth=, alpha= inline on every plot call.", 0, SLATE, False),
        ("Custom severity palette (green→yellow→orange→red), shared with the map's risk scale.", 0, SLATE, False),
    ],
    code_lines=[
        '<div className="rounded-xl border border-slate-200 bg-white p-5">',
        '  <div className="text-3xl font-bold text-slate-900">{value}</div>',
        '</div>',
    ],
    chose="plain CSS / CSS Modules (naming + file-hopping) and CSS-in-JS (runtime cost, awkward with Server Components).",
)

# ── SLIDE 13 — Component 6: State & Hooks ──────────────────────────────────────
component(
    6, "State & Interactivity (hooks)",
    "State = a variable that, when changed, auto-re-runs the cell that displays it.",
    [
        ("useState → changeable data.  useEffect(fn,[deps]) → side-effect when inputs change ([] = run-once __init__).", 0, SLATE, False),
        ("useMemo ≈ lru_cache for a derived value.", 0, SLATE, False),
        ('Pick "Flood" → setFilterType("Flood") → re-render → map redraws. You never wrote "redraw".', 0, GREEN, True),
    ],
    code_lines=[
        'const [points, setPoints] = useState<RiskMapPoint[]>([])',
        'useEffect(() => { endpoints.regions.riskMap().then(setPoints) }, [])',
        'const filtered = useMemo(() => points.filter(...), [points, filterType])',
    ],
    chose="TanStack Query / SWR — we hand-roll loading/error/result with useState. Simple enough now; the first thing I'd add if it grew.",
    note="This is the heart of the talk for a data-science audience. The mental model: "
         "change a variable, the UI reconciles itself. You never imperatively redraw.",
)

# ── SLIDE 14 — Component 7: Talking to the Backend ─────────────────────────────
component(
    7, "Talking to the Backend (Axios)",
    "Two Axios instances + typed wrappers — the rest of the app never writes a raw URL.",
    [
        ("api → Server Components (public, no token).  apiClient → Client Components.", 0, SLATE, False),
        ("apiClient interceptor auto-attaches the JWT from the NextAuth session.", 0, SLATE, False),
        ("A response interceptor normalizes every failure to one shape: ApiError {status, detail, original}.", 0, SLATE, False),
        ("lib/endpoints.ts = one typed function per route.", 0, SLATE, False),
    ],
    code_lines=[
        'predict: (body, client = apiClient) =>',
        '  client.post<PredictionResult>("/predictions/predict", body)',
        '        .then(r => r.data)',
    ],
    chose="scattered raw fetch() calls — centralizing gives one place for auth headers, timeouts, and error normalization.",
)

# ── SLIDE 15 — Component 8: Authentication ─────────────────────────────────────
component(
    8, "Authentication (NextAuth + middleware)",
    "Authn = who are you (a signed JWT).  Authz = what may you do.",
    [
        ("Backend tokens live inside the NextAuth JWT, stored in an HttpOnly cookie — JS can't read it. Never localStorage.", 0, SLATE, False),
        ("Access token expires in 30 min → the JWT callback proactively refreshes 60s early. User never notices.", 0, SLATE, False),
        ("middleware.ts guards /dashboard, /admin, … → redirect if not allowed.", 0, SLATE, False),
        ("Golden rule: middleware is UX only — the backend Depends() guards are the real boundary.", 0, AMBER, True),
    ],
    chose="rolling our own cookie/refresh (security-sensitive) and Clerk/Auth0 (cost + duplicates our Postgres users).",
)

# ── SLIDE 16 — Component 9: Maps ───────────────────────────────────────────────
component(
    9, "Maps (Leaflet)",
    "react-leaflet = folium for the browser: coordinates → clickable markers + popups, free OSM tiles, no API key.",
    [
        ("RiskMap.tsx: fetch /regions/risk-map → render each point as a CircleMarker colored by risk band.", 0, SLATE, False),
        ("Hover → tooltip (precomputed, no network).  Click → live prediction (logged-in) or sign-up CTA (guest).", 0, SLATE, False),
        ("Honest: leaflet.heat is still a dependency but UNUSED — discrete markers guarantee color = exactly one risk level.", 0, AMBER, True),
    ],
    chose="Mapbox / Google Maps — both need an API key + billing. Leaflet + OpenStreetMap is genuinely free (a hard requirement here).",
)

# ── SLIDE 17 — Component 10: Charts ────────────────────────────────────────────
component(
    10, "Charts (Recharts)",
    "Recharts ≈ matplotlib / plotly, but interactive and in the browser.",
    [
        ("Describe the chart as components, feed it data rows — data prep is pure DS muscle memory.", 0, SLATE, False),
        ("Used for: forecast line chart, analytics bars, and a ComposedChart with a client-side linear regression.", 0, SLATE, False),
        ("Slope → Increasing / Decreasing / Stable: np.polyfit overlaid in the browser.", 0, SLATE, False),
    ],
    code_lines=[
        'const rows = days.map(d => ({',
        '  day: d.forecast_day_offset + 1,',
        '  probability_pct: +(d.probability_score * 100).toFixed(1),',
        '}))',
    ],
    chose="D3 (too low-level — hand-build axes), Chart.js (canvas, less React-idiomatic), Plotly (heavier).",
)

# ── SLIDE 18 — One Request End to End (money slide) ────────────────────────────
s = slide(); bg(s, WHITE)
title_bar(s, "One Screen, End to End", kicker="The Money Slide")
tf = box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "A logged-in user opens /map and clicks a Flood marker:"
r.font.name = "Calibri"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DEEP
steps = [
    ("1", "URL → gate", "middleware.ts: /map is public → passes through"),
    ("2", "Server shell", "map/page.tsx (Server Component) renders Nav + 'Loading map…'"),
    ("3", "Browser-only boot", "dynamic({ssr:false}) loads RiskMap.tsx (Leaflet needs window)"),
    ("4", "Fetch points", "useEffect → endpoints.regions.riskMap() → setPoints → markers paint"),
    ("5", "Click marker", "setPending({lat, lon, 'Flood'}) → <Popup> opens"),
    ("6", "Live predict", "endpoints.predictions.predict(...) via apiClient — JWT auto-attached"),
    ("7", "JSON → UI", "setResult → re-render → SeverityBadge + probability + risk score"),
    ("8", "On error", "429 → rate-limit msg, 401 → sign-in msg (from lib/strings.ts)"),
]
y = 1.95
for num, head, detail in steps:
    band(s, BLUE, Inches(0.6), Inches(y), Inches(0.5), Inches(0.42))
    tf = box(s, Inches(0.6), Inches(y - 0.02), Inches(0.5), Inches(0.46)); p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    tf = box(s, Inches(1.25), Inches(y - 0.04), Inches(11.4), Inches(0.5)); p = tf.paragraphs[0]
    r = p.add_run(); r.text = head + "  —  "
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
    r = p.add_run(); r.text = detail
    r.font.size = Pt(13); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    y += 0.5
tf = box(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.7)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Every interactive screen is this loop:  URL → middleware → server shell → client component → fetch → state → re-render → pixels."
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
notes(s, "Walk this slowly, naming the real file at each hop. This is where it clicks. "
         "If short on time, cut a component slide — NOT this one.")

# ── SLIDE 19 — DS Translation Table (killer slide) ─────────────────────────────
s = slide(); bg(s, WHITE)
title_bar(s, "The Data-Scientist Translation Table")
table(s, [
    ["Frontend concept", "Closest data-science equivalent"],
    ["Component", "Reusable plotting function returning UI"],
    ["Props", "Function arguments (read-only)"],
    ["State (useState)", "A variable that re-runs the cell when it changes"],
    ["useEffect(fn,[deps])", "Callback fired when watched inputs change ([] = __init__)"],
    ["Re-render", "plt.show() firing again automatically"],
    ["Server Component", "A cell run once on a server; you get the output"],
    ["TypeScript types", "mypy + Pydantic schemas as a backend contract"],
    ["Axios + endpoints.ts", "A typed requests wrapper, one fn per endpoint"],
    ["JSON from the API", "A record / DataFrame you render to the screen"],
    ["Recharts / Leaflet", "matplotlib·plotly / folium — in the browser"],
], Inches(0.9), Inches(1.4), Inches(11.5),
   col_widths=[Inches(4.3), Inches(7.2)], fsize=14)
notes(s, "For this audience this is the slide that makes everything land. Pause on it. "
         "Invite them to map any frontend term onto something they already do daily.")

# ── SLIDE 20 — Trade-offs ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Engineering Decisions & Trade-offs")
table(s, [
    ["Decision", "Chosen", "Over", "Why"],
    ["Framework", "Next.js 14", "Vite / Remix", "SSR + routing + deploy, free"],
    ["Public pages", "ISR (revalidate)", "live SSR", "Precomputed JSON; 0 runtime DB"],
    ["Map", "Leaflet + OSM", "Mapbox / Google", "No API key, no billing"],
    ["Map render", "Discrete markers", "Heatmap", "Color = exactly one risk level"],
    ["Charts", "Recharts", "D3 / Chart.js", "Declarative, React-idiomatic"],
    ["Auth", "NextAuth JWT cookie", "Roll-your-own / Auth0", "Secure, free, fits backend"],
    ["Types", "TS strict, mirror API", "plain JS", "Compile-time contract"],
], Inches(0.5), Inches(1.45), Inches(12.3),
   col_widths=[Inches(2.2), Inches(3.0), Inches(2.9), Inches(4.2)], fsize=13)
tf = box(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Theme: match tools to a single-developer, zero-cost constraint — and know what each choice trades away."
r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
notes(s, "Engineering maturity is not 'I used the trendiest tool' but 'I matched tools to "
         "constraints and I know the trade-offs.'")

# ── SLIDE 21 — Honest Engineering ──────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Honest Engineering: What We Left Out")
tf = box(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Deliberately NOT added — portfolio-scale choices, each with an upgrade path:"
r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = AMBER; r.font.name = "Calibri"
bullets(s, [
    ("No data-fetching library (TanStack Query/SWR) → hand-rolled loading/error in each useEffect.", 0, SLATE, False),
    ("Upgrade: add Query for caching, retries, request dedup.", 1, GREY, False),
    ("No form library (react-hook-form/zod) → plain useState + minimal manual validation.", 0, SLATE, False),
    ("Upgrade: zod mirrors Pydantic nicely.", 1, GREY, False),
    ("No component library (shadcn/Radix) or icon pack → hand-built Tailwind atoms.", 0, SLATE, False),
    ("Upgrade: shadcn (= Radix + Tailwind) for accessible primitives.", 1, GREY, False),
    ("leaflet.heat dependency is now unused → remove.", 0, SLATE, False),
], top=Inches(2.0), size=16, gap=7)
band(s, LIGHT, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7))
tf = box(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Knowing what you traded away is the engineering — not a weakness."
r.font.size = Pt(17); r.font.bold = True; r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"

# ── SLIDE 22 — Results ─────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Results")
bullets(s, [
    ("Live in production: safeearth.tech (Vercel) ↔ api.safeearth.tech (Render)", 0, GREEN, True),
    ("17 routes, mostly prerendered (static / ISR) — fast first paint, free hosting", 0, SLATE, False),
    ("Full auth flow: register → verify → login → role-aware UI, JWT auto-refresh", 0, SLATE, False),
    ("Interactive risk map, 4-tab analytics, prediction dashboard, 30-day forecast", 0, SLATE, False),
    ("TypeScript strict; types mirror the backend 1:1; zero hardcoded UI strings (i18n-ready)", 0, SLATE, False),
    ("PWA: installable + offline caching of public read-only endpoints", 0, SLATE, False),
], top=Inches(1.7), size=19, gap=15)

# ── SLIDE 23 — Future Work ─────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Future Work")
bullets(s, [
    ("Adopt TanStack Query for caching, retries, and request dedup", 0, SLATE, False),
    ("Add react-hook-form + zod for richer, schema-validated forms", 0, SLATE, False),
    ("Bring in shadcn/ui (Radix) for accessible dialogs/menus + keyboard support", 0, SLATE, False),
    ("Multi-line forecast chart (needs a backend endpoint returning all 8 type probabilities)", 0, SLATE, False),
    ("Remove the dead leaflet.heat dependency; audit accessibility (a11y) end-to-end", 0, SLATE, False),
], top=Inches(1.8), size=20, gap=18)

# ── SLIDE 24 — Takeaways ───────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("The frontend turns JSON into something a human can act on — notebook vs product.", 0, DEEP, True),
    ("Server-by-default, interactive-where-needed keeps it fast and cheap.", 0, SLATE, False),
    ("One contract end to end: TS types mirror the backend, so mismatches fail at compile time.", 0, SLATE, False),
    ("Every screen is the same loop: URL → gate → render → fetch → state → re-render.", 0, SLATE, False),
    ("Every trade-off was deliberate, documented, and reversible.", 0, SLATE, False),
], top=Inches(1.9), size=20, gap=20)

# ── SLIDE 25 — Thank You ───────────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, BLUE, 0, Inches(3.0), W, Inches(0.08))
tf = box(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(4.6), Inches(11.3), Inches(1.4)); p = tf.paragraphs[0]
for i, txt in enumerate([
    "SafeEarth Intelligence — Next.js · React · TypeScript · Tailwind · Leaflet · Recharts",
    "Full technical deep-dive: docs/FRONTEND_EXPLAINED.md",
    "safeearth.tech",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Likely Qs: Why server+client split? (speed + SEO for content; interactivity only "
         "where needed). Why no Redux/Query? (app state is small; component state + a token "
         "bridge suffice). Is the frontend the security boundary? (No — UX only; backend "
         "Depends() re-checks every request). Why Leaflet over Mapbox? (free, no API key/"
         "billing). Heatmap vs markers? (markers guarantee color = one risk level).")

out = Path(__file__).resolve().parent / "SafeEarth_Frontend_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides)")
