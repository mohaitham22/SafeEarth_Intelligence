"""
build_preprocessing_pptx.py — generate an editable PowerPoint (with charts) from the
SafeEarth preprocessing deep dive.

Run:  py -3.12 docs/build_preprocessing_pptx.py
Out:  docs/SafeEarth_Preprocessing_Presentation.pptx

Native python-pptx (editable text boxes + tables) + embedded PNG charts rendered by
make_preprocessing_charts.py from the real artifacts. Content mirrors
docs/PREPROCESSING_PITCH_DECK.md, derived from docs/PREPROCESSING_DEEP_DIVE.md.
Every number is from this project's own metrics/ and streamlit_eda/experiments/.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
import make_preprocessing_charts  # noqa: E402

ASSETS = Path(__file__).resolve().parent / "preprocessing_assets"

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xDC, 0x26, 0x26)
SLATE  = RGBColor(0x33, 0x41, 0x55)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
PINK   = RGBColor(0xFE, 0xE2, 0xE2)

W, H = Inches(13.333), Inches(7.5)

# Ensure chart PNGs exist before any slide embeds them.
if not (ASSETS / "01_regression_winner.png").exists():
    print("Charts not found — rendering them first...")
    make_preprocessing_charts.generate_all()

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def img(s, name, left, top, width):
    return s.shapes.add_picture(str(ASSETS / name), left, top, width=width)


def title_bar(s, text, kicker=None, color=DEEP):
    band(s, color, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = (RGBColor(0xFE, 0xCA, 0xCA) if color == RED else GREEN)
        r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    first = True
    for it in items:
        text, level, color, bold = (it + (0, SLATE, False))[:4] if isinstance(it, tuple) else (it, 0, SLATE, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0)); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None, hcolor=DEEP):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val); r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hcolor
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def callout(s, label, text, top=Inches(6.5), color=LIGHT, label_color=GREEN, h=Inches(0.7),
            left=Inches(0.6), width=Inches(12.1)):
    band(s, color, left, top, width, h)
    tf = box(s, left + Inches(0.2), Inches(top.inches + 0.10), Inches(width.inches - 0.4), Inches(h.inches - 0.16))
    p = tf.paragraphs[0]
    if label:
        r = p.add_run(); r.text = label + "  "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = label_color; r.font.name = "Calibri"
    r = p.add_run(); r.text = text
    r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"


def chip(s, left, top, text, color):
    """Small rounded box used to draw the pipeline flow."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.55), Inches(0.62))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.color.rgb = color; sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    return sh


def arrow(s, left, top):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.42), Inches(0.30))
    sh.fill.solid(); sh.fill.fore_color.rgb = GREY; sh.line.fill.background(); sh.shadow.inherit = False


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.45), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
for i, (txt, sz, col, bold) in enumerate([
    ("Preprocessing Deep Dive", 28, GREEN, True),
    ("Every transform, and how it moved classification F1 and regression error", 18, RGBColor(0xCB, 0xD5, 0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.35), Inches(11.3), Inches(1.1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "EM-DAT (1900–2021)  ·  XGBoost + CatBoost classifier  ·  per-type XGB/RF regressors"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Thesis of the talk: preprocessing — not the algorithm — drove the wins here. "
         "The single biggest improvement in the whole project (impact error 111x -> 2.7x) is a "
         "missing-value decision, not a model change. Every number is from our own metrics/ and "
         "streamlit_eda/experiments/ JSON, evaluated on a held-out test set.")

# ── SLIDE 2 — Roadmap ─────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Roadmap")
bullets(s, [
    ("The split & the pipeline order — where every fitted transform sits vs the split", 0, SLATE, False),
    ("Missing values — 81.7% of rows have no coordinates (the governing fact)", 0, SLATE, False),
    ("CLASSIFICATION lever — class weights beat SMOTE & inverse-frequency (+0.0123 macro-F1)", 0, GREEN, True),
    ("REGRESSION lever — drop-null + per-type cut impact error up to 111× → 2.7×", 0, GREEN, True),
    ("Decision discipline — every choice ablated against a noise floor; leaks caught & rejected", 0, SLATE, False),
    ("Two honest gaps + prioritised fixes (one validated win not yet shipped)", 0, RED, True),
], top=Inches(1.7), size=20, gap=18)
notes(s, "Set expectations: I separate CLASSIFICATION effects (macro/weighted F1) from REGRESSION "
         "effects (log-MAE / error factor) the whole way through, because different steps help "
         "different models. I'll also show negative results and a validated-but-unshipped win — "
         "evidence of real ablation, not cherry-picking.")

# ── SLIDE 3 — Split strategy ──────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Split — Pre-Materialised, Not Programmatic")
bullets(s, [
    ("Two CSV files ARE the split — no train_test_split for the main model", 0, SLATE, False),
    ("Train 1900–2021 (14,476 rows, 8 types) · Holdout 1970–2021 (13,070 rows)", 0, SLATE, False),
    ("random_state=42 on every downstream CV / sub-split / Optuna sampler / SMOTENC", 0, SLATE, False),
    ("Neither temporal nor geographic structure is respected:", 0, AMBER, True),
    ("test window (1970–2021) ⊂ train window (1900–2021); same regions in both (by design — encoders must see them)", 1, SLATE, False),
], top=Inches(1.45), size=17, gap=10)
table(s, [
    ["Sub-split", "Used by", "Scheme"],
    ["Optuna HP search", "run_training.py", "StratifiedKFold(3) on TRAIN, macro-F1 objective"],
    ["Minority-strategy CV", "train_minority_strategies.py", "StratifiedKFold(5) on TRAIN"],
    ["Weight / threshold tuning", "tune_*.py", "stratified 50/50 of TEST → val 6,535 / holdout 6,535"],
], Inches(0.6), Inches(4.55), Inches(12.1),
   col_widths=[Inches(3.2), Inches(4.0), Inches(4.9)], fsize=12.5)
callout(s, "Validity:", "this tests generalization across event records, not forecasting the future — read macro-F1 as in-distribution (see Problem on Slide 14).",
        top=Inches(6.5), h=Inches(0.7), color=LIGHT, label_color=AMBER)

# ── SLIDE 4 — Pipeline order of operations ────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Order of Operations (where leakage would hide)")
# Row 1 — shared front
y1 = Inches(1.45)
chip(s, Inches(0.5), y1, "latin-1 load +\ntype filter", DEEP)
arrow(s, Inches(3.12), Inches(1.60))
chip(s, Inches(3.6), y1, "parse_coord +\nrange mask", DEEP)
arrow(s, Inches(6.22), Inches(1.60))
chip(s, Inches(6.7), y1, "lat/lon median\nimpute", AMBER)
arrow(s, Inches(9.32), Inches(1.60))
chip(s, Inches(9.8), y1, "16-feature\nengineering", DEEP)
# train/test boundary banner
band(s, GREEN, Inches(0.5), Inches(2.35), Inches(12.1), Inches(0.42))
tf = box(s, Inches(0.6), Inches(2.38), Inches(12.0), Inches(0.36)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "LabelEncoders · region_freq_map · month-mode · class weights — ALL fit on TRAIN only (test → unseen-category 0)"
r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
# Two branches
tf = box(s, Inches(0.5), Inches(3.0), Inches(6.0), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "CLASSIFIER PATH"; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = "Calibri"
tf = box(s, Inches(6.9), Inches(3.0), Inches(6.0), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "REGRESSOR PATH"; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
y2 = Inches(3.5)
chip(s, Inches(0.5), y2, "class sample\nweights", BLUE)
arrow(s, Inches(3.12), Inches(3.65))
chip(s, Inches(3.6), y2, "Optuna 3-fold\n(macro-F1)", BLUE)
chip(s, Inches(6.9), y2, "targets:\nNaN PRESERVED", GREEN)
arrow(s, Inches(9.52), Inches(3.65))
chip(s, Inches(10.0), y2, "per-type\ndrop-null fit", GREEN)
y3 = Inches(4.35)
chip(s, Inches(0.5), y3, "fit XGB / CAT\n+ weights", BLUE)
arrow(s, Inches(3.12), Inches(4.50))
chip(s, Inches(3.6), y3, "grid-search\nensemble wts", BLUE)
chip(s, Inches(6.9), y3, "log1p target", GREEN)
arrow(s, Inches(9.52), Inches(4.50))
chip(s, Inches(10.0), y3, "expm1 + EM-DAT\nblend + p99", GREEN)
callout(s, "No feature scaling anywhere:", "every model is a tree (XGB / CatBoost / RandomForest) → scale-invariant. Inference rebuilds the identical 16-feature vector; coord imputation never runs at serve time (lat/lon is supplied).",
        top=Inches(5.45), h=Inches(0.95), color=LIGHT, label_color=GREEN)
notes(s, "Walk the flow left-to-right. The green banner is the leakage guarantee — every fitted "
         "transform is fit on train only. The amber 'lat/lon median impute' chip is the one I flag "
         "(it's fit per-split offline, but never runs at serve time since the caller supplies coords).")

# ── SLIDE 5 — Data cleaning ───────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Data Cleaning — the Correctness Preconditions")
table(s, [
    ["Step", "What", "Why"],
    ["Encoding", "read_csv(encoding='latin-1')", "raw export has °C, mojibake Â°C, accented names"],
    ["Type normalize", "Disaster Type .str.strip(); keep 8 types", "CSV stores 'Extreme temperature ' (trailing space)"],
    ["Coordinate parse", "regex '34.01 N' / '78.46 W' → signed float", "274 rows store directional strings, not floats"],
    ["Impossible values", "lat∉[−90,90] / lon∉[−180,180] → NaN", "DMS artefacts stored as 36100.0 etc."],
    ["Magnitude → numeric", "to_numeric(coerce).fillna(0) + has_magnitude", "present in only ~38% of rows; absence is signal"],
    ["Target sign", "targets .clip(lower=0)", "no negative deaths / damage"],
    ["Deduplication", "NONE", "not performed — flagged on Slide 14"],
], Inches(0.5), Inches(1.4), Inches(12.3),
   col_widths=[Inches(2.7), Inches(4.7), Inches(4.9)], fsize=12.5)
callout(s, "Not an ablation arm:", "cleaning is a correctness precondition, not a measured delta — but the coord parse + range mask feed directly into the imputation quality measured next.",
        top=Inches(6.5), h=Inches(0.7), label_color=AMBER)

# ── SLIDE 6 — Missing values (chart) ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Missing Values — the Governing Fact")
bullets(s, [
    ("81.7% of TRAIN rows have no usable lat/lon", 0, RED, True),
    ("filled by country → continent → global median chain", 1, SLATE, False),
    ("So raw latitude/longitude are mostly imputed constants per country", 0, SLATE, False),
    ("→ a BETTER coordinate fill moves the needle (next slide);", 1, GREY, False),
    ("→ lat/lon-derived geo features add little (Slide 13)", 1, GREY, False),
    ("has_magnitude missingness indicator preserves the 'value absent' bit", 0, GREEN, True),
    ("→ top-2 SHAP feature; drives Earthquake F1 = 0.976", 1, SLATE, False),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.5), size=16, gap=11)
img(s, "04_coord_missing.png", Inches(6.3), Inches(1.7), Inches(6.6))
notes(s, "This one fact governs the rest of the imputation story. When 4 in 5 rows have no real "
         "coordinates, how you fill them is an outcome-relevant modelling decision, not cosmetic.")

# ── SLIDE 7 — Centroid imputation (chart) ─────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Coordinate Imputation — a Validated, Unshipped Win", kicker="Measured · classification + regression")
bullets(s, [
    ("Clean A/B: median chain vs curated ISO3 centroids (fills 11,707 / 11,822)", 0, SLATE, False),
    ("CLASSIFICATION: macro-F1 0.7285 → 0.7452", 0, GREEN, True),
    ("+0.0167 > CV-std 0.008 → KEEP", 1, SLATE, False),
    ("REGRESSION: mean log-MAE 2.3526 → 2.3215", 0, GREEN, True),
    ("(all four impact targets improved)", 1, SLATE, False),
    ("Minorities gain most: Wildfire +0.056, Volcanic +0.029", 0, DEEP, True),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.3), size=16, gap=11)
img(s, "05_centroid.png", Inches(6.1), Inches(1.85), Inches(6.9))
callout(s, "Status:", "validated but NOT in run_training.py yet — the highest-confidence un-deployed preprocessing win in the repo (Recommendation P0).",
        top=Inches(6.55), h=Inches(0.62), color=LIGHT, label_color=RED)
notes(s, "This is a strength to show, not hide: an experiment with a clean KEEP verdict that hasn't "
         "shipped. Risk-free to deploy because lat/lon is caller-supplied at inference — only the "
         "training fill changes.")

# ── SLIDE 8 — Outliers & skew ─────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Outliers & Skew — Cap Statistics, Not the Rows")
bullets(s, [
    ("Impact is right-skewed ~100×+ (Flood mean deaths 1,735 vs median 16)", 0, SLATE, False),
    ("Genuine extreme events carry the signal — they must NOT be clipped away", 0, AMBER, True),
    ("Three mechanisms, none touching the rows fed to the trees:", 0, DEEP, True),
    ("log1p target transform — the primary regression skew fix (expm1 at predict)", 1, SLATE, False),
    ("per-type p99 cap on EM-DAT MEDIANS — one tsunami can't inflate a small country's median", 1, SLATE, False),
    ("per-type p99 ceilings on impact OUTPUTS — guards against absurd over-prediction", 1, SLATE, False),
    ("No input-feature winsorizing — trees are split-invariant, so none is needed", 0, GREEN, True),
], top=Inches(1.5), size=17, gap=12)
callout(s, "Classification vs regression:", "skew handling is a REGRESSION concern (log1p targets). The classifier is unaffected — trees split on thresholds regardless of input skew.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)

# ── SLIDE 9 — Encoding + scaling ──────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Categorical Encoding & the Scaling No-Op")
bullets(s, [
    ("Label-encode continent(5) / region(23) / country(225); unseen → 0 at test", 0, SLATE, False),
    ("historical_freq = FREQUENCY encoding (train event count per region) + log1p", 0, SLATE, False),
    ("No target/mean encoding → no per-fold target-encoding leakage surface", 0, GREEN, True),
    ("Ordinal country codes are safe ONLY because every consumer is a tree", 0, AMBER, True),
    ("CatBoost gets the same float matrix — native cat_features raises on float arrays", 0, GREY, False),
], left=Inches(0.55), top=Inches(1.5), width=Inches(6.6), size=16, gap=12)
table(s, [
    ["Model", "Needs scaling?", "In project?"],
    ["XGBoost (clf + reg)", "No", "✓ used"],
    ["CatBoost (clf)", "No", "✓ used"],
    ["RandomForest (reg)", "No", "✓ used"],
    ["Logistic / SVM / SVR / KNN", "YES", "✗ not used"],
], Inches(7.4), Inches(1.65), Inches(5.3),
   col_widths=[Inches(2.9), Inches(1.4), Inches(1.0)], fsize=12.5)
callout(s, "Scaling delta = 0:", "no StandardScaler/RobustScaler anywhere — a no-op for the all-tree stack. It becomes mandatory the day a linear/distance model is added (Recommendation P7).",
        top=Inches(6.4), h=Inches(0.75), label_color=GREEN)

# ── SLIDE 10 — Class imbalance (chart) ────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Class Imbalance — Weights Beat Resampling", kicker="CLASSIFICATION lever")
bullets(s, [
    ("Macro-F1 is PRIMARY so minorities (Volcanic 222, Landslide 713) count equally", 0, SLATE, False),
    ("Hand-tuned sample weights SHIPPED: +0.0123 macro-F1 (the only v4.1→v4.2 change)", 0, GREEN, True),
    ("Drought +0.056, Landslide +0.013, Wildfire +0.010; no majority regression >0.02", 1, SLATE, False),
    ("SMOTE was done RIGHT — SMOTENC, cat-idx [5,6,7,11,14], in-fold only…", 0, DEEP, True),
    ("…and still HURT (−0.0324; Landslide −0.12): scattered minorities → implausible synthetics", 1, SLATE, False),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.6), size=15.5, gap=11)
img(s, "03_strategies.png", Inches(6.4), Inches(2.0), Inches(6.6))
callout(s, "Why minorities are hard:", "it's a feature-separability problem (Landslide overlaps Flood/Storm in lat/lon/season), not a row-count problem — synthetic interpolation makes it worse.",
        top=Inches(6.55), h=Inches(0.62), color=PINK, label_color=RED)
notes(s, "Research-maturity slide: a tested-and-rejected method WITH the reason it fails. Stress that "
         "SMOTE was implemented correctly (SMOTENC, in-fold) — it still lost, so the conclusion is "
         "about the data, not a bug.")

# ── SLIDE 11 — Target transformation (chart) ──────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Target Handling — 'Null' Is Not 'Zero'", kicker="REGRESSION lever")
bullets(s, [
    ("Old regressors filled null targets with 0 — but 73% of injuries/damage are nulls", 0, SLATE, False),
    ("Model learns ≈0, then mispredicts real events (injuries=1 next to $1B damage)", 0, RED, True),
    ("Drop-null (train only on observed rows) wins on EVERY target", 0, GREEN, True),
    ("The trap: scoring on a median-FILLED holdout makes a broken model look great", 0, AMBER, True),
    ("(damage MAE 0.83 'illusion' vs 1.66 honest)", 1, SLATE, False),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.4), size=15.5, gap=12)
img(s, "02_fill_ablation.png", Inches(6.1), Inches(1.95), Inches(6.9))
notes(s, "Evaluation-integrity teaching moment: the only HONEST way to score is on observed-target "
         "rows. Filling the test set the same way you filled training manufactures a fake score.")

# ── SLIDE 12 — Regression money slide (2 charts) ──────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Two Independent Regression Wins, Combined", kicker="The biggest improvement in the project")
img(s, "07_typeaware.png", Inches(0.45), Inches(1.5), Inches(6.0))
img(s, "01_regression_winner.png", Inches(6.7), Inches(1.5), Inches(6.3))
callout(s, "Combined (per-type + drop-null), shipped:", "damage error 111× → 2.7×, affected 24.5× → 4.3×, injuries 9× → 2.8×, deaths 2.8× → 1.9× — entirely a preprocessing/structure change, classifier untouched.",
        top=Inches(6.35), h=Inches(0.85), color=LIGHT, label_color=GREEN)
notes(s, "Left: adding the disaster-TYPE axis (type-blind regressors gave Flood and Earthquake the "
         "same impact at the same coords). Right: dropping fake zeros. Independent axes that compound. "
         "This is the single strongest result in the repo and it is 100% preprocessing.")

# ── SLIDE 13 — Decision discipline (chart) ────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Decision Discipline — Ablate Against a Noise Floor")
img(s, "06_decisions.png", Inches(3.0), Inches(1.35), Inches(7.4))
callout(s, "Rule:", "keep a step iff holdout_gain > baseline CV-std. Tempting gains were rejected when they were leaks (per-scale magnitude +0.0097 — Dis Mag Scale ≈ the label for 93% of rows) or within noise (geo −0.0006).",
        top=Inches(6.45), h=Inches(0.85), label_color=GREEN)
notes(s, "This single chart is the 'I do controlled experiments' slide. Blue = validated-not-shipped "
         "(centroid). Green = shipped. Red = a measured +0.0097 we REJECTED because it's leakage. "
         "Grey = discarded within noise or because it regressed. Discipline over chasing deltas.")

# ── SLIDE 14 — Leakage audit ──────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Leakage Audit — Two Honest Open Items", kicker="Cross-cutting", color=RED)
table(s, [
    ["Surface", "Status"],
    ["Encoders / freq-map / weights fit on TRAIN only", "✓ clean"],
    ["SMOTENC in training fold only; val & holdout never resampled", "✓ clean"],
    ["Target drop-null / log1p (targets only)", "✓ clean"],
    ["Per-scale magnitude leak (Dis Mag Scale ≈ label)", "✓ caught & rejected"],
    ["frac_* label-distribution features", "✓ caught & dropped"],
    ["lat/lon median impute fit PER-SPLIT (not train-only)", "⚠ minor — never runs at serve time"],
    ["Cross-file train↔test row deduplication", "✗ NOT done — unquantified"],
], Inches(0.5), Inches(1.4), Inches(12.3),
   col_widths=[Inches(8.6), Inches(3.7)], fsize=13)
callout(s, "Earthquake F1 = 0.976 is NOT a leak:", "it's the dense Richter magnitude channel (has_magnitude / dis_mag_value). The genuinely leaky signal (Dis Mag Scale) was identified and kept out of production.",
        top=Inches(6.45), h=Inches(0.75), color=PINK, label_color=RED)
notes(s, "Own the two open items confidently. The dedup gap is the biggest threat to the absolute "
         "numbers (train 1900-2021 and test 1970-2021 overlap in time). It threatens absolute "
         "macro-F1, not the RELATIVE decisions, which all used CV or the honest observed-row holdout.")

# ── SLIDE 15 — Impact matrix (centerpiece) ────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Preprocessing → Performance, in One Matrix")
table(s, [
    ["Step", "CLASSIFICATION (Δ macro-F1)", "REGRESSION (Δ log-MAE / factor)", "Measured?"],
    ["Centroid coord impute", "+0.0167 (KEEP, unshipped)", "−0.0311 mean (all 4 down)", "✓"],
    ["Hand-tuned class weights", "+0.0123 (SHIPPED)", "—", "✓"],
    ["Drop-null targets (vs fill-0)", "—", "damage 111× → 4.1× (global)", "✓"],
    ["Per-type regressors", "—", "mean 2.353 → 1.836", "✓"],
    ["Per-type + drop-null (shipped)", "—", "damage 111× → 2.7×", "✓"],
    ["Feature scaling", "0 (no-op)", "0 (no-op)", "n/a"],
    ["SMOTENC / balanced-inv-freq", "−0.032 / −0.077 (discard)", "—", "✓"],
    ["Per-scale magnitude (leak)", "+0.0097 — REJECTED", "—", "✓"],
], Inches(0.4), Inches(1.35), Inches(12.5),
   col_widths=[Inches(3.6), Inches(3.5), Inches(3.6), Inches(1.8)], fsize=12)
callout(s, "No RMSE/R² stored:", "regression is reported as log1p-MAE on observed rows (↔ multiplicative error factor) — RMSE would be swamped by mega-events; R² is uninformative under 27–68% coverage. Adding back-transformed RMSE/R² is Recommendation P3.",
        top=Inches(6.5), h=Inches(0.78), label_color=AMBER)

# ── SLIDE 16 — Recommendations ────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Recommendations — Highest Ceiling First")
bullets(s, [
    ("P0 · Ship centroid coordinate imputation — already validated (+0.0167 macro, −0.0311 log-MAE); risk-free at inference", 0, GREEN, True),
    ("P1 · Deduplicate train vs test before trusting the holdout — biggest threat to every absolute number", 0, RED, True),
    ("P2 · Re-tune the inference blend (still 70% EM-DAT) now that per-type regressors are strong; reconcile coverage constants", 0, SLATE, False),
    ("P3 · Add back-transformed RMSE/R² per type to the regressor validation", 0, SLATE, False),
    ("P4 · Keep class_weight; do NOT revisit SMOTE — data already answered this; next minority lever is better features / focal loss", 0, SLATE, False),
    ("P5 · Fix two train/serve skews (month-mode at serve; heavy-vs-light regressor bundle)", 0, SLATE, False),
    ("P6 · Encapsulate build_features in ONE transformer (it's hand-duplicated across ≥6 files)", 0, SLATE, False),
    ("P7 · Only if a linear/distance model is added: RobustScaler + target/hashing encoders, in-CV", 0, GREY, False),
], top=Inches(1.45), size=15, gap=10)
notes(s, "Ordered by expected impact × confidence. P0 and P1 are the two that change the real picture: "
         "P0 ships a known win; P1 tells us how much of 0.7052 is real. The rest is rigor and hygiene.")

# ── SLIDE 17 — Takeaways ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("Preprocessing, not the algorithm, drove the wins. The biggest single improvement (impact error 111× → 2.7×) is a missing-value + structure decision.", 0, DEEP, True),
    ("Classification and regression need different levers: class weights lifted macro-F1; drop-null + per-type + log1p cut regression error. Scaling helped neither (all trees).", 0, SLATE, False),
    ("Negative results are results: SMOTE, inverse-frequency, geo features, per-class thresholds — all tested, measured, and rejected against a noise floor.", 0, SLATE, False),
    ("Honesty: a validated win not yet shipped (centroid), a leak caught (Dis Mag Scale), and a dedup gap I won't hide — each with a concrete fix.", 0, SLATE, False),
], top=Inches(1.7), size=18, gap=20)
notes(s, "Land the separation of classification vs regression effects, and the discipline narrative. "
         "For an academic audience the honesty (negative results + self-found gaps) is the differentiator.")

# ── SLIDE 18 — Thank you / Q&A ────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(2.55), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(2.7), Inches(11.3), Inches(0.6)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
tf = box(s, Inches(0.9), Inches(3.65), Inches(11.5), Inches(3.1)); p = tf.paragraphs[0]
qa = [
    ("Anticipated questions:", DEEP, True, 16),
    ("Is 0.7052 macro-F1 real?  →  In-distribution upper bound; train/test windows overlap and aren't deduped (P1).", SLATE, False, 14),
    ("Why not oversample?  →  Tested SMOTENC (done right, in-fold): −0.0324 macro; Landslide −0.12. Weights won.", SLATE, False, 14),
    ("Biggest single win?  →  Per-type + drop-null regressors: impact error 111× → 2.7×. Pure preprocessing.", SLATE, False, 14),
    ("Why no scaling?  →  Every model is a tree (XGB/CatBoost/RF) — scale-invariant; scaling is a measured no-op.", SLATE, False, 14),
    ("One thing to ship tomorrow?  →  Centroid coordinate imputation — already validated +0.0167 macro-F1.", SLATE, False, 14),
]
for i, (txt, col, bold, sz) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
    p.space_after = Pt(7)
    r = p.add_run(); r.text = ("•  " if i else "") + txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0) if i else GREEN
    r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(6.95), Inches(11.3), Inches(0.4)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Full write-up: docs/PREPROCESSING_DEEP_DIVE.md   ·   Data: EM-DAT (1900–2021)   ·   safeearth.tech"
r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Keep the contamination/dedup answer ready — most likely probe. Owning it (with the fix) "
         "is the win, same as the centroid 'validated but unshipped' answer.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "SafeEarth_Preprocessing_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides, charts embedded)")
