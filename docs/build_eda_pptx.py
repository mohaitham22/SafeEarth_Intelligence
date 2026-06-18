"""
build_eda_pptx.py — generate an editable PowerPoint (with charts) from the SafeEarth EDA deck.

Run:  py -3.12 docs/build_eda_pptx.py
Out:  docs/SafeEarth_EDA_Presentation.pptx

Native python-pptx (editable text boxes + tables) + embedded PNG charts rendered by
make_eda_charts.py from the real data. Content mirrors docs/EDA_PITCH_DECK.md, which is
derived from docs/EDA_DEEP_DIVE.md. Every number is from this project's own data.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# make_eda_charts lives next to this file
sys.path.insert(0, str(Path(__file__).resolve().parent))
import make_eda_charts  # noqa: E402

ASSETS = Path(__file__).resolve().parent / "eda_assets"

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xDC, 0x26, 0x26)
SLATE  = RGBColor(0x33, 0x41, 0x55)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)
PINK   = RGBColor(0xFE, 0xE2, 0xE2)

W, H = Inches(13.333), Inches(7.5)

# Ensure chart PNGs exist before any slide embeds them.
if not (ASSETS / "05_thesis_scatter.png").exists():
    print("Charts not found — rendering them first...")
    make_eda_charts.generate_all()

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def img(s, name, left, top, width):
    """Add a chart PNG by file name; height auto from aspect ratio."""
    return s.shapes.add_picture(str(ASSETS / name), left, top, width=width)


def title_bar(s, text, kicker=None, color=DEEP):
    band(s, color, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = (RGBColor(0xFE,0xCA,0xCA) if color == RED else GREEN)
        r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    first = True
    for it in items:
        text, level, color, bold = (it + (0, SLATE, False))[:4] if isinstance(it, tuple) else (it, 0, SLATE, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0)); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None, hcolor=DEEP):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val); r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hcolor
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def callout(s, label, text, top=Inches(6.5), color=LIGHT, label_color=GREEN, h=Inches(0.7),
            left=Inches(0.6), width=Inches(12.1)):
    band(s, color, left, top, width, h)
    tf = box(s, left + Inches(0.2), Inches(top.inches + 0.10), Inches(width.inches - 0.4), Inches(h.inches - 0.16))
    p = tf.paragraphs[0]
    if label:
        r = p.add_run(); r.text = label + "  "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = label_color; r.font.name = "Calibri"
    r = p.add_run(); r.text = text
    r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.55), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
for i, (txt, sz, col, bold) in enumerate([
    ("Exploratory Data Analysis — Deep Dive", 26, GREEN, True),
    ("What the data says, and how it shapes (and limits) the models", 19, RGBColor(0xCB,0xD5,0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "EM-DAT (1900–2021, 16,126 events)  ·  XGBoost + CatBoost + RandomForest"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94,0xA3,0xB8); r.font.name = "Calibri"
notes(s, "Open: I'll argue you can predict this model's per-class accuracy from two data "
         "properties alone (correlation 0.85), and that one headline metric is inflated by "
         "a data leak I'll show you. This is a data-science talk, not a product demo.")

# ── SLIDE 2 — Roadmap ─────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Roadmap")
bullets(s, [
    ("The data — what one row is, and why that single fact governs everything", 0, SLATE, False),
    ("Three data-quality stories — imbalance, skew, class-conditional missingness", 0, SLATE, False),
    ("The thesis — data quality predicts per-class accuracy (r = 0.85)", 0, GREEN, True),
    ("What we did — imbalance handling & the regressor fix (with the negative results)", 0, SLATE, False),
    ("Two problems I found — train/test contamination, and a median-aggregation bug", 0, RED, True),
    ("Prioritised gaps — what to fix next, highest-ceiling first", 0, SLATE, False),
], top=Inches(1.7), size=20, gap=18)
notes(s, "Tell them the talk has a critical spine on purpose — negative results and "
         "self-found bugs are the evidence of real analysis. Every number is from this "
         "project's own data, extracted or computed by replaying the feature pipeline.")

# ── SLIDE 3 — Dataset & unit of analysis ──────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Dataset & the Unit of Analysis")
bullets(s, [
    ("Source: EM-DAT International Disaster Database — the standard global registry", 0, SLATE, False),
    ("One row = one declared disaster event, curated AFTER the fact", 0, SLATE, False),
    ("Impact columns (deaths, affected, damage) are POST-EVENT accounting totals —", 0, AMBER, True),
    ("knowable only after the event we claim to predict (drives the leakage analysis)", 1, SLATE, False),
    ("Scope: global — ~225 countries, 23 regions, 5 continents, 1900–2021", 0, SLATE, False),
], top=Inches(1.45), size=18, gap=9)
table(s, [
    ["File", "Raw", "Modeled (8 types)", "Dropped"],
    ["Train (1900–2021)", "16,126 × 45", "14,476", "1,650 (10.2%)"],
    ["Test (1970–2021)", "14,644 × 47", "13,070", "—"],
], Inches(0.6), Inches(4.7), Inches(12.1),
   col_widths=[Inches(3.6), Inches(2.8), Inches(3.4), Inches(2.3)], fsize=15)
callout(s, "Why it matters:", "post-event totals are the half of the data the model is NOT allowed to know at prediction time.",
        top=Inches(6.45), h=Inches(0.7))

# ── SLIDE 4 — Provenance ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Where the EDA Actually Lives (provenance)")
bullets(s, [
    ("notebooks/01_eda.ipynb is a one-cell stub — ZERO executed analysis", 0, RED, True),
    ("The real EDA is distributed across four artifacts:", 0, DEEP, True),
], top=Inches(1.45), size=18, gap=8)
table(s, [
    ["Artifact", "What it actually contains"],
    ["scripts/inspect_data.py", "Inventory / missingness / cardinality (data-dictionary substitute)"],
    ["streamlit_eda/", "Interactive EDA app — the quality→accuracy thesis"],
    ["SafeEarth_ML_Report.ipynb", "62 cells; LOADS the production .pkl artifacts; reproduces every metric"],
    ["streamlit_eda/experiments/ + metrics/", "Isolated experiments & the decision paper-trail"],
], Inches(0.6), Inches(2.7), Inches(12.1),
   col_widths=[Inches(4.0), Inches(8.1)], fsize=14)
callout(s, "Note:", "figures tagged “computed” were produced fresh by replaying engineer_features on the raw CSV.",
        top=Inches(6.3), h=Inches(0.7))
notes(s, "Faculty respect provenance. Admitting the stub notebook up front buys "
         "credibility for everything after it.")

# ── SLIDE 5 — 45 cols → 16 features ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "From 45 Columns to 16 Features")
bullets(s, [
    ("45 raw columns → a 16-feature model input (27 object / 14 float / 4 int dtypes)", 0, SLATE, False),
    ("Most “numeric” impact + coordinate fields arrive as text → need coercion", 1, GREY, False),
    ("Predictive: geography (lat/lon/encoded), time (decade, month sin/cos), magnitude + has_magnitude flag, regional frequency", 0, SLATE, False),
    ("Deliberately EXCLUDED to avoid leakage:", 0, AMBER, True),
    ("impact columns · Dis Mag Scale · Disaster Subtype · Associated Dis (all post-event or target-revealing)", 1, SLATE, False),
    ("day_offset is constant 0 in training → a dead feature at fit time (revisited Slide 13)", 0, GREY, False),
], top=Inches(1.6), size=18, gap=13)
notes(s, "The 'I know which columns are features vs leaks vs metadata' slide. Keep brisk.")

# ── SLIDE 6 — Class imbalance (+ chart) ───────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Class Imbalance — the 21:1 Problem")
bullets(s, [
    ("Imbalance ≈ 21:1 (Flood : Volcanic)", 0, DEEP, True),
    ("Flood + Storm = 69% of all rows", 0, SLATE, False),
    ("Constant-“Flood” baseline ≈ 38% accuracy", 0, SLATE, False),
    ("→ accuracy / weighted-F1 are useless here; macro-F1 is the binding metric", 1, AMBER, True),
    ("Subtlety: hardest ≠ smallest —", 0, DEEP, True),
    ("Volcanic (265) BEATS Landslide (776) on F1 (0.67 vs 0.48). Size is the weaker driver.", 1, SLATE, False),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.1), size=17, gap=12)
img(s, "01_imbalance.png", Inches(5.9), Inches(2.0), Inches(7.1))
notes(s, "Set up the punchline: it's not row count, it's signal — next two slides prove it.")

# ── SLIDE 7 — Skew → median (+ chart) ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Impact Is Skewed 100×–1,400× → Median, Never Mean")
table(s, [
    ["Target", "mean", "median", "ratio"],
    ["Deaths", "1,576", "6", "263×"],
    ["Affected", "558,521", "400", "1,396×"],
    ["Injured", "531", "0", "∞"],
    ["Damage ('000$)", "262,559", "0", "∞"],
], Inches(0.55), Inches(1.6), Inches(5.3),
   col_widths=[Inches(2.0), Inches(1.4), Inches(1.0), Inches(0.9)], fsize=13)
img(s, "02_skew.png", Inches(6.1), Inches(1.5), Inches(6.9))
callout(s, "Median 0 = missingness, not zeros:", "injuries/damage are 26%/36% present, coerced to 0 — the time-bomb the regressors step on (Slide 14).",
        top=Inches(5.05), h=Inches(0.75), color=PINK, label_color=RED)
notes(s, "Two ideas: skew justifies median; and the zero-spike is masked missingness.")

# ── SLIDE 8 — Missingness (+ chart) ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Missingness Is the Headline Data-Quality Story")
bullets(s, [
    ("Reliable: Country 100%, Start Month 98%, Affected 71%, Deaths 70%", 0, GREEN, True),
    ("Sparse: Damage 36%, Magnitude 34%, Injuries 26%, Coords ~19%, Insured 8%", 0, RED, True),
    ("Most missingness is MNAR (recorded when severe), not MCAR", 0, DEEP, True),
    ("→ biases impact models low; imputation is not innocent", 1, SLATE, False),
    ("Only ~19% have real coordinates — 81% imputed to centroids (Slide 12)", 0, AMBER, True),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.2), size=16, gap=13)
img(s, "03_missingness.png", Inches(6.0), Inches(1.7), Inches(7.0))

# ── SLIDE 9 — Magnitude coverage (+ chart) ────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Magnitude Coverage Is Class-Conditional (key insight)")
bullets(s, [
    ("Magnitude carries the MOST information about type (mutual information 0.503 — top of 16 features)", 0, SLATE, False),
    ("…but recorded on different scales (Richter / Kph / Km² / °C) and present 94% for Earthquake, ~1% for Landslide/Volcanic", 1, SLATE, False),
    ("The engineered has_magnitude flag turns MNAR-missingness into the strongest LINEAR separator (ANOVA F = 655)", 0, GREEN, True),
], left=Inches(0.55), top=Inches(1.5), width=Inches(5.2), size=16, gap=14)
img(s, "04_magcoverage.png", Inches(6.0), Inches(1.7), Inches(7.1))
callout(s, "Hold this chart —", "the next slide shows it predicts the model's per-class accuracy almost perfectly.",
        top=Inches(6.5), h=Inches(0.65), label_color=AMBER)

# ── SLIDE 10 — THE THESIS (chart-dominant) ────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Thesis: Data Quality Predicts Accuracy", kicker="The Money Slide")
img(s, "05_thesis_scatter.png", Inches(2.55), Inches(1.32), Inches(8.2))
callout(s, "The takeaway:", "signal availability — not algorithm or row count (corr with size only 0.46) — sets the per-class ceiling. The model is only as strong as its rarest feature.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)
notes(s, "The money slide. Earthquake wins on a near-complete distinct Richter signal; "
         "Landslide/Volcanic lose because their feature is absent and they overlap "
         "floods/storms spatially. This frames why resampling fails and what to fix.")

# ── SLIDE 11 — Bivariate MI / ANOVA ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Feature → Target Separability (bivariate evidence)")
table(s, [
    ["Rank by Mutual Information", "MI", "Rank by ANOVA F (linear)", "F"],
    ["dis_mag_value", "0.503", "has_magnitude", "655"],
    ["latitude", "0.430", "abs_latitude", "182"],
    ["abs_latitude", "0.420", "lon_cos", "176"],
    ["longitude / lon_sin", "~0.42", "decade", "153"],
    ["country_enc", "0.363", "continent_enc", "118"],
    ["day_offset", "0.000", "day_offset", "NaN"],
], Inches(0.6), Inches(1.5), Inches(12.1),
   col_widths=[Inches(3.8),Inches(1.6),Inches(4.7),Inches(2.0)], fsize=13.5)
bullets(s, [
    ("Magnitude dominates non-linear information; its presence FLAG dominates linear separation", 0, SLATE, False),
    ("Geography is a strong secondary; season is weak-but-real; day_offset is dead at training (MI 0 / F NaN)", 0, SLATE, False),
], top=Inches(5.2), height=Inches(1.9), size=17, gap=12)
notes(s, "Pre-empts 'did you check relevance beyond importances?' — yes: MI and F-test.")

# ── SLIDE 12 — Geospatial ─────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Geography — 81% Imputed, and It Matters")
bullets(s, [
    ("Only 18.8% of events have raw coordinates → 81% imputed country→continent→global median", 0, SLATE, False),
    ("Coverage is class-conditional: Earthquake 97%, but Extreme-temp / Wildfire / Drought = 0% raw coords", 0, SLATE, False),
    ("→ pinned to one country centroid, collapsing within-country variance", 1, GREY, False),
    ("Imputation is OUTCOME-relevant, not cosmetic: centroid imputation moved holdout macro-F1", 0, DEEP, True),
    ("0.7285 → 0.7452 (+0.0167 > CV-std 0.008 → KEEP), filling 11,707 of 11,822 missing-coord rows", 1, GREEN, True),
    ("Reporting drift: floods 524 (1980s) → 1,725 (2000s) = 3.3× — climate AND reporting, confounded; decade bakes the era-bias in", 0, AMBER, True),
], top=Inches(1.55), size=17, gap=12)

# ── SLIDE 13 — Imbalance handling (+ chart) ───────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Imbalance Handling — Weights Beat Resampling")
bullets(s, [
    ("Every resampling strategy REGRESSED holdout macro-F1 —", 0, AMBER, True),
    ("the bottleneck is missing signal, not missing rows", 1, SLATE, False),
    ("SMOTE is methodologically invalid here:", 0, RED, True),
    ("5 of 16 features are encoded categoricals — interpolating country_enc(Egypt)↔(Peru) invents a country that doesn't exist", 1, SLATE, False),
    ("Chosen weights up-weight hard classes: Drought 4.0, Landslide/Volcanic 3.0, Wildfire 2.5", 0, GREEN, True),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.4), size=15.5, gap=11)
img(s, "06_strategies.png", Inches(6.2), Inches(2.0), Inches(6.8))
notes(s, "Research-maturity slide: a tested-and-rejected result WITH the reason it fails.")

# ── SLIDE 14 — Regressor fix (+ chart) ────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Regressor Fix — Per-Type + Drop-Null")
bullets(s, [
    ("Original regressors were type-blind AND trained on fake zeros (median injuries/damage = 0)", 0, SLATE, False),
    ("Two independent wins: add the disaster-TYPE axis + DROP fake zeros", 0, GREEN, True),
    ("Validated on an honest holdout (observed-target rows only)", 0, SLATE, False),
    ("Exposed an evaluation trap: median-filling the HOLDOUT too makes a broken model look great (injuries 0.45 vs honest 1.28)", 0, RED, True),
], left=Inches(0.55), top=Inches(1.55), width=Inches(5.3), size=15.5, gap=13)
img(s, "07_regressors.png", Inches(6.1), Inches(2.0), Inches(6.9))
notes(s, "Strongest single piece of analysis in the repo. The 'evaluation trap' is a great "
         "teaching moment — how you fill the test set can manufacture a false score.")

# ── SLIDE 15 — Contamination (red theme + chart) ──────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Train/Test Contamination", kicker="Problem #1", color=RED)
bullets(s, [
    ("The “holdout” 1970–2021 file is a STRICT SUBSET of the 1900–2021 train file — same export, windowed twice", 0, SLATE, False),
    ("Computed: 100.0% of the 14,644 test rows share an EM-DAT (Year, Seq) event id with a train row", 0, RED, True),
    ("→ the model sees every evaluation event during fit()", 1, SLATE, False),
    ("So Macro-F1 0.7052 is IN-SAMPLE — an upper bound, not generalization", 0, AMBER, True),
    ("It voids the absolute numbers, not the relative decisions (those used CV / honest holdout)", 0, SLATE, False),
], left=Inches(0.55), top=Inches(1.5), width=Inches(5.7), size=15.5, gap=12)
img(s, "08_contamination.png", Inches(6.4), Inches(1.9), Inches(6.6))
callout(s, "Fix (Slide 19 #1):", "re-split by (Year, Seq) or by time, then re-evaluate; expect macro-F1 to fall toward ~0.6.",
        top=Inches(6.55), h=Inches(0.62), color=PINK, label_color=RED)
notes(s, "Deliver calmly and confidently — finding it is a strength, not a confession. "
         "This is the most likely probe; owning it is the win.")

# ── SLIDE 16 — Median bug ─────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Median Fails on Small Cells (a live analytics bug)", kicker="Problem #2", color=RED)
bullets(s, [
    ("The “median, never mean” rule protects GLOBAL aggregates — but it's applied per-cell with no minimum-support guard", 0, SLATE, False),
    ("continent_stats.json reports Europe drought deaths = 600,001", 0, RED, True),
    ("That “median” is over only ~2 observed European-drought death rows, dominated by the 1921 Soviet famine (1.2 M deaths)", 1, SLATE, False),
    ("Median is robust to outliers — NOT to tiny samples", 0, AMBER, True),
    ("Fix: apply the same n ≥ 5 guard the EM-DAT impact lookup already uses — but doesn't apply to continent analytics", 0, GREEN, True),
], top=Inches(1.6), size=18, gap=16)
notes(s, "Short slide, big credibility: 'I audit my own outputs.'")

# ── SLIDE 17 — EDA → model matrix ─────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "EDA → Model Impact, in One Matrix")
table(s, [
    ["Insight", "Effect", "Mitigation"],
    ["21:1 imbalance", "majority bias → macro-F1 the honest read", "class weights; report macro-F1"],
    ["class-conditional magnitude (r=0.85)", "sets the per-class ceiling", "has_magnitude flag; BACKFILL magnitude"],
    ["100×–1400× impact skew", "breaks linear/distance models; trees survive", "log1p targets; tree ensemble"],
    ["fake-zero, type-blind regressors", "error up to 111×", "per-type + drop-null → 2.7×"],
    ["VIF 12–21, r up to 0.94", "harmless for trees, fatal for LR/SVM", "drop day_offset; prune if linear"],
    ["outliers = genuine catastrophes", "clipping destroys the key signal", "log1p + per-type P99 display cap"],
    ["train/test contamination", "inflates every headline", "re-split, re-evaluate"],
], Inches(0.5), Inches(1.4), Inches(12.3),
   col_widths=[Inches(4.0),Inches(4.6),Inches(3.7)], fsize=12.5)
callout(s, "The recurring lesson:", "trees were the right call — collinear, mixed-unit, outlier-dominated, high-cardinality data breaks LR/SVM/SVR/KNN; trees absorb it.",
        top=Inches(6.55), h=Inches(0.65), label_color=GREEN)

# ── SLIDE 18 — Model scoreboard (+ 2 charts) ──────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Model Scoreboard — What the Data Permits")
table(s, [
    ["Model", "Task", "Headline", "Read"],
    ["1 · Disaster Type", "8-class", "Macro-F1 0.705 · Wtd 0.759 · AUC 0.962", "strong; EQ 0.98, Landslide 0.48"],
    ["2 · Impact", "4 regressors", "mean log-R² 0.377 · orig-R² ≈ 0", "location signal; type from EM-DAT medians"],
    ["3 · Risk level", "4-class", "Macro-F1 0.534 vs Wtd 0.946", "97% Low; recall 0.92 / precision 0.24"],
], Inches(0.5), Inches(1.35), Inches(12.3),
   col_widths=[Inches(2.4),Inches(1.8),Inches(4.2),Inches(3.9)], fsize=12)
img(s, "09_perclass_f1.png", Inches(0.45), Inches(3.35), Inches(6.35))
img(s, "10_risklevel.png", Inches(6.95), Inches(3.45), Inches(6.1))
notes(s, "Left: Model 1 per-class F1 (macro vs weighted gap). Right: Model 3's 97%-Low "
         "imbalance — why risk level is COMPUTED, not classified. The data, not the "
         "algorithm, dictates each model's shape.")

# ── SLIDE 19 — Recommendations ────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Recommendations — Highest Ceiling First")
bullets(s, [
    ("1. Fix the split — re-evaluate without contamination. Group by (Year, Seq) or split by time. The only number that matters first.", 0, RED, True),
    ("2. Backfill magnitude for Landslide / Volcanic / Drought (VEI, volume index, SPI) — directly attacks the macro-F1 floor (r=0.85)", 0, GREEN, True),
    ("3. Significance & power: bootstrap CIs on per-class F1; flag underpowered classes (test Critical n=27)", 0, SLATE, False),
    ("4. Distribution-match / drift audit (PSI/KS) after re-splitting; quantify the decade era-bias", 0, SLATE, False),
    ("5. Multicollinearity hygiene: compute VIF (done); drop the dead day_offset", 0, SLATE, False),
    ("6. Real skew/kurtosis + dip-test for bimodality (median can hide bimodal impact)", 0, SLATE, False),
    ("7. Calibration: the class-weighted ensemble is likely miscalibrated, yet thresholds drive alerts → reliability diagrams + isotonic", 0, SLATE, False),
    ("8. Grouped / temporal CV for spatial-overfitting; fix the small-cell median bug", 0, SLATE, False),
], top=Inches(1.5), size=15.5, gap=11)
notes(s, "Ordered by expected impact. #1 and #2 are the only two that move the real "
         "ceiling; the rest are rigor.")

# ── SLIDE 20 — Takeaways ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("The data explains the model. Per-class accuracy is predicted by magnitude coverage at r = 0.85 — signal availability is the ceiling, not the algorithm.", 0, DEEP, True),
    ("Median, never mean — impact is skewed 100×–1,400× — but median needs a minimum-support guard.", 0, SLATE, False),
    ("Negative results are results: resampling/SMOTE tested and rejected with reasons; the regressor fix came from a controlled experiment.", 0, SLATE, False),
    ("I audited my own pipeline and found a contaminated evaluation and a live aggregation bug — with concrete, prioritised fixes.", 0, SLATE, False),
], top=Inches(1.9), size=19, gap=22)
notes(s, "End on rigor and honesty — the differentiator for an academic audience.")

# ── SLIDE 21 — Thank you / Q&A ────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(2.6), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.55), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(2.75), Inches(11.3), Inches(0.6)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
tf = box(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(3.2)); p = tf.paragraphs[0]
qa = [
    ("Anticipated questions:", DEEP, True, 16),
    ("Is 0.705 real?  →  No — upper bound; contaminated holdout. Honest CV ≈ 0.59.", SLATE, False, 14),
    ("Why not oversample?  →  Tested; SMOTE −0.04 macro & invalid on encoded categoricals.", SLATE, False, 14),
    ("Why trees over a linear model?  →  VIF 21, mixed-unit magnitude, 100× outliers, 225-level categoricals all break LR/SVM/KNN.", SLATE, False, 14),
    ("Biggest single improvement?  →  Backfill magnitude for the 3 signal-poor classes (r = 0.85).", SLATE, False, 14),
    ("Median vs mean concretely?  →  Flood deaths mean 1,735 vs median 16; affected mean/median = 1,396×.", SLATE, False, 14),
]
for i, (txt, col, bold, sz) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
    p.space_after = Pt(7)
    r = p.add_run(); r.text = ("•  " if i else "") + txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = RGBColor(0xE2,0xE8,0xF0) if i else GREEN
    r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(6.95), Inches(11.3), Inches(0.4)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Full write-up: docs/EDA_DEEP_DIVE.md   ·   Data: EM-DAT (1900–2021)   ·   safeearth.tech"
r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x94,0xA3,0xB8); r.font.name = "Calibri"
notes(s, "Keep the contamination answer ready — most likely probe; owning it is the win.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "SafeEarth_EDA_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides, charts embedded)")
