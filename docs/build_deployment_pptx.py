"""
build_deployment_pptx.py — generate an editable PowerPoint from the SafeEarth
DEPLOYMENT pitch deck.

Run:  py -3.12 docs/build_deployment_pptx.py
Out:  docs/SafeEarth_Deployment_Presentation.pptx

Native python-pptx (editable text boxes + tables), no browser needed.
Companion to build_pptx.py (backend, green) and build_frontend_pptx.py (frontend,
blue) — same palette + helpers, plus an INDIGO accent to signal the
"infrastructure / deployment chapter". Content sourced from DEPLOYMENT_EXPLAINED.md
and DEPLOYMENT_PITCH_DECK.md, grounded in the repo's real config files
(render.yaml, vercel.json, backend/scripts/render_build.sh, n8n/, .env.example).
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette (matched to the other decks, + INDIGO for the deployment flavor) ────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)   # deep green (brand — title bars)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)   # accent green (positive / live)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # accent amber (honest trade-offs)
INDIGO = RGBColor(0x4F, 0x46, 0xE5)   # accent indigo (deployment / infra)
SLATE  = RGBColor(0x33, 0x41, 0x55)   # body text
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # light panel
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)

W, H = Inches(13.333), Inches(7.5)     # 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def title_bar(s, text, kicker=None):
    band(s, DEEP, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = INDIGO; r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(28); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    tf.word_wrap = True
    first = True
    for it in items:
        # tuple (text, level, color, bold)
        if isinstance(it, tuple):
            text, level, color, bold = (it + (0, SLATE, False))[:4]
        else:
            text, level, color, bold = it, 0, SLATE, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0))
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = DEEP
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def code(s, lines, left, top, width, height, fsize=13):
    band(s, RGBColor(0x0F, 0x17, 0x2A), left, top, width, height)
    tf = box(s, left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = ln
        r.font.name = "Consolas"; r.font.size = Pt(fsize)
        r.font.color.rgb = RGBColor(0x7E, 0xE7, 0xB0)
    return tf


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def node(s, l, t, w, h, text, fill, fg=WHITE, fs=13):
    sh = band(s, fill, l, t, w, h)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = fg; r.font.name = "Calibri"
    return sh


# ── Component slide helper (mirrors the backend / frontend decks) ──────────────
def component(num, title, what, body_items, code_lines=None, chose=None,
              honest=None, note=None):
    s = slide(); bg(s, WHITE)
    title_bar(s, title, kicker=f"Component {num}")
    top = Inches(1.4)
    if what:
        tf = box(s, Inches(0.6), top, Inches(12.1), Inches(0.7))
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = what
        r.font.size = Pt(16); r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
        top = Inches(2.12)
    if code_lines:
        code(s, code_lines, Inches(0.6), top, Inches(12.1), Inches(0.4 * len(code_lines) + 0.28), fsize=13)
        top = Inches(top.inches + 0.4 * len(code_lines) + 0.5)
    bullets(s, body_items, top=top, size=16, gap=8)
    y = 6.5
    if honest:
        band(s, RGBColor(0xFE, 0xF3, 0xC7), Inches(0.6), Inches(y - 0.78), Inches(12.1), Inches(0.66))
        tf = box(s, Inches(0.8), Inches(y - 0.72), Inches(11.7), Inches(0.55)); p = tf.paragraphs[0]
        r = p.add_run(); r.text = "Honest note:  "
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = AMBER; r.font.name = "Calibri"
        r = p.add_run(); r.text = honest
        r.font.size = Pt(13); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    if chose:
        band(s, LIGHT, Inches(0.6), Inches(y), Inches(12.1), Inches(0.66))
        tf = box(s, Inches(0.8), Inches(y + 0.08), Inches(11.7), Inches(0.5)); p = tf.paragraphs[0]
        r = p.add_run(); r.text = "Chose over: "
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = INDIGO; r.font.name = "Calibri"
        r = p.add_run(); r.text = chose
        r.font.size = Pt(14); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    if note:
        notes(s, note)
    return s


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, INDIGO, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.65), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
for i, (txt, sz, col, bold) in enumerate([
    ("Getting an ML Model Online", 27, INDIGO, True),
    ("Deployment & Infrastructure", 21, RGBColor(0xCB, 0xD5, 0xE1), False),
    ("How a notebook becomes a service that runs 24/7 for the internet", 16, RGBColor(0x94, 0xA3, 0xB8), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.5), Inches(11.3), Inches(1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "Vercel · Render · Neon · HuggingFace · Groq · n8n · Docker"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "I trained a disaster-prediction model and wrapped it in an API. This talk is the "
         "part nobody sees in a notebook demo: the seven cloud services, the secrets, the "
         "build scripts, and the trade-offs that take it from 'runs on my laptop' to 'live "
         "on the internet'. It's the difference between a result and a product.")

# ── SLIDE 2 — The Problem ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Problem")
bullets(s, [
    ("\"It works on my laptop\" is not a product.", -1, DEEP, True),
    ("In a notebook, I am the runtime: I open it, run the cells, watch the output — it stops when I close the lid.", 0, SLATE, False),
    ("localhost:3000 literally means \"this computer, right here\". Nobody else on Earth can reach it.", 0, SLATE, False),
    ("A model that only runs while I'm watching helps exactly one person — me.", 0, AMBER, True),
], top=Inches(1.6), size=20, gap=16)
band(s, LIGHT, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.4))
tf = box(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Goal: "
r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = INDIGO; r.font.name = "Calibri"
r = p.add_run(); r.text = ("make the prediction reachable by anyone with a browser, always on, safe, and ship "
                           "new versions without breaking the live one.")
r.font.size = Pt(19); r.font.color.rgb = SLATE; r.font.name = "Calibri"
notes(s, "Land the mental shift: the job changes from 'produce the right answer once' to 'keep "
         "producing it forever, on a machine I don't control, for people I'll never meet'. "
         "Everything after this slide solves that.")

# ── SLIDE 3 — What Deployment Even Is ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What Deployment Even Is")
tf = box(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Analogy: cooking at home vs. opening a restaurant."
r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
table(s, [
    ["Cooking at home (notebook)", "Opening a restaurant (deployment)"],
    ["Your private kitchen, set up your way", "A building on a public street = a HOST"],
    ["You eat what you cook", "Strangers walk in and order"],
    ["Closes when you sleep", "Open 24/7, no human watching"],
    ["Burnt toast — only you know", "A crash is one EVERYONE sees"],
    ["New dish whenever", "A safe way to roll out new dishes = CI/CD"],
], Inches(0.6), Inches(1.95), Inches(12.1),
   col_widths=[Inches(5.6), Inches(6.5)], fsize=15)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Host = a computer someone else runs, always on, public.   Production = the real live version;  dev = the private rehearsal."
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = SLATE; r.font.name = "Calibri"
notes(s, "This restaurant analogy carries the whole talk — I reuse kitchen / dining room / "
         "utilities / new dishes for every piece. Localhost = home kitchen. Deployment = "
         "restaurant: building (host), utilities (DB, config), safe rollout (CI/CD).")

# ── SLIDE 4 — The Deployment Topology ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Deployment Topology")
# user
node(s, Inches(5.35), Inches(1.4), Inches(2.6), Inches(0.7), "User's browser", SLATE)
# hosts row
node(s, Inches(2.5), Inches(2.5), Inches(3.0), Inches(0.9), "Vercel\nFrontend (Next.js)\nsafeearth.tech", INDIGO, fs=12)
node(s, Inches(7.8), Inches(2.5), Inches(3.0), Inches(0.9), "Render (free, Oregon)\nFastAPI · uvicorn\n/api/v1/*", DEEP, fs=12)
# backend dependencies row
node(s, Inches(0.6), Inches(4.05), Inches(2.7), Inches(0.85), "Neon\nPostgreSQL\n(managed DB)", SLATE, fs=11)
node(s, Inches(3.45), Inches(4.05), Inches(2.7), Inches(0.85), "HuggingFace\n.pkl models\n(pull at startup)", SLATE, fs=11)
node(s, Inches(6.3), Inches(4.05), Inches(2.7), Inches(0.85), "Groq\nLLM API\n(RAG)", AMBER, DEEP, fs=11)
node(s, Inches(9.15), Inches(4.05), Inches(2.7), Inches(0.85), "Gmail / Resend\nemail", SLATE, fs=11)
# n8n
node(s, Inches(7.8), Inches(5.25), Inches(3.0), Inches(0.8), "n8n (scheduled)\nMon 08:00 UTC →\nPOST /alerts/dispatch", AMBER, DEEP, fs=11)
# captions
tf = box(s, Inches(2.5), Inches(2.05), Inches(3.0), Inches(0.4)); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "loads page"; r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GREY; r.font.name = "Calibri"
tf = box(s, Inches(7.8), Inches(2.05), Inches(3.0), Inches(0.4)); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "API calls (JSON / HTTPS)"; r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GREY; r.font.name = "Calibri"
tf = box(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.7)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Seven hosted pieces — each a different company's computer, in a different data center, free-tiered separately, connected only by network calls."
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = SLATE; r.font.name = "Calibri"
notes(s, "Key insight for a notebook person: opposite of one process where everything shares "
         "memory. Here model / data store / API / UI are four different machines in four data "
         "centers. That separation lets each scale and be replaced independently — and it's "
         "also why deployment feels complicated at first.")

# ── SLIDE 5 — Stack at a Glance ────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Deployment Stack at a Glance")
table(s, [
    ["Platform / Tool", "What it hosts or does", "Config in the repo"],
    ["Vercel", "The Next.js frontend (the website)", "vercel.json · next.config.js"],
    ["Render", "The FastAPI backend (the API) + n8n", "render.yaml · render_build.sh"],
    ["Neon", "Managed / serverless PostgreSQL", "DATABASE_URL · alembic/"],
    ["HuggingFace Hub", "Trained .pkl models, pulled at startup", "ml/predictor.py · HUGGINGFACE_*"],
    ["Groq", "External LLM API (recommendations)", "rag/recommender.py · GROQ_API_KEY"],
    ["chapters.json", "RAG knowledge base — built at deploy", "rag/extract_chapters.py"],
    ["n8n", "Self-hosted scheduled alert dispatch", "n8n/weekly_dispatch.json"],
    ["Alembic", "Runs DB migrations on each deploy", "alembic.ini (in render_build.sh)"],
    ["Docker (local only)", "Local Postgres for dev — NOT the host", "docker-compose.yml"],
    ["Git auto-deploy", "Push → both platforms rebuild", "dashboards (no GitHub Actions)"],
], Inches(0.5), Inches(1.35), Inches(12.3),
   col_widths=[Inches(2.9), Inches(5.0), Inches(4.4)], fsize=12)
notes(s, "Don't read every row. Two takeaways: everything is declared as a file in the repo "
         "(infrastructure-as-code), and everything is free tier. That constraint shaped every "
         "later decision.")

# ── SLIDE 6 — Component 1: Vercel ──────────────────────────────────────────────
component(
    1, "Frontend Hosting (Vercel)",
    "A host for your website. A PaaS — like Colab/Deepnote, but for sites: hand it code, it builds + serves + HTTPS-es.",
    [
        ("vercel.json → { rootDirectory: \"frontend\", framework: \"nextjs\" }", 0, SLATE, False),
        ("Vercel runs npm install → npm run build (next build) → serves .next", 0, SLATE, False),
        ("Why: Vercel makes Next.js — zero-config, free Hobby tier, global CDN, auto-HTTPS", 0, SLATE, False),
    ],
    code_lines=['{ "rootDirectory": "frontend", "framework": "nextjs" }'],
    chose="Netlify (Next.js server features lag), Cloudflare Pages (edge-only runtime), self-host (you babysit a Linux box + SSL).",
    honest=".env.production points at the raw Render URL (safeearth-backend.onrender.com), not the intended api.safeearth.tech — one line to swap.",
    note="Free Hobby tier is non-commercial with bandwidth caps — perfect for a portfolio. "
         "The 'push and it's live with HTTPS in 90s' experience is the selling point.",
)

# ── SLIDE 7 — Component 2: Render ──────────────────────────────────────────────
component(
    2, "Backend Hosting (Render)",
    "A host for your API — the Python program that does the real work. Render runs it for you and restarts it if it crashes.",
    [
        ("startCommand: uvicorn main:app --host 0.0.0.0 --port $PORT", 0, SLATE, False),
        ("--host 0.0.0.0 = accept connections from anywhere (laptop default is localhost-only)", 1, SLATE, False),
        ("--port $PORT = the host assigns the port; READ it, don't hardcode 8000", 1, SLATE, False),
        ("buildCommand runs render_build.sh: install · gen JSON · chapters.json · alembic upgrade head · seed", 0, SLATE, False),
        ("Free tier SLEEPS after ~15 min idle → 30-60s cold start; UptimeRobot pings /health every 14 min", 0, AMBER, True),
    ],
    chose="Railway (time-limited credits), Fly.io (Docker curve), AWS (DevOps wall), bare VPS (you own the OS + SSL).",
    note="Same uvicorn command as the notebook, two new flags, and the host keeps it alive. "
         "The 512 MB RAM limit is the recurring villain — flag it here, pay it off on the RAG slide.",
)

# ── SLIDE 8 — Component 3: Neon ─────────────────────────────────────────────────
component(
    3, "Managed Database (Neon Postgres)",
    "A live, concurrent data store — not a dead CSV. Managed = someone else runs the DB server; you just get a connection string.",
    [
        ("One env var: DATABASE_URL=postgresql+asyncpg://user:pass@host:5432/safeearth", 0, SLATE, False),
        ("In render.yaml it's sync:false → the real string is pasted in the dashboard, never in git", 0, SLATE, False),
        ("Local dev uses the SAME Postgres 15 via Docker (docker-compose.yml)", 0, SLATE, False),
        ("Why a DB not a CSV: a CSV can't do safe concurrent writes and dies on host restart", 0, SLATE, False),
    ],
    chose="Supabase (bundles auth/storage we don't need), AWS RDS (no lasting free tier), self-managed Postgres (you own backups + the 2 a.m. page).",
    note="DS framing: a CSV is a dead file you load into pandas; a DB is a live queryable store "
         "many users hit at once, with guarantees two simultaneous signups won't collide. Same "
         "Postgres locally and in prod = no surprises.",
)

# ── SLIDE 9 — Component 4: HuggingFace ─────────────────────────────────────────
component(
    4, "Model Hosting (HuggingFace Hub)",
    "\"GitHub for models\" — a store built for large model files. The .pkl's live here, NOT in the git repo.",
    [
        ("Why not git: history is forever + binaries diff poorly; a 50-100 MB pkl bloats every clone", 0, SLATE, False),
        ("At startup the app downloads the 2 pkl's via hf_hub_download (HUGGINGFACE_REPO_ID/TOKEN)", 0, SLATE, False),
        ("Memory trick: the 109 MB SHAP explainer is REBUILT from the in-RAM model, not downloaded (~300 MB saved)", 0, AMBER, True),
    ],
    code_lines=[
        'hf_hub_download(repo_id=repo_id, filename=filename,',
        '                local_dir=str(saved_models_dir),',
        '                token=token or None)   # None = public repo',
    ],
    chose="bundling in git/image (repo bloat, LFS pain), S3 / R2 / GCS (the 'real' answer, but you manage buckets + IAM keys).",
    honest="The HF repo must be PUBLIC, or the startup download 401s on Render. Keeping it public sidesteps that.",
    note="DS analogy: a model registry / artifact store — the thing you'd reach for instead of "
         "a .pkl sitting in the repo.",
)

# ── SLIDE 10 — Component 5: Groq + the pivot ───────────────────────────────────
component(
    5, "LLM API (Groq) + an Honest Pivot",
    "An external API dependency — a feature you rent. Recommendations are written by Groq's llama-3.1-8b-instant; the API key is a secret.",
    [
        ("Degrade-not-fail: Groq down/empty/malformed → fall back to a seeded recommendations DB table", 0, GREEN, True),
        ("The pivot — where the 'vector store' went:", 0, DEEP, True),
        ("Original: ChromaDB + sentence-transformers embeddings → pulls in PyTorch (~2 GB) → OOM on 512 MB", 1, SLATE, False),
        ("Now: build-time extract_chapters.py (PyMuPDF) → chapters.json; runtime = dict lookup → Groq", 1, SLATE, False),
        ("No embeddings, no vector DB, nothing extra to host", 1, SLATE, False),
    ],
    chose="OpenAI / Anthropic (paid + slower for tiny tasks), self-hosted LLM (needs a GPU — impossible at 512 MB), full vector RAG (the OOM we just hit).",
    note="The slide professors probe. Lesson = engineering maturity: match the solution to a "
         "real constraint and degrade gracefully. Whole-chapter retrieval is less precise than "
         "vectors, but for a fixed 8-category PDF it's fine — and it's the difference between "
         "'RAG runs on free tier' and 'the server won't boot'.",
)

# ── SLIDE 11 — Component 6: n8n ─────────────────────────────────────────────────
component(
    6, "Automation / Scheduled Jobs (n8n)",
    "Work that runs on a timer, not when a user clicks. n8n = \"visual cron with a UI\", running as its own always-on service.",
    [
        ("weekly_dispatch.json: [Schedule cron \"0 8 * * 1\"] → [HTTP POST /api/v1/alerts/dispatch]", 0, SLATE, False),
        ("Architecture rule: n8n is ONLY a trigger — FastAPI does all fan-out + email. n8n just rings the doorbell.", 0, AMBER, True),
        ("X-Dispatch-Secret = a shared password (ALERT_DISPATCH_SECRET) so a stranger can't trigger a mass email", 0, SLATE, False),
    ],
    chose="OS cron (needs an always-on box — Render's free service sleeps), GitHub Actions schedule (best-effort timing), Celery + Redis (real infra, out of scope for v1).",
    note="The 'trigger only, no logic' separation is the point — the schedule lives outside the "
         "code so I change it without redeploying, and the secret stops user-spam.",
)

# ── SLIDE 12 — Component 7: Env vars & secrets ─────────────────────────────────
component(
    7, "Environment Variables & Secrets",
    "Config handed to the program from OUTSIDE the code at startup. A secret = a sensitive one (password, API key).",
    [
        ("Why never in git: history is forever — commit a key, delete it next commit, it's STILL in history", 0, AMBER, True),
        (".env.example is committed as a <fill_me> template; the real .env is gitignored", 0, SLATE, False),
        ("Render: value: (non-secret config) · sync:false (paste in dashboard) · generateValue (Render invents a random)", 0, SLATE, False),
        ("Vercel: dashboard vars; NEXT_PUBLIC_ prefix = a deliberate opt-in for anything the browser may see", 0, SLATE, False),
    ],
    chose="hardcoding config in source — leaks secrets into git history and ties every environment to one set of values.",
    note="The NEXT_PUBLIC_ rail is a good security detail: it forces a conscious choice to expose "
         "anything to the browser, so a backend secret can't leak into the client bundle by accident.",
)

# ── SLIDE 13 — Component 8: CI/CD ──────────────────────────────────────────────
component(
    8, "CI/CD & Git-based Deploys",
    "\"When I push to GitHub, the right things happen automatically.\" DS analogy: an automated pipeline, like scheduled retraining — but it ships the app.",
    [
        ("No .github/workflows — this project does NOT use GitHub Actions", 0, SLATE, False),
        ("It uses the platforms' built-in git auto-deploy: push to main → Render AND Vercel rebuild + redeploy", 0, SLATE, False),
        ("Vercel preview deploys: every branch/PR gets its own temporary URL to click before merging", 0, SLATE, False),
        ("Owned trade-off: no cloud test gate — the ~190 pytest tests run LOCALLY, not in CI", 0, AMBER, True),
    ],
    chose="a full GitHub Actions test-gated pipeline — simpler git auto-deploy is enough for a portfolio; the CI gate is named future work.",
    note="Be upfront: this is continuous deployment, not full CI with a test gate. Naming the "
         "gap shows I know what a mature pipeline adds — and it's a clean future-work item.",
)

# ── SLIDE 14 — Component 9: Alembic ────────────────────────────────────────────
component(
    9, "Database Migrations (Alembic)",
    "A migration = a versioned, reversible script that changes the DB schema safely. Git, but for your database's shape.",
    [
        ("Why not hand-edit the live DB: hand-edits are unrepeatable, untracked, irreversible", 0, AMBER, True),
        ("On EVERY deploy, render_build.sh runs: alembic upgrade head", 0, SLATE, False),
        ("Brings prod schema up to date BEFORE the new code starts; idempotent (no-op if already current)", 0, SLATE, False),
        ("New column? Write a migration, push — the deploy applies it. You never touch the live DB by hand.", 0, SLATE, False),
    ],
    code_lines=['alembic upgrade head   # apply every migration up to the latest'],
    chose="manual SQL on the production database — unrepeatable, unreviewable, and impossible to roll back cleanly.",
    note="Real example: an ads table and an unsubscribe token were both added AFTER launch, "
         "safely, via migrations.",
)

# ── SLIDE 15 — Component 10: CORS / Docker / HTTPS ─────────────────────────────
s = slide(); bg(s, WHITE)
title_bar(s, "CORS, Containers & HTTPS (the glue)", kicker="Component 10")
bullets(s, [
    ("CORS — why two domains need a permission slip:", 0, DEEP, True),
    ("Browsers block JS on safeearth.tech from calling a different origin (...onrender.com) by default", 1, SLATE, False),
    ("The backend posts a 'this origin is welcome' sign: CORSMiddleware reads the CORS_ORIGINS env var", 1, SLATE, False),
    ("Forget it → every browser call fails, but curl works (classic first-deploy gotcha)", 1, AMBER, False),
    ("Containers (Docker) — used locally, deliberately NOT to host the app:", 0, DEEP, True),
    ("A container = a frozen conda env that also includes the OS — runs identically anywhere", 1, SLATE, False),
    ("The only Docker here is docker-compose.yml → one local PostgreSQL for dev", 1, SLATE, False),
    ("The app uses Render's native Python buildpack, not a Dockerfile — less to maintain for one service", 1, SLATE, False),
    ("Domains & HTTPS — Vercel and Render issue + renew SSL certificates automatically", 0, DEEP, True),
    ("Add the domain, point DNS, and the padlock just works", 1, SLATE, False),
], top=Inches(1.4), size=15, gap=6)
notes(s, "Three small 'glue' topics. Honest Docker framing: a container would give bit-for-bit "
         "portability if I outgrew Render, but for one Python service a buildpack is simpler. I "
         "chose simplicity and can name what I traded.")

# ── SLIDE 16 — One Deploy, End to End (money slide) ────────────────────────────
s = slide(); bg(s, WHITE)
title_bar(s, "One Deploy, End to End", kicker="The Money Slide")
tf = box(s, Inches(0.6), Inches(1.28), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "git push origin main"
r.font.name = "Consolas"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = DEEP
r = p.add_run(); r.text = "   →   a user feels the change. Nine automated cloud steps:"
r.font.name = "Calibri"; r.font.size = Pt(16); r.font.color.rgb = SLATE
steps = [
    ("1", "You push", "one command IS the whole deploy — no SSH, no file copy"),
    ("2", "GitHub notifies both hosts", "webhooks fire Render AND Vercel"),
    ("3", "Render builds backend", "render_build.sh: deps · EM-DAT JSON · chapters.json · alembic upgrade head · seed"),
    ("4", "Render starts backend", "uvicorn ...; lifespan downloads models from HuggingFace, loads chapters.json"),
    ("5", "Health gate", "Render hits /api/v1/health — only models_loaded:true switches traffic"),
    ("6", "Vercel builds frontend", "npm install && npm run build → push .next to the global CDN"),
    ("7", "User loads safeearth.tech", "served by Vercel"),
    ("8", "Browser calls the backend", "...onrender.com/api/v1/... — allowed by CORS"),
    ("9", "Backend reads/writes Neon, replies", "all running the new code — the user never knew a deploy happened"),
]
y = 1.92
for num, head, detail in steps:
    band(s, INDIGO, Inches(0.6), Inches(y), Inches(0.5), Inches(0.42))
    tf = box(s, Inches(0.6), Inches(y - 0.02), Inches(0.5), Inches(0.46)); p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    tf = box(s, Inches(1.25), Inches(y - 0.04), Inches(11.4), Inches(0.5)); p = tf.paragraphs[0]
    r = p.add_run(); r.text = head + "  —  "
    r.font.size = Pt(13.5); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
    r = p.add_run(); r.text = detail
    r.font.size = Pt(12.5); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    y += 0.5
tf = box(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "I changed one line and typed git push — nine steps across five services happened on their own. THAT automation is what \"deployed\" means."
r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
notes(s, "Where it clicks — walk it slowly, name the real file at each hop. If short on time, "
         "cut a component slide, never this one. The health gate (step 5) is worth a beat: a "
         "failed build never switches traffic, so the old version stays live.")

# ── SLIDE 17 — Trade-offs ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Decisions & Trade-offs")
table(s, [
    ["Decision", "Chosen", "Over", "Why"],
    ["Frontend host", "Vercel", "Netlify / CF / VPS", "Makes Next.js; free CDN + HTTPS"],
    ["Backend host", "Render (free)", "Railway / Fly / AWS", "Long-lived Python + IaC, free"],
    ["Database", "Neon Postgres", "Supabase / RDS", "Serverless, scales to zero, free"],
    ["Models", "HuggingFace Hub", "In-git / S3", "Built for models; pull at startup"],
    ["RAG retrieval", "Chapter JSON + Groq", "ChromaDB vectors", "PyTorch (~2 GB) OOM'd 512 MB"],
    ["Scheduling", "n8n (trigger-only)", "cron / GH Actions", "Free, survives the sleeping app"],
    ["CI/CD", "Git auto-deploy", "GH Actions gate", "Simplest path; gate is future work"],
    ["Containers", "Buildpack (app)", "Dockerfile", "Less to maintain for one service"],
], Inches(0.5), Inches(1.4), Inches(12.3),
   col_widths=[Inches(2.3), Inches(3.0), Inches(2.7), Inches(4.3)], fsize=13)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Maturity isn't 'I used the trendiest tool' — it's 'I matched tools to a single-dev, zero-cost, 512 MB constraint, and I know what I traded.'"
r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
notes(s, "This slide shows engineering maturity for the room. Every row pairs a choice with "
         "what it cost.")

# ── SLIDE 18 — Challenges I Solved ─────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Challenges I Solved (Deployment Edition)")
bullets(s, [
    ("512 MB RAM ceiling → dropped PyTorch/ChromaDB RAG for build-time chapters.json; rebuilt SHAP from the in-memory model (~300 MB saved)", 0, SLATE, False),
    ("Big models vs git → .pkl's live on HuggingFace, downloaded at startup; the repo stays light", 0, SLATE, False),
    ("Free-tier cold starts → UptimeRobot pings /api/v1/health every 14 minutes to keep Render awake", 0, SLATE, False),
    ("Secrets across 3 dashboards → one .env.example template; sync:false / generateValue; NEXT_PUBLIC_ opt-in", 0, SLATE, False),
    ("Schema changes on a live DB → alembic upgrade head baked into the build script, never a manual edit", 0, SLATE, False),
    ("Cross-origin calls → an explicit CORS_ORIGINS allowlist read by CORSMiddleware", 0, SLATE, False),
], top=Inches(1.7), size=18, gap=15)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "The through-line: almost every challenge traces back to the free-tier RAM/sleep limits. Constraints drove the architecture."
r.font.size = Pt(15); r.font.italic = True; r.font.color.rgb = GREY; r.font.name = "Calibri"

# ── SLIDE 19 — Honest Engineering ──────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Honest Engineering: Portfolio vs Production")
tf = box(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Production-grade here:"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
bullets(s, [
    ("Infrastructure-as-code (render.yaml, vercel.json) · migrations on every deploy · health-checked rollouts", 0, SLATE, False),
    ("Graceful LLM degradation · secrets never in git · automatic HTTPS · same Postgres locally and in prod", 0, SLATE, False),
], top=Inches(1.95), size=16, gap=6)
tf = box(s, Inches(0.6), Inches(3.15), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Conscious free-tier shortcuts (each with an upgrade path):"
r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = AMBER; r.font.name = "Calibri"
bullets(s, [
    ("Backend sleeps & cold-starts → a ~$7/mo tier removes it", 0, SLATE, False),
    ("512 MB RAM forced the simplified RAG → a larger tier brings vector search back", 0, SLATE, False),
    ("No cloud test gate before deploy → add a GitHub Actions CI job", 0, SLATE, False),
    (".env.production points at the raw Render URL → swap to api.safeearth.tech", 0, SLATE, False),
], top=Inches(3.7), size=16, gap=9)
band(s, LIGHT, Inches(0.6), Inches(6.35), Inches(12.1), Inches(0.75))
tf = box(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Knowing what you traded away is the engineering — not a weakness."
r.font.size = Pt(17); r.font.bold = True; r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
notes(s, "Deliberate slide: it pre-empts every 'but isn't X a problem?' by naming X first, with "
         "the fix. Professors reward the engineer who critiques their own system.")

# ── SLIDE 20 — Results ──────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Results")
bullets(s, [
    ("Live in production: safeearth.tech (Vercel) ↔ FastAPI backend (Render) ↔ Neon Postgres", 0, GREEN, True),
    ("Seven coordinated cloud services — each free-tier, each declared as config in the repo", 0, SLATE, False),
    ("One-command deploys: git push → automatic build + migrate + ship on both hosts", 0, SLATE, False),
    ("Self-healing rollouts: health-checked startup, models auto-pulled from HuggingFace", 0, SLATE, False),
    ("Resilient by design: an LLM failure falls back to the DB — a hiccup never takes down predictions", 0, SLATE, False),
    ("Kept awake & observable: /api/v1/health reports models_loaded + rag_loaded; UptimeRobot watches it", 0, SLATE, False),
], top=Inches(1.7), size=18, gap=15)

# ── SLIDE 21 — Future Work ──────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Future Work")
bullets(s, [
    ("Pay off the free tier: upgrade Render to kill cold starts + raise RAM; re-introduce vector RAG for finer retrieval", 0, SLATE, False),
    ("Add a real CI gate: GitHub Actions runs ~190 pytest tests + the frontend build, and BLOCKS deploys on failure", 0, SLATE, False),
    ("Attach custom domains end-to-end: api.safeearth.tech for the backend; update .env.production", 0, SLATE, False),
    ("Harden observability: structured logs + error tracking (e.g. Sentry) instead of dashboard log-reading", 0, SLATE, False),
    ("Containerize the backend (Dockerfile) for bit-for-bit portability if it outgrows Render", 0, SLATE, False),
], top=Inches(1.8), size=19, gap=17)
tf = box(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Every item maps directly to a trade-off named earlier — the shortcuts were conscious, with a known path forward."
r.font.size = Pt(15); r.font.italic = True; r.font.color.rgb = GREY; r.font.name = "Calibri"

# ── SLIDE 22 — Key Takeaways ────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("Deployment = model.predict() made always-on, public, and self-shipping for strangers.", 0, DEEP, True),
    ("SafeEarth runs as seven free-tier cloud services, wired together and declared as code.", 0, SLATE, False),
    ("The 512 MB / sleeping free tier was the dominant constraint — and it DROVE the architecture.", 0, SLATE, False),
    ("Resilience + migrations are non-negotiable: graceful fallback and schema-as-code keep the live system safe.", 0, SLATE, False),
    ("Every trade-off was deliberate, documented, and reversible.", 0, SLATE, False),
], top=Inches(2.0), size=20, gap=20)

# ── SLIDE 23 — DS Translation Table ────────────────────────────────────────────
s = slide(); bg(s, WHITE)
title_bar(s, "The Data-Scientist Translation Table")
table(s, [
    ["Deployment concept", "Closest data-science equivalent"],
    ["Deploying an app", "Turning a notebook model into a callable endpoint anyone can hit"],
    ["Host / server", "A remote machine running your code 24/7 vs. your watched laptop"],
    ["localhost vs. production", "A scratch experiment vs. the published result others depend on"],
    ["Environment variables / secrets", "Config + API keys kept OUT of a shared notebook"],
    [".env.example", "A requirements.txt, but for which secrets to provide"],
    ["CI/CD", "An automated pipeline that rebuilds + ships on every change"],
    ["Container (Docker)", "A frozen conda env that runs identically anywhere"],
    ["Managed database (Neon)", "A hosted, concurrent data store vs. a local CSV"],
    ["Migration (Alembic)", "A tracked, re-runnable transform for your schema"],
    ["Model hosting (HuggingFace)", "A model registry instead of a .pkl in your repo"],
    ["Cold start", "A kernel that shut down on idle and must re-import everything"],
    ["Health check", "assert models_loaded — proof the pipeline initialized"],
], Inches(0.9), Inches(1.3), Inches(11.5),
   col_widths=[Inches(4.3), Inches(7.2)], fsize=13)
notes(s, "For this audience, the slide that makes everything land. Pause on it. Invite them to "
         "map any deployment term onto something they already do every day.")

# ── SLIDE 24 — Thank You ────────────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, INDIGO, 0, Inches(3.0), W, Inches(0.08))
tf = box(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = INDIGO; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(4.6), Inches(11.3), Inches(1.4)); p = tf.paragraphs[0]
for i, txt in enumerate([
    "SafeEarth Intelligence — Vercel · Render · Neon · HuggingFace · Groq · n8n",
    "Full deep-dive: docs/DEPLOYMENT_EXPLAINED.md  ·  Script: docs/DEPLOYMENT_PITCH_DECK.md",
    "safeearth.tech",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Likely Qs: Why not AWS/Kubernetes? (one small app; PaaS free tiers do the job with a "
         "fraction of the ops). Isn't the free tier a problem? (cold start + 512 MB are real; "
         "UptimeRobot mitigates one, the RAM limit DROVE the RAG design; both have a paid-tier "
         "fix). Keeping secrets safe? (never in git; gitignored .env; sync:false/generateValue; "
         "NEXT_PUBLIC_ gates the browser). Groq/HF down? (Groq → DB fallback; HF only matters at "
         "cold start). Why two URLs / CORS? (different origins; backend allowlists the frontend). "
         "Migration fails? (runs in the build; a failed build doesn't switch traffic — old "
         "version stays live).")

out = Path(__file__).resolve().parent / "SafeEarth_Deployment_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides)")
