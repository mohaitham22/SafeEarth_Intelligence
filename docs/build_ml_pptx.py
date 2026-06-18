"""
build_ml_pptx.py — generate an editable PowerPoint (with charts) from the SafeEarth ML deep dive.

Run:  py -3.12 docs/build_ml_pptx.py
Out:  docs/SafeEarth_ML_Presentation.pptx

Native python-pptx (editable text boxes + tables) + embedded PNG charts rendered by
make_ml_charts.py. Content mirrors docs/ML_PRESENTATION.md, which is derived from
docs/ML_DEEP_DIVE.md. Every number is from this project's own metric files, saved bundles,
and reproducible runs (re-run 2026-06-15).

Caveat carried on every metric slide: the "holdout" test set is a 100% subset of train, so
holdout numbers are resubstitution [resub]; the honest generalization estimate is the
5-fold CV [CV]. See docs/ML_DEEP_DIVE.md.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
import make_ml_charts  # noqa: E402

ASSETS = Path(__file__).resolve().parent / "ml_assets"

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xDC, 0x26, 0x26)
SLATE  = RGBColor(0x33, 0x41, 0x55)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)
PINK   = RGBColor(0xFE, 0xE2, 0xE2)

W, H = Inches(13.333), Inches(7.5)

# Ensure chart PNGs exist before embedding.
if not (ASSETS / "05_confusion.png").exists():
    print("Charts not found — rendering them first...")
    make_ml_charts.generate_all()

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def img(s, name, left, top, width):
    return s.shapes.add_picture(str(ASSETS / name), left, top, width=width)


def title_bar(s, text, kicker=None, color=DEEP):
    band(s, color, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = (RGBColor(0xFE, 0xCA, 0xCA) if color == RED else GREEN)
        r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    first = True
    for it in items:
        text, level, color, bold = (it + (0, SLATE, False))[:4] if isinstance(it, tuple) else (it, 0, SLATE, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0)); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None, hcolor=DEEP):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val); r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hcolor
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def callout(s, label, text, top=Inches(6.5), color=LIGHT, label_color=GREEN, h=Inches(0.7),
            left=Inches(0.6), width=Inches(12.1)):
    band(s, color, left, top, width, h)
    tf = box(s, left + Inches(0.2), Inches(top.inches + 0.10), Inches(width.inches - 0.4), Inches(h.inches - 0.16))
    p = tf.paragraphs[0]
    if label:
        r = p.add_run(); r.text = label + "  "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = label_color; r.font.name = "Calibri"
    r = p.add_run(); r.text = text
    r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
for i, (txt, sz, col, bold) in enumerate([
    ("Machine-Learning Deep Dive", 28, GREEN, True),
    ("Every model, every metric — and a rigorous floor/ceiling analysis", 19, RGBColor(0xCB, 0xD5, 0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "EM-DAT (1900–2021)  ·  XGBoost + CatBoost ensemble  ·  per-type impact regressors"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Thesis of this talk: I'll give the real per-model numbers, then argue precisely WHY "
         "they sit where they do — and show that the headline macro-F1 is inflated by a "
         "train/test leak I found and quantified. Every number is from this project's own runs.")

# ── SLIDE 2 — Roadmap ─────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Roadmap")
bullets(s, [
    ("Two ML problems — 8-class disaster type, and 4 impact regressors", 0, SLATE, False),
    ("Every model trained, with baselines that anchor the floor", 0, SLATE, False),
    ("Per-model metrics — classification (macro-F1, per-class) and regression (log-MAE, R²)", 0, SLATE, False),
    ("What actually moved the numbers (and what didn't)", 0, GREEN, True),
    ("The contamination I found: the headline 0.705 is in-sample; honest ≈ 0.586", 0, RED, True),
    ("Floor & ceiling — the bias/variance verdict per problem", 0, DEEP, True),
    ("Recommendations — highest-ceiling-first", 0, SLATE, False),
], top=Inches(1.6), size=19, gap=15)
notes(s, "Tell them the spine is honesty: real numbers, negative results, a self-found leak, "
         "and a concrete bias-vs-variance diagnosis.")

# ── SLIDE 3 — Two problems ─────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Two ML Problems, One Feature Pipeline")
table(s, [
    ["", "Classification", "Regression"],
    ["Target", "Disaster type (8 classes)", "Deaths · Injuries · Affected · Damage"],
    ["Model", "XGBoost + CatBoost soft-vote", "Per-type XGB / RandomForest"],
    ["Loss", "softmax / multiclass logloss", "squared error on log1p(target)"],
    ["Metric that counts", "macro-F1 (21:1 imbalance)", "log1p-MAE on observed rows"],
    ["Why that metric", "rare-disaster miss = costly", "raw R² ≈ 0 — tail is unpredictable"],
], Inches(0.55), Inches(1.5), Inches(12.2),
   col_widths=[Inches(2.7), Inches(4.4), Inches(5.1)], fsize=14)
callout(s, "Design rule:", "impact is right-skewed 100×–1,400× → median / log1p everywhere, never mean (CLAUDE.md).",
        top=Inches(5.5), h=Inches(0.7))
callout(s, "8 classes:", "Drought · Earthquake · Extreme temperature · Flood · Landslide · Storm · Volcanic activity · Wildfire.",
        top=Inches(6.4), h=Inches(0.7), color=LIGHT, label_color=DEEP)

# ── SLIDE 4 — Data & splits ────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Data, Splits & the 21:1 Imbalance")
table(s, [
    ["Class", "Train n", "Train %", "Holdout n"],
    ["Flood", "5,551", "38.3%", "5,272"],
    ["Storm", "4,496", "31.1%", "4,005"],
    ["Earthquake", "1,544", "10.7%", "1,137"],
    ["Landslide", "776", "5.4%", "713"],
    ["Drought", "770", "5.3%", "685"],
    ["Extreme temperature", "603", "4.2%", "584"],
    ["Wildfire", "471", "3.3%", "452"],
    ["Volcanic activity", "265", "1.8%", "222"],
], Inches(0.55), Inches(1.4), Inches(6.6),
   col_widths=[Inches(3.0), Inches(1.3), Inches(1.2), Inches(1.1)], fsize=12.5)
bullets(s, [
    ("Train 14,476 · holdout 13,070 (8 modelled types)", 0, SLATE, False),
    ("Imbalance ≈ 21:1 (Flood : Volcanic)", 0, DEEP, True),
    ("CV = 5-fold StratifiedKFold, random_state=42", 0, SLATE, False),
    ("Encoders + freq-maps fit on the train fold only", 0, SLATE, False),
    ("Target coverage: deaths 70%, injuries 26%,", 0, AMBER, True),
    ("affected 56%, damage 36% — nulls ≠ zeros", 1, SLATE, False),
], left=Inches(7.4), top=Inches(1.5), width=Inches(5.5), size=15, gap=10)
callout(s, "Foreshadow:", "the test file is a 100% subset of train — quantified on Slide 14.",
        top=Inches(6.5), h=Inches(0.65), color=PINK, label_color=RED)

# ── SLIDE 5 — Models inventory ─────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Models Inventory")
tf = box(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Classification"; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
table(s, [
    ["Model", "Role", "Outcome"],
    ["Logistic Regression", "candidate", "rejected — macro-F1 0.137 (≈ dummy)"],
    ["Random Forest", "candidate", "competitive 0.700, no edge"],
    ["XGBoost", "final (primary)", "best single learner — 0.60 weight"],
    ["CatBoost", "final (member)", "diversity — 0.40 weight"],
    ["LightGBM", "dropped", "optimal weight tuned to 0.00"],
    ["Soft-vote ensemble (v4.2)", "deployed", "0.60·XGB + 0.40·CAT"],
], Inches(0.6), Inches(1.65), Inches(6.1),
   col_widths=[Inches(2.5), Inches(1.7), Inches(1.9)], fsize=11.5)
tf = box(s, Inches(7.0), Inches(1.25), Inches(6), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Regression (per-type + drop-null)"; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
table(s, [
    ["Target", "Model"],
    ["deaths, damage", "XGBRegressor"],
    ["injuries, affected", "RandomForest"],
    ["(light variant)", "XGB for all 4 (23 MB)"],
    ["Ridge", "rejected (worst / exploded)"],
    ["SVR / KNN", "never built → §recs"],
], Inches(7.0), Inches(1.65), Inches(5.9),
   col_widths=[Inches(2.6), Inches(3.3)], fsize=11.5)
callout(s, "Why LightGBM went:", "weakest member (0.654) AND the macro-F1-optimal ensemble weight for it was exactly 0.00 → v4.2 dropped it (bundle 26 MB → 10 MB).",
        top=Inches(5.7), h=Inches(0.85))
notes(s, "SVM/SVR/KNN were named in the brief but never built — given LogReg's collapse, a "
         "kernel SVM is unlikely to beat trees and scales badly at n≈14k. Moved to recommendations.")

# ── SLIDE 6 — Baselines ────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Baselines — the Floor Reference")
tf = box(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Classification dummies (holdout)"; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
table(s, [
    ["Strategy", "Accuracy", "Macro-F1"],
    ["most_frequent (Flood)", "0.403", "0.072"],
    ["stratified", "0.271", "0.129"],
], Inches(0.6), Inches(1.65), Inches(5.9),
   col_widths=[Inches(3.2), Inches(1.4), Inches(1.3)], fsize=13)
tf = box(s, Inches(0.6), Inches(3.35), Inches(12), Inches(0.4)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Regression — train-median predictor (log-MAE)"; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
table(s, [
    ["Target", "median", "log-MAE", "factor"],
    ["deaths", "17", "1.22", "3.4×"],
    ["injuries", "50", "1.54", "4.7×"],
    ["affected", "15,000", "2.36", "10.6×"],
    ["damage ('000$)", "60,000", "2.12", "8.3×"],
], Inches(0.6), Inches(3.75), Inches(5.9),
   col_widths=[Inches(2.4), Inches(1.4), Inches(1.2), Inches(0.9)], fsize=12.5)
bullets(s, [
    ("Macro-F1 floor ≈ 0.07–0.13", 0, DEEP, True),
    ("Honest [CV] macro-F1 0.586 sits +0.46 above", 0, SLATE, False),
    ("the stratified dummy — 4.5× its value.", 1, SLATE, False),
    ("", -1, SLATE, False),
    ("Every deployed regressor beats its", 0, GREEN, True),
    ("median baseline (e.g. damage 2.66× vs 8.3×).", 1, SLATE, False),
    ("Judge regression lift vs THIS, not raw R².", 0, AMBER, True),
], left=Inches(7.0), top=Inches(1.6), width=Inches(5.9), size=15, gap=9)

# ── SLIDE 7 — Candidate sweep (chart) ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Model Family Is the Lever", kicker="Classification")
img(s, "01_candidates.png", Inches(2.65), Inches(1.35), Inches(8.0))
callout(s, "The point:", "swapping linear → trees buys ~+0.57 macro-F1. All three tree models cluster at 0.700–0.704; the choice among them is noise (≈ CV std).",
        top=Inches(6.45), h=Inches(0.7), label_color=GREEN)
notes(s, "Logistic Regression collapses to ~the stratified dummy — the type boundaries are "
         "massively non-linear (collinear, mixed-unit magnitude, 225-level categoricals). "
         "This single decision is essentially the entire score.")

# ── SLIDE 8 — Ensemble (chart) ─────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Ensemble — and Why LightGBM Was Dropped", kicker="Classification")
img(s, "02_ensemble.png", Inches(2.5), Inches(1.35), Inches(8.3))
callout(s, "Honest read:", "v4.2 (XGB+CAT) beats its best single member by +0.0113 — real, but ~2× the CV std. The ensemble is a modest lever, not a transformation.",
        top=Inches(6.45), h=Inches(0.7), label_color=GREEN)
notes(s, "A step-0.05 simplex grid over (XGB,LGB,CAT) put LightGBM's optimal weight at 0.00. "
         "The v4.1 3-way ensemble actually scored BELOW its own best member; removing LGB and "
         "reweighting to 0.6/0.4 is what made the ensemble net-positive.")

# ── SLIDE 9 — Tuning ───────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Hyperparameter Tuning — and an Honest Caveat")
bullets(s, [
    ("Optuna TPE (seed 42), maximizing 3-fold CV macro-F1: XGB 40 trials · LGB 30 · CAT 20", 0, SLATE, False),
    ("Search space e.g. XGB: depth 4–10, lr 0.01–0.2 (log), reg_alpha 0–2, reg_lambda 0.5–5, subsample/colsample", 1, GREY, False),
    ("Selected XGB: n=481, depth=8, lr=0.052, reg_alpha=1.85, reg_lambda=2.40, gamma=0.50", 0, SLATE, False),
    ("Selected CatBoost: 426 iters, depth 6, lr 0.169, l2_leaf_reg 9.71", 0, SLATE, False),
    ("Measured gain from tuning over a sensible default: +0.0018 macro-F1 (0.7020 → 0.7038)", 0, AMBER, True),
    ("→ within fold noise (CV std ≈ 0.006–0.008). Tuning is NOT where the points came from.", 1, RED, True),
], top=Inches(1.55), size=17, gap=13)
callout(s, "Takeaway:", "model family bought ~0.57; Optuna bought ~0.002. Spend effort on features & evaluation, not the 40th trial.",
        top=Inches(6.45), h=Inches(0.7), label_color=GREEN)

# ── SLIDE 10 — Classification scoreboard ───────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Per-Model Scoreboard — Classification")
table(s, [
    ["Model / config", "Macro-F1 [resub]", "Wtd-F1", "Accuracy", "Macro-F1 [CV]"],
    ["XGBoost only (v4.1)", "0.6939", "—", "—", "—"],
    ["LightGBM only (v4.1)", "0.6540", "—", "—", "—"],
    ["CatBoost only (v4.1)", "0.6817", "—", "—", "—"],
    ["Ensemble v4.1 (X·L·C)", "0.6929", "0.7519", "0.7467", "—"],
    ["Ensemble v4.2 (X·C) — deployed", "0.7052", "0.7587", "0.7467", "0.5857 ±0.006"],
], Inches(0.5), Inches(1.45), Inches(12.3),
   col_widths=[Inches(4.3), Inches(2.4), Inches(1.8), Inches(1.9), Inches(1.9)], fsize=13)
bullets(s, [
    ("[resub] = full test set, which is a 100% subset of train → a training-set fit.", 0, RED, True),
    ("[CV] = 5-fold StratifiedKFold on train → the honest generalization estimate.", 0, GREEN, True),
    ("The deployed model is 0.7052 in-sample but ~0.586 out-of-fold — a 0.12 gap (Slides 14–15).", 0, DEEP, True),
    ("ROC-AUC / PR-AUC / log-loss were never computed — flagged in recommendations.", 0, GREY, False),
], top=Inches(3.9), size=16, gap=13)

# ── SLIDE 11 — Confusion matrix (chart) ────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Where the Classifier Fails", kicker="Error Analysis")
img(s, "05_confusion.png", Inches(0.5), Inches(1.35), Inches(7.5))
bullets(s, [
    ("Earthquake near-perfect:", 0, DEEP, True),
    ("1,110 / 1,137 — the magnitude signal.", 1, SLATE, False),
    ("Hydro-met cluster bleeds:", 0, AMBER, True),
    ("565 Flood→Storm, 518 Storm→Flood,", 1, SLATE, False),
    ("168 Landslide→Flood, 103 Landslide→Storm.", 1, SLATE, False),
    ("Minority precision collapses:", 0, RED, True),
    ("Landslide P=0.43, Drought P=0.48 —", 1, SLATE, False),
    ("class weights trade precision for recall.", 1, SLATE, False),
], left=Inches(8.2), top=Inches(1.5), width=Inches(4.9), size=14.5, gap=8)
notes(s, "The four hydro-meteorological classes overlap in (continent, region, season, "
         "coordinate) space — ~1,800 mutual misclassifications. The features have no "
         "hydrology/meteorology to separate them. Earthquake is the one class with a clean feature.")

# ── SLIDE 12 — Regression metrics ──────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Per-Model Metrics — Regression", kicker="Regression")
table(s, [
    ["Target", "best model", "MAE (non-zero)", "raw R²", "log-MAE (factor)"],
    ["deaths", "XGB (300,d6)", "368", "0.0202", "0.63 (1.9×)"],
    ["injuries", "XGB (300,d6)", "1,830", "0.0151", "1.02 (2.8×)"],
    ["affected", "XGB (300,d6)", "990,998", "0.0158", "1.46 (4.3×)"],
    ["damage ('000$)", "XGB (300,d6)", "747,942", "0.0142", "0.98 (2.7×)"],
], Inches(0.5), Inches(1.45), Inches(7.2),
   col_widths=[Inches(1.9), Inches(1.7), Inches(1.7), Inches(1.0), Inches(0.9)], fsize=11.5)
img(s, "07_reg_r2.png", Inches(7.95), Inches(1.5), Inches(5.0))
callout(s, "The honest read:", "raw R² ≈ 0.01–0.02 even IN-SAMPLE → the features cannot fit the heavy tail. log-MAE captures type-conditioned central tendency; magnitude itself is unpredictable here.",
        top=Inches(5.7), h=Inches(0.85), color=PINK, label_color=RED)

# ── SLIDE 13 — Regressor fix (chart) ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Regressor Fix — Per-Type + Drop-Null", kicker="Regression")
img(s, "06_reg_factors.png", Inches(0.5), Inches(1.4), Inches(7.2))
bullets(s, [
    ("Old regressors were type-blind AND", 0, SLATE, False),
    ("trained on fake zeros (64% of damage", 1, SLATE, False),
    ("rows are NULL → filled 0 → model", 1, SLATE, False),
    ("learns ≈0).", 1, SLATE, False),
    ("Two independent wins:", 0, GREEN, True),
    ("add the disaster-TYPE axis +", 1, SLATE, False),
    ("DROP the fake zeros.", 1, SLATE, False),
    ("Caveat: this un-poisons the target;", 0, AMBER, True),
    ("it does NOT crack magnitude (R²≈0).", 1, SLATE, False),
], left=Inches(7.9), top=Inches(1.5), width=Inches(5.1), size=14, gap=7)
notes(s, "Validated on an honest holdout (observed-target rows only). The headline 111×→2.7× "
         "for damage is mostly about not poisoning the model with fake zeros, plus conditioning "
         "on type — not about new feature signal.")

# ── SLIDE 14 — What improved (levers) ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What Actually Moved the Numbers")
table(s, [
    ["Lever", "Δ", "Evidence"],
    ["Model family: linear → trees", "+0.565 macro-F1", "benchmark sweep"],
    ["Hand weights vs balanced inv-freq", "+0.090", "minority strategies"],
    ["Centroid imputation (validated, NOT shipped)", "+0.017", "preprocessing exp."],
    ["Drop LightGBM + reweight (v4.2)", "+0.011", "weight tuning"],
    ["Ensemble weight tuning (LGB→0)", "+0.0014", "weight tuning"],
    ["Optuna over a sensible default", "+0.0018", "candidate sweep"],
    ["Per-class threshold tuning", "−0.0027 → DISCARDED", "overfit val"],
    ["Regression: per-type + drop-null", "111× → 2.7× (damage)", "combined exp."],
], Inches(0.5), Inches(1.4), Inches(12.3),
   col_widths=[Inches(5.6), Inches(3.0), Inches(3.7)], fsize=12.5)
callout(s, "The verdict:", "one lever (trees) bought ~0.57; everything after lives in the ±0.01–0.09 band, and several 'improvements' were noise or negative. Hand weights mostly rebalance precision/recall.",
        top=Inches(6.45), h=Inches(0.75), label_color=GREEN)

# ── SLIDE 15 — Contamination (chart, red theme) ────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Headline Metric Is Inflated", kicker="Problem I Found", color=RED)
img(s, "04_floor_ceiling.png", Inches(0.5), Inches(1.4), Inches(7.0))
bullets(s, [
    ("The 'holdout' 1970–2021 file is a", 0, SLATE, False),
    ("STRICT SUBSET of the 1900–2021 train file.", 1, SLATE, False),
    ("Verified: 100.0% of the 13,070 test", 0, RED, True),
    ("rows share an EM-DAT Year-Seq-Country", 1, SLATE, False),
    ("key with a train row.", 1, SLATE, False),
    ("→ macro-F1 0.7052 is IN-SAMPLE.", 0, AMBER, True),
    ("Honest 5-fold CV ≈ 0.586.", 0, GREEN, True),
    ("It voids absolute numbers, not the", 0, SLATE, False),
    ("relative decisions (those used CV).", 1, SLATE, False),
], left=Inches(7.8), top=Inches(1.5), width=Inches(5.2), size=14, gap=7)
notes(s, "Deliver this calmly and confidently — finding and quantifying it is the strength. "
         "Most likely faculty probe. The relative experiment decisions used CV / honest "
         "holdout, so the engineering choices still stand; only the headline number is inflated.")

# ── SLIDE 16 — Floor & ceiling (chart) ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Floor & Ceiling — the Bias/Variance Verdict")
img(s, "03_resub_vs_cv.png", Inches(0.5), Inches(1.35), Inches(7.4))
bullets(s, [
    ("Earthquake gap +0.012 → bias-limited,", 0, DEEP, True),
    ("already at its signal ceiling (0.96).", 1, SLATE, False),
    ("Minorities gap +0.16–0.26 → variance-", 0, RED, True),
    ("limited: tiny support (Volcanic n=265),", 1, SLATE, False),
    ("F1 variance-dominated, memorized.", 1, SLATE, False),
    ("Regression: raw R²≈0 in-sample →", 0, AMBER, True),
    ("bias/signal-limited (features don't", 1, SLATE, False),
    ("determine magnitude).", 1, SLATE, False),
], left=Inches(8.1), top=Inches(1.5), width=Inches(5.0), size=14, gap=7)
notes(s, "This slide IS the answer to 'why not better.' Majority classes: bias-limited → need "
         "better features, not more data. Minorities: variance-limited → more data / merging / "
         "two-stage. Regression: signal-limited → exogenous features (exposure, GDP, population).")

# ── SLIDE 17 — Explainability (chart) ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Explainability — Feature Importance & a Red Flag")
img(s, "08_importance.png", Inches(0.5), Inches(1.4), Inches(7.3))
bullets(s, [
    ("Two magnitude features = 45.7% of gain.", 0, DEEP, True),
    ("Red flag: has_magnitude is essentially", 0, AMBER, True),
    ("'is this an Earthquake/Storm' — one", 1, SLATE, False),
    ("binary carries 20% of decisions and", 1, SLATE, False),
    ("is why Earthquake hits F1 0.98.", 1, SLATE, False),
    ("(Not leakage — magnitude is known", 1, GREY, False),
    ("pre-event.)", 1, GREY, False),
    ("day_offset = 0.0%: constant in training,", 0, RED, True),
    ("so the 30-day forecast varies only by season.", 1, SLATE, False),
], left=Inches(8.0), top=Inches(1.5), width=Inches(5.1), size=13.5, gap=6)
notes(s, "SHAP is computed per-prediction (top-3) at inference; a formal per-class SHAP "
         "table isn't persisted. The gain importances shown are from the deployed v4.2 XGBoost.")

# ── SLIDE 18 — Recommendations ─────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Recommendations — Highest Ceiling First")
bullets(s, [
    ("P0 · Fix the evaluation: de-duplicate by Year-Seq-Country, adopt a TEMPORAL split, re-measure. The prerequisite for trusting any gain.", 0, RED, True),
    ("P0 · Report ROC-AUC / PR-AUC / log-loss + learning curves to confirm the bias/variance verdicts empirically.", 0, SLATE, False),
    ("P1 · Variance-limited minorities: ship centroid imputation (+0.017, already validated, unshipped); two-stage common-vs-rare classifier; consider class merging.", 0, GREEN, True),
    ("P2 · Bias-limited hydro-met cluster: add external features (ERA5/CHIRPS precip, elevation, Köppen) — the only real lever (exp. +0.03–0.06).", 0, SLATE, False),
    ("P3 · Probability calibration (isotonic/temperature) — severity thresholds drive alerts but run on uncalibrated probabilities.", 0, SLATE, False),
    ("P4 · Regression: add exposure features (population, GDP, building codes) or reframe to ordinal risk bands; rebalance the EM-DAT blend now that regressors are stronger.", 0, SLATE, False),
    ("Skip SVM/SVR/KNN — given LogReg's collapse, they won't beat trees and scale badly.", 0, GREY, False),
], top=Inches(1.45), size=15, gap=12)
notes(s, "Ordered by expected impact. P0 (fix the split) and P2/P1 (features for the signal-poor "
         "classes) are the only levers that move the real ceiling; the rest is rigor.")

# ── SLIDE 19 — Takeaways ───────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("Model family is the lever: linear → trees bought ~0.57 macro-F1; tuning/ensembling bought ~0.01.", 0, DEEP, True),
    ("The honest number is 0.586, not 0.705 — I found and quantified a 100% train/test overlap.", 0, RED, True),
    ("Classification: majority classes bias-limited, minorities variance-limited. Regression: signal-limited (raw R²≈0).", 0, SLATE, False),
    ("The biggest regressor win was un-poisoning the target (drop fake zeros + per-type), not a better algorithm.", 0, SLATE, False),
    ("Negative results are results: SMOTE, balanced weights, threshold tuning, extra geo features — all tested and rejected with reasons.", 0, SLATE, False),
], top=Inches(1.7), size=18, gap=20)
notes(s, "End on rigor and honesty — the differentiator for an academic audience.")

# ── SLIDE 20 — Thank you / Q&A ─────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(2.55), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(2.7), Inches(11.3), Inches(0.6)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
tf = box(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(3.2)); p = tf.paragraphs[0]
qa = [
    ("Anticipated questions:", DEEP, True, 16),
    ("Is 0.705 real?  →  No — in-sample upper bound; test ⊂ train. Honest CV ≈ 0.586.", SLATE, False, 14),
    ("Why not oversample?  →  Tested; SMOTE −0.04 macro & invalid on encoded categoricals.", SLATE, False, 14),
    ("Why trees over linear?  →  LogReg collapsed to 0.137; mixed-unit magnitude + 225-level categoricals + outliers break LR/SVM.", SLATE, False, 14),
    ("Biggest single win available?  →  Fix the split, then add magnitude/exposure features for the signal-poor classes.", SLATE, False, 14),
    ("Why is regression R² ≈ 0?  →  Impact needs exposure/GDP/population — absent from the feature set.", SLATE, False, 14),
]
for i, (txt, col, bold, sz) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
    p.space_after = Pt(7)
    r = p.add_run(); r.text = ("•  " if i else "") + txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0) if i else GREEN
    r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(6.95), Inches(11.3), Inches(0.4)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Full write-up: docs/ML_DEEP_DIVE.md   ·   Data: EM-DAT (1900–2021)   ·   safeearth.tech"
r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Keep the contamination answer ready — most likely probe; owning it is the win.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "SafeEarth_ML_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides, charts embedded)")
