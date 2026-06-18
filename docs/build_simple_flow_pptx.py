"""
build_simple_flow_pptx.py — generate an editable, plain-language PowerPoint that
explains the SafeEarth RAG pipeline AND recommendation system as a simple,
step-by-step sequence anyone can follow.

Run:  py -3.12 docs/build_simple_flow_pptx.py
Out:  docs/SafeEarth_RAG_RecSys_Simple_Flow.pptx

Content mirrors docs/RAG_AND_RECSYS_SIMPLE_FLOW_PITCH.md, which is the simple
companion to docs/RAG_AND_RECSYS_SIMPLE_FLOW.md. Numbers are pulled from this
project's own source (backend/rag/recommender.py, extract_chapters.py,
services/recommendation_service.py, backend/rag/chapters.json).

This is the SIMPLE deck (plain words, native flow diagrams, 3 friendly charts).
It is deliberately separate from the dense SafeEarth_RAG_Presentation.pptx
deep-dive deck.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
import make_simple_flow_charts  # noqa: E402

ASSETS = Path(__file__).resolve().parent / "simple_flow_assets"

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xDC, 0x26, 0x26)
SLATE  = RGBColor(0x33, 0x41, 0x55)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x93, 0x33, 0xEA)
PINK   = RGBColor(0xFE, 0xE2, 0xE2)

W, H = Inches(13.333), Inches(7.5)

# Ensure chart PNGs exist before embedding.
if not (ASSETS / "01_why_chapters.png").exists():
    print("Charts not found — rendering them first...")
    make_simple_flow_charts.generate_all()

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def img(s, name, left, top, width):
    return s.shapes.add_picture(str(ASSETS / name), left, top, width=width)


def title_bar(s, text, kicker=None, color=DEEP):
    band(s, color, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = (RGBColor(0xFE, 0xCA, 0xCA) if color == RED else GREEN)
        r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    first = True
    for it in items:
        text, level, color, bold = (it + (0, SLATE, False))[:4] if isinstance(it, tuple) else (it, 0, SLATE, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0)); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None, hcolor=DEEP):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val); r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hcolor
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def callout(s, label, text, top=Inches(6.5), color=LIGHT, label_color=GREEN, h=Inches(0.7),
            left=Inches(0.6), width=Inches(12.1)):
    band(s, color, left, top, width, h)
    tf = box(s, left + Inches(0.2), Inches(top.inches + 0.10), Inches(width.inches - 0.4), Inches(h.inches - 0.16))
    p = tf.paragraphs[0]
    if label:
        r = p.add_run(); r.text = label + "  "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = label_color; r.font.name = "Calibri"
    r = p.add_run(); r.text = text
    r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"


def flow_node(s, text, left, top, w, h, fill, fontcolor=WHITE, fsize=12, sub=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill; sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fsize); r.font.bold = True; r.font.color.rgb = fontcolor; r.font.name = "Calibri"
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(fsize - 2.5); r2.font.color.rgb = fontcolor; r2.font.name = "Calibri"
    return sh


def arrow(s, x1, y, x2, color=SLATE):
    from pptx.oxml.ns import qn
    cn = s.shapes.add_connector(2, x1, y, x2, y)  # straight
    cn.line.color.rgb = color; cn.line.width = Pt(2.25)
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return cn


def step_card(s, num, head, body, left, top, w, h, accent):
    """A simple numbered step card: accent number chip + heading + plain body."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT
    sh.line.color.rgb = accent; sh.line.width = Pt(1.25); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{num}.  "
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = accent; r.font.name = "Calibri"
    r = p.add_run(); r.text = head
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = body
    r.font.size = Pt(11.5); r.font.color.rgb = SLATE; r.font.name = "Calibri"
    return sh


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.45), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
for i, (txt, sz, col, bold) in enumerate([
    ("RAG Pipeline & Recommendation System", 28, GREEN, True),
    ("Explained simply — step by step, no jargon left behind", 18, RGBColor(0xCB, 0xD5, 0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "Part A: how the AI answers from our document   ·   Part B: how we pick the 6 safety actions"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Two systems explained in plain words. Part A is RAG: how we let the AI answer using our own "
         "safety document instead of its memory. Part B is the recommender: how we choose and order the "
         "6 safety actions we show a user. A worked example for each.")

# ── SLIDE 2 — What this talk covers ────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What This Talk Covers")
bullets(s, [
    ("Part A — RAG: how the AI answers using OUR trusted document, not its memory.", 0, DEEP, True),
    ("Part B — Recommender: how we pick and order the 6 safety actions for a person.", 0, GREEN, True),
    ("A real worked example for each system, so it clicks end to end.", 0, SLATE, False),
    ("A one-line glossary at the end for every key word.", 0, SLATE, False),
], top=Inches(1.8), size=19, gap=18)
callout(s, "The promise:", "both systems told as a simple ordered sequence — what happens first, second, third, all the way to the result.",
        top=Inches(6.4), h=Inches(0.75))
notes(s, "Set expectations: one idea per slide, plain language, two clear sequences.")

# ── SLIDE 3 — PART A divider ───────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.45), W, Inches(0.08))
tf = box(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Part A — The RAG Pipeline"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.7)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Turning a safety PDF into grounded, AI-written advice"
r.font.size = Pt(20); r.font.color.rgb = GREEN; r.font.name = "Calibri"
notes(s, "RAG = Retrieval-Augmented Generation: retrieve the right text, then generate the answer from it.")

# ── SLIDE 4 — The problem RAG solves ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Problem RAG Solves")
bullets(s, [
    ("A plain AI is smart, but it can MAKE THINGS UP — dangerous for safety advice.", 0, RED, True),
    ("RAG's idea: give the AI the right page FIRST, then ask the question.", 0, DEEP, True),
    ("So it answers from OUR document, not from memory.", 1, SLATE, False),
    ("Our document: an official Natural Disaster Safety Guidelines PDF.", 0, SLATE, False),
    ("One chapter per disaster type (Flood, Earthquake, Storm, …).", 1, SLATE, False),
], top=Inches(1.6), size=18, gap=14)
callout(s, "Analogy:", "don't ask a friend from memory — hand them the exact manual chapter and say \"answer using this.\"",
        top=Inches(6.4), h=Inches(0.75), label_color=GREEN)
notes(s, "RAG stands for Retrieval-Augmented Generation. Retrieve the right text, then generate from it. "
         "That is the whole idea in one line.")

# ── SLIDE 5 — RAG flow at a glance (native straight line) ───────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The RAG Flow at a Glance")
y = Inches(2.15); hN = Inches(1.0)
nodes = [
    ("Safety\nPDF", SLATE),
    ("Split into\nchapters", BLUE),
    ("Find the right\nchapter", GREEN),
    ("Build a\nprompt", AMBER),
    ("Groq AI\nwrites answer", PURPLE),
    ("Return\nanswer", DEEP),
]
wN = Inches(1.78); gap = Inches(0.22)
x = Inches(0.55)
positions = []
for i, (txt, col) in enumerate(nodes):
    flow_node(s, txt, x, y, wN, hN, col, fsize=13)
    positions.append((x, col))
    x = Emu(int(x) + int(wN) + int(gap))
for i in range(len(nodes) - 1):
    x1 = Emu(int(positions[i][0]) + int(wN))
    x2 = Emu(int(positions[i + 1][0]))
    arrow(s, x1, Emu(int(y) + int(hN) // 2), x2)
callout(s, "Read it left to right:", "everything before \"find the chapter\" is one-time preparation; everything after is the live answer.",
        top=Inches(4.1), h=Inches(0.8))
callout(s, "In short:", "documents → chapters → find the match → prompt → AI writes → answer. No detours.",
        top=Inches(5.3), h=Inches(0.75), color=LIGHT, label_color=DEEP)
notes(s, "The straight line is the point: a clear sequence with no branching to follow.")

# ── SLIDE 6 — Prepare the knowledge (done once) ────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Step by Step — Prepare the Knowledge (done once)", kicker="Build time")
step_card(s, 1, "Collect", "Start with one trusted source: the official safety PDF.",
          Inches(0.6), Inches(1.55), Inches(3.9), Inches(1.7), SLATE)
step_card(s, 2, "Clean it up", "Pull the text out with PyMuPDF and strip the repeating page headers.",
          Inches(4.7), Inches(1.55), Inches(3.9), Inches(1.7), BLUE)
step_card(s, 3, "Split into chunks", "Cut into one chunk per disaster type → saved as chapters.json (8 keys, ~60 KB).",
          Inches(8.8), Inches(1.55), Inches(3.95), Inches(1.7), GREEN)
callout(s, "Analogy:", "like tabbing the manual so each disaster has its own tab you can flip straight to.",
        top=Inches(3.6), h=Inches(0.7), label_color=GREEN)
callout(s, "At startup:", "the app loads chapters.json into memory once and opens one Groq connection — so every later answer is instant.",
        top=Inches(4.55), h=Inches(0.75), color=LIGHT, label_color=DEEP)
callout(s, "Heads-up (next slide):", "the first design used embeddings + a vector database here. We changed it — and the reason is a good story.",
        top=Inches(5.5), h=Inches(0.75), color=PINK, label_color=RED)
notes(s, "We split by chapter because the document is already organised by disaster — simple and reliable. "
         "chapters.json is committed to the repo so the server needs nothing heavy at boot.")

# ── SLIDE 7 — Why chapters, not a vector database (chart) ───────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Why Chapters, Not a Vector Database", kicker="An honest trade-off")
img(s, "01_why_chapters.png", Inches(0.45), Inches(1.4), Inches(7.2))
bullets(s, [
    ("First design: embeddings + ChromaDB", 0, SLATE, False),
    ("(a vector database) for similarity search.", 1, SLATE, False),
    ("That needs PyTorch — about 2 GB of memory.", 0, RED, True),
    ("Our free server has only 512 MB → it crashed.", 1, SLATE, False),
    ("Fix: ship a 60 KB JSON and look up the", 0, GREEN, True),
    ("chapter by name. It actually runs.", 1, SLATE, False),
    ("Cost: coarser search (whole chapter, not", 0, AMBER, True),
    ("tiny pieces). Benefit: it deploys.", 1, SLATE, False),
], left=Inches(7.8), top=Inches(1.6), width=Inches(5.3), size=14, gap=8)
notes(s, "This is the headline engineering decision and a great teaching moment about real constraints. "
         "The embedding version still exists as legacy code; it just isn't deployed.")

# ── SLIDE 7b — We tested how to split the document (chart) ─────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "We Tested Our Choices", kicker="Why semantic splitting")
img(s, "03_chunking_pick.png", Inches(0.45), Inches(1.4), Inches(7.3))
bullets(s, [
    ("Splitting a document can go wrong:", 0, SLATE, False),
    ("cut in the wrong place and a chunk starts", 1, SLATE, False),
    ("mid-sentence and loses its meaning.", 1, SLATE, False),
    ("So we tried 4 splitting methods and scored", 0, DEEP, True),
    ("them on 30 test questions.", 1, SLATE, False),
    ("Winner: \"Semantic\" — it cuts at natural", 0, GREEN, True),
    ("topic changes, keeping whole sentences", 1, SLATE, False),
    ("together.", 1, SLATE, False),
    ("(This guided the design; the live app sends", 0, GREY, False),
    ("the whole chapter, so it skips splitting.)", 1, GREY, False),
], left=Inches(7.85), top=Inches(1.55), width=Inches(5.25), size=14, gap=7)
notes(s, "Shows we evaluated the choice instead of guessing. Higher score = better mix of relevance and "
         "keeping sentences whole. The production app uses whole chapters, so this informed the legacy "
         "embedding path more than the live one — but it's good evidence of rigor.")

# ── SLIDE 8 — Answer the question (live) ───────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Step by Step — Answer the Question (live)", kicker="Request time")
step_card(s, 4, "Request arrives", "Disaster type, severity, region — e.g. Flood / High / Cairo.",
          Inches(0.55), Inches(1.5), Inches(4.0), Inches(1.5), DEEP)
step_card(s, 5, "Find the chunk", "Look up the chapter by disaster type. Instant and exact.",
          Inches(4.65), Inches(1.5), Inches(4.0), Inches(1.5), GREEN)
step_card(s, 6, "Build the prompt", "Our rules + the chapter + the request, in one message.",
          Inches(8.75), Inches(1.5), Inches(4.0), Inches(1.5), AMBER)
step_card(s, 7, "The AI writes", "Groq llama-3.1-8b-instant, low temperature (focused), JSON output.",
          Inches(0.55), Inches(3.2), Inches(4.0), Inches(1.5), PURPLE)
step_card(s, 8, "Check it", "Must be valid JSON with EXACTLY 6 items — or we use the backup.",
          Inches(4.65), Inches(3.2), Inches(4.0), Inches(1.5), BLUE)
step_card(s, 9, "Return", "Sort the 6 items and send them back to the user.",
          Inches(8.75), Inches(3.2), Inches(4.0), Inches(1.5), SLATE)
callout(s, "The key point:", "step 6 forces a strict format — we're not having an open chat, we're asking for a structured 6-item checklist.",
        top=Inches(5.05), h=Inches(0.75), label_color=GREEN)
callout(s, "Low temperature = ", "the AI stays focused and predictable, not random or creative. Right choice for safety.",
        top=Inches(6.0), h=Inches(0.75), color=LIGHT, label_color=DEEP)
notes(s, "Retrieval here is a simple chapter lookup by disaster type — no math needed because the chapters "
         "are already labelled. The prompt also tells the model to base everything on the provided text and "
         "not invent procedures.")

# ── SLIDE 9 — Follow one real question ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Follow One Real Question")
callout(s, "The situation:", "a subscriber runs a Flood prediction in Cairo, and it comes back High severity.",
        top=Inches(1.45), h=Inches(0.7), color=LIGHT, label_color=DEEP)
bullets(s, [
    ("Request = Flood / High / Cairo.", 0, DEEP, True),
    ("Grab the \"Flood\" chapter from memory (already loaded at startup).", 0, SLATE, False),
    ("Drop the chapter into the prompt together with the request.", 0, SLATE, False),
    ("Groq writes 6 actions: \"evacuate to higher ground\", \"pack a go-bag\", …", 0, GREEN, True),
    ("We confirm it's valid JSON with exactly 6 items, sort them, and show them.", 0, SLATE, False),
    ("The user sees 6 clear, source-grounded flood actions for Cairo.  ✅", 0, DEEP, True),
], top=Inches(2.45), size=18, gap=15)
callout(s, "End to end:", "well under a second — and every line traces back to the official chapter.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)
notes(s, "Walk it slowly so the sequence clicks. Emphasise that the advice is grounded in the document, "
         "not invented.")

# ── SLIDE 10 — PART B divider ──────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.45), W, Inches(0.08))
tf = box(s, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Part B — The Recommendation System"
r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.7), Inches(11.3), Inches(0.7)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Choosing and ordering the 6 safety actions"
r.font.size = Pt(20); r.font.color.rgb = GREEN; r.font.name = "Calibri"
notes(s, "Important framing: this recommender is content/context-based, and it reuses the RAG stack.")

# ── SLIDE 11 — What the recommender does (and isn't) ───────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What the Recommender Does (and Isn't)")
bullets(s, [
    ("It recommends the 6 most relevant SAFETY ACTIONS for a situation.", 0, DEEP, True),
    ("It is content / context-based:", 0, GREEN, True),
    ("we match advice to the disaster + severity + region —", 1, SLATE, False),
    ("not to \"people like you.\"", 1, SLATE, False),
    ("It is NOT collaborative filtering (no \"others also viewed\").", 0, RED, True),
    ("There is no user–item history and no item catalogue to mine.", 1, SLATE, False),
    ("Same situation → same 6 actions for everyone.", 0, SLATE, False),
], top=Inches(1.55), size=18, gap=12)
callout(s, "So what's personal?", "only a friendly NOTE if you've been warned before — the list itself doesn't change. Shown in a moment.",
        top=Inches(6.45), h=Inches(0.75), label_color=GREEN)
notes(s, "Be precise about the method: content-based, not collaborative. The signal is the situation, "
         "not the person's taste or history.")

# ── SLIDE 12 — Recommender flow at a glance (native) ───────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Recommender Flow at a Glance")
y = Inches(2.2); hN = Inches(1.0)
r1 = flow_node(s, "Inputs", Inches(0.6), y, Inches(2.0), hN, DEEP, sub="disaster·severity·region")
r2 = flow_node(s, "Try RAG (Groq)", Inches(2.95), y, Inches(2.3), hN, PURPLE, sub="write 6 actions")
r3 = flow_node(s, "Sort by category", Inches(5.6), y, Inches(2.3), hN, AMBER, sub="fixed order")
r4 = flow_node(s, "Add notice", Inches(8.25), y, Inches(2.1), hN, BLUE, sub="if warned before")
r5 = flow_node(s, "Top 6", Inches(10.7), y, Inches(2.0), hN, GREEN, sub="returned to user")
for ax1, ax2 in [(Inches(2.6), Inches(2.95)), (Inches(5.25), Inches(5.6)),
                 (Inches(7.9), Inches(8.25)), (Inches(10.35), Inches(10.7))]:
    arrow(s, ax1, Emu(int(y) + int(hN) // 2), ax2)
# Fallback branch
fb = flow_node(s, "Backup: DB table", Inches(2.95), Inches(3.7), Inches(2.3), Inches(0.8), RED,
               fsize=11, sub="if the AI fails")
arrow(s, Emu(int(Inches(4.1))), Emu(int(y) + int(hN)), Emu(int(Inches(4.1))), color=RED)
callout(s, "Notice:", "it reuses the RAG stack — there is no second machine-learning model hiding here.",
        top=Inches(5.0), h=Inches(0.7))
callout(s, "In short:", "inputs → try the AI (or the backup) → sort → add a note → return the best 6.",
        top=Inches(5.95), h=Inches(0.7), color=LIGHT, label_color=DEEP)
notes(s, "The recommender IS the RAG pipeline plus a sort, a fallback, and a personalisation note.")

# ── SLIDE 13 — The recommender steps (chart) ───────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Recommendation Steps")
bullets(s, [
    ("1.  Goal: 6 safety actions for one situation.", 0, DEEP, True),
    ("2.  Inputs: disaster type, severity, region, and the user (if logged in).", 0, SLATE, False),
    ("3.  Represent the situation — the key is simply the disaster type label.", 0, SLATE, False),
    ("4.  Item pools ready: the chapters (for the AI) + a static DB table (the backup).", 0, SLATE, False),
    ("5.  A request comes in.", 0, SLATE, False),
    ("6.  Match by the same key on both sides.", 0, SLATE, False),
    ("7.  Score: the AI generates the best-fitting 6 actions from the chapter.", 0, GREEN, True),
    ("8.  Rank: fixed order — evacuation → kit → shelter → medical → contact.", 0, SLATE, False),
    ("9.  Filter & rules: cap at 6; on failure use the DB; add personalisation.", 0, SLATE, False),
    ("10. Return the top 6 (plus an optional note).", 0, DEEP, True),
], left=Inches(0.6), top=Inches(1.45), width=Inches(7.0), size=14.5, gap=8)
img(s, "02_recommendation_contract.png", Inches(7.75), Inches(3.0), Inches(5.4))
callout(s, "Why this order?", "evacuation comes first because it's the most time-critical. The order is a rule, not a model.",
        top=Inches(6.55), h=Inches(0.65), width=Inches(7.0), label_color=GREEN)
notes(s, "Scoring is done by the LLM generating the items; ranking is a deterministic category sort in the "
         "service. No learned ranking model.")

# ── SLIDE 14 — The safety net ──────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Safety Net — Always Return Something")
y = Inches(2.0); hN = Inches(1.1)
flow_node(s, "A request\nfor advice", Inches(0.7), y, Inches(2.3), hN, DEEP, fsize=13)
flow_node(s, "Try RAG (Groq)\nwrites 6 actions", Inches(3.6), y, Inches(2.9), hN, PURPLE, fsize=13)
flow_node(s, "Works?\n6 actions ✅", Inches(7.4), Inches(1.3), Inches(2.6), Inches(0.95), GREEN, fsize=12)
flow_node(s, "Fails? → DB table\npre-written 6", Inches(7.4), Inches(2.85), Inches(2.6), Inches(0.95), RED, fsize=12)
flow_node(s, "User gets\nadvice", Inches(10.55), y, Inches(2.2), hN, DEEP, fsize=13)
arrow(s, Inches(3.0), Emu(int(y) + int(hN) // 2), Inches(3.6))
arrow(s, Inches(6.5), Emu(int(y) + int(hN) // 2), Inches(7.4), color=GREEN)
arrow(s, Inches(6.5), Emu(int(y) + int(hN) // 2), Inches(7.4), color=RED)
arrow(s, Inches(10.0), Inches(1.78), Inches(10.55), color=GREEN)
arrow(s, Inches(10.0), Inches(3.33), Inches(10.55), color=RED)
bullets(s, [
    ("Primary path: the AI writes the 6 actions.", 0, SLATE, False),
    ("If anything fails — no API key, network error, bad JSON, wrong count — we fall back to the DB table.", 0, RED, True),
    ("The user NEVER sees an error caused by the AI being down.", 0, DEEP, True),
], top=Inches(4.4), size=17, gap=12)
callout(s, "This is also \"cold start\":", "when the usual method can't run, we still hand the user a solid pre-written checklist.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)
notes(s, "Every RAG failure degrades to the static recommendations table. The client never gets a 500 "
         "caused by Groq being unavailable.")

# ── SLIDE 15 — Personalisation ─────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Personalisation — a Note, Not a New List")
bullets(s, [
    ("If a logged-in user was ALREADY alerted about this disaster in this region…", 0, DEEP, True),
    ("…we add a line on top:", 0, SLATE, False),
], top=Inches(1.6), size=18, gap=12)
band(s, LIGHT, Inches(1.2), Inches(3.0), Inches(10.9), Inches(1.0))
tf = box(s, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.7)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "“You were previously warned about a Flood risk in Cairo — review carefully.”"
r.font.size = Pt(18); r.font.italic = True; r.font.color.rgb = DEEP; r.font.name = "Calibri"
bullets(s, [
    ("The 6 actions DON'T change — it's a UI hint that adds urgency.", 0, AMBER, True),
    ("Guests (not logged in) never get a note.", 0, SLATE, False),
], top=Inches(4.4), size=18, gap=12)
callout(s, "Why this way?", "cheap, useful, and honest — we don't pretend to re-rank based on history we don't actually have.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)
notes(s, "The notice is computed in the router by checking past alerts for the same disaster + region. "
         "It changes the framing, never the list.")

# ── SLIDE 16 — Follow one real case ────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Follow One Real Case")
callout(s, "The situation:", "a subscriber who was alerted last week about Cairo flooding reopens flood advice.",
        top=Inches(1.45), h=Inches(0.7), color=LIGHT, label_color=DEEP)
bullets(s, [
    ("Inputs = Flood / High / Cairo, and we know who they are.", 0, DEEP, True),
    ("RAG writes 6 flood actions (or the DB does, if Groq is down).", 0, SLATE, False),
    ("Sorted: evacuation first, then kit, shelter, medical, contact.", 0, SLATE, False),
    ("Past alert found → personalisation note added on top.", 0, GREEN, True),
    ("They see a \"you were warned before\" line + 6 ordered, grounded actions.  ✅", 0, DEEP, True),
], top=Inches(2.45), size=18, gap=16)
notes(s, "Same example as Part A, now from the recommender's side, so both halves connect.")

# ── SLIDE 17 — Mini glossary ───────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Mini Glossary")
table(s, [
    ["Word", "In one line"],
    ["Chunk", "A small, meaningful piece of a document — here, one chapter."],
    ["Embedding", "Numbers that capture text meaning (legacy path only)."],
    ["Vector database", "Store that finds similar embeddings fast — ChromaDB (legacy only)."],
    ["Retrieval", "Fetching the right text — here, looking up the chapter."],
    ["Prompt", "The full message we give the AI: rules + chapter + request."],
    ["LLM", "The model that writes the answer — Groq's Llama."],
    ["Ranking", "Putting items in order — our fixed category order."],
    ["Cold start", "What we do when the usual method can't run → DB fallback."],
], Inches(0.55), Inches(1.5), Inches(12.25),
   col_widths=[Inches(3.1), Inches(9.15)], fsize=14)
notes(s, "Keep this as a reference slide — every word used in the talk, defined in one line.")

# ── SLIDE 18 — Key takeaways ───────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("RAG = retrieve the right chapter, then let the AI write FROM it — grounded, not invented.", 0, DEEP, True),
    ("We chose chapter lookup over a vector database so it fits a 512 MB server.", 0, SLATE, False),
    ("The recommender is content-based RAG-as-a-checklist: 6 items, 5 categories, fixed order.", 0, SLATE, False),
    ("A DB fallback guarantees an answer every time; personalisation adds a friendly note.", 0, GREEN, True),
], top=Inches(1.85), size=18, gap=20)
callout(s, "One simple spine:", "get the right text → produce a clean 6-item checklist → never fail silently.",
        top=Inches(6.4), h=Inches(0.75))
notes(s, "Two systems, one idea: ground the answer in our document and deliver a tidy, reliable checklist.")

# ── SLIDE 19 — Thank you / Q&A ─────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(2.55), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(2.7), Inches(11.3), Inches(0.6)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
tf = box(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(3.2)); p = tf.paragraphs[0]
qa = [
    ("Likely questions:", DEEP, True, 16),
    ("Is it real RAG?  →  Yes — retrieve, then generate; retrieval is a chapter lookup.", SLATE, False, 14),
    ("Why no vector database live?  →  PyTorch ~2 GB crashes the 512 MB free server; the JSON is 60 KB.", SLATE, False, 14),
    ("How do you stop the AI making things up?  →  Strict prompt + format check today; a faithfulness test is next.", SLATE, False, 14),
    ("Why content-based, not collaborative?  →  We have no user history; the situation is the signal.", SLATE, False, 14),
]
for i, (txt, col, bold, sz) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
    p.space_after = Pt(8)
    r = p.add_run(); r.text = ("•  " if i else "") + txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0) if i else GREEN
    r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(6.95), Inches(11.3), Inches(0.4)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Read more: docs/RAG_AND_RECSYS_SIMPLE_FLOW.md   ·   safeearth.tech"
r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Keep the deployability and \"how do you stop hallucination\" answers ready — the two most likely "
         "professor questions.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "SafeEarth_RAG_RecSys_Simple_Flow.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides, charts embedded)")
