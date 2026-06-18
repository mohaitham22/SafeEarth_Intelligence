"""
build_rag_pptx.py — generate an editable PowerPoint (with charts) from the SafeEarth
RAG & recommendation-system deep dive.

Run:  py -3.12 docs/build_rag_pptx.py
Out:  docs/SafeEarth_RAG_Presentation.pptx

Native python-pptx (editable text boxes + tables + shapes) + embedded PNG charts rendered
by make_rag_charts.py. Content mirrors docs/RAG_PRESENTATION.md, which is derived from
docs/RAG_AND_RECSYS_DEEP_DIVE.md. Every parameter/number is extracted from this project's
own source (backend/rag/*.py, chunking_report.md, services/recommendation_service.py,
backend/rag/chapters.json), re-verified 2026-06-15.

Honest-evidence caveat carried through the deck: the ONLY measured evaluation is the
4-strategy chunking benchmark (Precision@5 + structural coherence + an LLM-quality proxy);
no retrieval-recall, faithfulness, hallucination, or recsys metric was measured.
The production pipeline is chapter-based Groq (no runtime embeddings / vector store);
the ChromaDB + sentence-transformers pipeline is implemented but NOT deployed.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
import make_rag_charts  # noqa: E402

ASSETS = Path(__file__).resolve().parent / "rag_assets"

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP   = RGBColor(0x0B, 0x3D, 0x2E)
GREEN  = RGBColor(0x16, 0xA3, 0x4A)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xDC, 0x26, 0x26)
SLATE  = RGBColor(0x33, 0x41, 0x55)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x64, 0x74, 0x8B)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x93, 0x33, 0xEA)
PINK   = RGBColor(0xFE, 0xE2, 0xE2)

W, H = Inches(13.333), Inches(7.5)

# Ensure chart PNGs exist before embedding.
if not (ASSETS / "01_chunking_benchmark.png").exists():
    print("Charts not found — rendering them first...")
    make_rag_charts.generate_all()

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def img(s, name, left, top, width):
    return s.shapes.add_picture(str(ASSETS / name), left, top, width=width)


def title_bar(s, text, kicker=None, color=DEEP):
    band(s, color, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = (RGBColor(0xFE, 0xCA, 0xCA) if color == RED else GREEN)
        r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height)
    first = True
    for it in items:
        text, level, color, bold = (it + (0, SLATE, False))[:4] if isinstance(it, tuple) else (it, 0, SLATE, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0)); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None, hcolor=DEEP):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val); r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hcolor
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def callout(s, label, text, top=Inches(6.5), color=LIGHT, label_color=GREEN, h=Inches(0.7),
            left=Inches(0.6), width=Inches(12.1)):
    band(s, color, left, top, width, h)
    tf = box(s, left + Inches(0.2), Inches(top.inches + 0.10), Inches(width.inches - 0.4), Inches(h.inches - 0.16))
    p = tf.paragraphs[0]
    if label:
        r = p.add_run(); r.text = label + "  "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = label_color; r.font.name = "Calibri"
    r = p.add_run(); r.text = text
    r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"


def flow_node(s, text, left, top, w, h, fill, fontcolor=WHITE, fsize=11, sub=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill; sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(fsize); r.font.bold = True; r.font.color.rgb = fontcolor; r.font.name = "Calibri"
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(fsize - 2.5); r2.font.color.rgb = fontcolor; r2.font.name = "Calibri"
    return sh


def arrow(s, x1, y, x2, color=SLATE):
    from pptx.oxml.ns import qn
    cn = s.shapes.add_connector(2, x1, y, x2, y)  # straight
    cn.line.color.rgb = color; cn.line.width = Pt(2.25)
    # Inject a triangle tail-arrowhead into the line XML (not exposed by python-pptx).
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return cn


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.45), Inches(11.3), Inches(1.2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
for i, (txt, sz, col, bold) in enumerate([
    ("RAG Pipeline & Recommendation System", 28, GREEN, True),
    ("Knowledge-grounded safety recommendations — every component, every parameter", 18, RGBColor(0xCB, 0xD5, 0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.1))
for i, txt in enumerate(["Mohamed Haitham  ·  University Presentation",
                         "Chapter-based Groq RAG  ·  llama-3.1-8b-instant  ·  5-category recommendation contract"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Thesis: the recommendation system IS a RAG-as-structured-output pipeline. I'll give "
         "the exact type and parameters of every stage, be explicit about what was measured vs not, "
         "and explain the engineering pivot from ChromaDB+PyTorch to chapter-based Groq that made it "
         "deployable on a 512 MB free tier.")

# ── SLIDE 2 — Roadmap ──────────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Roadmap")
bullets(s, [
    ("What RAG solves here — and the (disaster_type, severity) → 6-item contract", 0, SLATE, False),
    ("The architecture pivot: ChromaDB + PyTorch → chapter-based Groq (the deploy story)", 0, GREEN, True),
    ("Every component, stated by TYPE — chunking, embedding, store, retrieval, generation", 0, SLATE, False),
    ("The one thing measured: the 4-strategy chunking benchmark (Semantic wins)", 0, DEEP, True),
    ("The recommendation system — knowledge-based, not collaborative/content-based", 0, SLATE, False),
    ("Honest evidence: what was NOT measured (retrieval recall, faithfulness, recsys metrics)", 0, RED, True),
    ("Levers that moved quality / latency / cost — and prioritized improvements", 0, SLATE, False),
], top=Inches(1.6), size=18, gap=15)
notes(s, "The spine is the same as the ML talk: real parameters, a self-described engineering "
         "trade-off, and honesty about the evaluation gap.")

# ── SLIDE 3 — What RAG solves ──────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What RAG Solves Here")
bullets(s, [
    ("The predictor answers 'what disaster, how bad?' — zero actionable guidance.", 0, SLATE, False),
    ("RAG answers 'what should the user DO?' — grounded in an official 51-page PDF,", 0, DEEP, True),
    ("not the LLM's parametric memory.", 1, SLATE, False),
    ("Question type is fixed: (disaster_type, severity, region) → a safety checklist.", 0, SLATE, False),
    ("Output is a hard product contract:", 0, GREEN, True),
    ("exactly 6 items · 5 categories {evacuation, kit, shelter, medical, contact}", 1, SLATE, False),
    ("title ≤ 60 chars · body ≤ 280 chars · sorted evacuation → … → contact", 1, SLATE, False),
    ("Surfaced inline in every prediction, in GET /recommendations, and in PDF reports.", 0, SLATE, False),
], top=Inches(1.55), size=18, gap=12)
callout(s, "Framing:", "this is RAG-as-structured-output — retrieval feeds a strict JSON schema, not free-form Q&A.",
        top=Inches(6.45), h=Inches(0.7))

# ── SLIDE 4 — Production architecture (native flow) ────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Production Architecture — Chapter-Based Groq", kicker="Deployed")
# Build-time row
tf = box(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.35)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "BUILD-TIME (run once, committed to repo)"
r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GREY; r.font.name = "Calibri"
y0 = Inches(1.75); hN = Inches(0.85); wN = Inches(2.5)
n1 = flow_node(s, "PDF", Inches(0.6), y0, Inches(1.6), hN, SLATE, sub="51 pages")
n2 = flow_node(s, "extract_chapters.py", Inches(2.55), y0, wN, hN, BLUE, sub="PyMuPDF · regex split")
n3 = flow_node(s, "chapters.json", Inches(5.4), y0, wN, hN, GREEN, sub="8 keys · ~60 KB · committed")
arrow(s, Inches(2.2), Emu(int(y0) + int(hN)//2), Inches(2.55))
arrow(s, Inches(5.05), Emu(int(y0) + int(hN)//2), Inches(5.4))
# Runtime row
tf = box(s, Inches(0.6), Inches(3.1), Inches(12), Inches(0.35)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "REQUEST-TIME"
r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = GREY; r.font.name = "Calibri"
y1 = Inches(3.55)
r1 = flow_node(s, "GET /recommendations", Inches(0.6), y1, Inches(2.3), hN, DEEP, sub="type·severity·region")
r2 = flow_node(s, "_chapters[type]", Inches(3.25), y1, Inches(2.1), hN, GREEN, sub="O(1) dict lookup")
r3 = flow_node(s, "prompt + [:6000]", Inches(5.7), y1, Inches(2.1), hN, BLUE, sub="chapter as context")
r4 = flow_node(s, "llama-3.1-8b-instant", Inches(8.15), y1, Inches(2.3), hN, PURPLE, sub="temp 0.3 · JSON mode")
r5 = flow_node(s, "validate + sort", Inches(10.8), y1, Inches(2.0), hN, AMBER, sub="6 items · 5 cats")
for ax1, ax2 in [(Inches(2.9), Inches(3.25)), (Inches(5.35), Inches(5.7)),
                 (Inches(7.8), Inches(8.15)), (Inches(10.45), Inches(10.8))]:
    arrow(s, ax1, Emu(int(y1) + int(hN)//2), ax2)
# Fallback node
fb = flow_node(s, "DB fallback → recommendations table", Inches(3.25), Inches(4.75), Inches(5.2), Inches(0.7),
               RED, fsize=11, sub="on GroqUnavailableError / any exception")
callout(s, "Cardinal rules:", "chapters + Groq client are startup singletons — never per-request; the router is logic-free; the service owns RAG + fallback.",
        top=Inches(6.45), h=Inches(0.7), label_color=GREEN)
notes(s, "Production has NO runtime embeddings and NO vector store. Retrieval is an O(1) dict lookup "
         "on disaster_type. The only heavy dependency is the groq SDK. chapters.json is committed so "
         "Render needs nothing at boot but json.loads.")

# ── SLIDE 5 — The pivot (memory chart) ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Pivot — Why We Left ChromaDB + PyTorch", kicker="Engineering Trade-off")
img(s, "02_memory_footprint.png", Inches(0.45), Inches(1.35), Inches(7.1))
bullets(s, [
    ("Original design: ChromaDB PersistentClient", 0, SLATE, False),
    ("+ sentence-transformers (all-MiniLM-L6-v2).", 1, SLATE, False),
    ("That pulls PyTorch + CUDA ≈ 2 GB resident.", 0, RED, True),
    ("Render free tier caps at 512 MB → OOM at boot.", 1, SLATE, False),
    ("Fix: pre-extract chapters offline; ship a", 0, GREEN, True),
    ("60 KB JSON; retrieve by key; let Groq read", 1, SLATE, False),
    ("the whole chapter as context.", 1, SLATE, False),
    ("Cost: coarser retrieval (whole chapter, not", 0, AMBER, True),
    ("top-k chunks). Benefit: it actually runs.", 1, SLATE, False),
], left=Inches(7.75), top=Inches(1.5), width=Inches(5.3), size=14, gap=7)
notes(s, "This is the headline engineering decision. The legacy chunking/embedding/vector pipeline is "
         "fully implemented and benchmarked (next slides) — it just isn't deployed. I keep both because "
         "the chunking benchmark is the only measured evaluation and still informs a future restore.")

# ── SLIDE 6 — Corpus & knowledge base ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Knowledge Base & Corpus")
table(s, [
    ["Property", "Value"],
    ["Source", "Natural_Disaster_Safety_Guidelines.pdf · 51 pages · 15 chapters"],
    ["Indexed (production)", "8 EM-DAT types (7 non-EM-DAT chapters excluded)"],
    ["Corpus size", "59,550 chars across 8 keys ≈ 14,887 tokens"],
    ["Per-chapter range", "6,094 (Landslide) – 8,461 (Wildfire) chars"],
    ["Legacy vector store", "167 chunks · 384-dim · ChromaDB (gitignored, not shipped)"],
    ["Update cadence", "STATIC — chapters.json committed; no n8n refresh"],
], Inches(0.55), Inches(1.5), Inches(12.2),
   col_widths=[Inches(2.9), Inches(9.3)], fsize=13.5)
callout(s, "Preprocessing:", "PyMuPDF page text → strip running headers (regex) → split on 'DISASTER TYPE N' → map PDF label to EM-DAT name. No lowercasing / stop-words / normalization beyond header removal.",
        top=Inches(5.5), h=Inches(0.85))
callout(s, "Note:", "the PDF has 15 chapters; only the 8 that match a predictor output class are kept in chapters.json.",
        top=Inches(6.5), h=Inches(0.65), color=LIGHT, label_color=DEEP)

# ── SLIDE 7 — Chunking type + params ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Chunking — Type & Parameters", kicker="Legacy pipeline (benchmarked)")
table(s, [
    ["Aspect", "Exact choice"],
    ["Type", "Semantic — cosine-similarity boundary detection"],
    ["Sentence splitter", "regex (?<=[.!?])\\s+(?=[A-Z]) — no NLTK/spaCy"],
    ["Boundary embedder", "all-MiniLM-L6-v2, batch_size=64, manual cosine"],
    ["similarity_threshold", "0.45 (split when sim < 0.45)"],
    ["min / max sentences", "3 / 15 (hard cap forces split)"],
    ["Overlap", "ZERO — each sentence in exactly one chunk"],
    ["Metadata / chunk", "{strategy, disaster_type, chunk_index}"],
], Inches(0.55), Inches(1.45), Inches(7.7),
   col_widths=[Inches(2.6), Inches(5.1)], fsize=12.5)
bullets(s, [
    ("Why semantic fits this corpus:", 0, DEEP, True),
    ("the PDF shifts topic sharply at section", 1, SLATE, False),
    ("headers (Before / During / After / Medical /", 1, SLATE, False),
    ("Evacuation). Cosine boundaries catch those", 1, SLATE, False),
    ("cleanly → coherence 0.994.", 1, SLATE, False),
    ("Trade-off:", 0, AMBER, True),
    ("3–15 sentences ≈ 100–400 words → precise", 1, SLATE, False),
    ("but a long 'During' section spans 2–3 chunks.", 1, SLATE, False),
], left=Inches(8.45), top=Inches(1.5), width=Inches(4.6), size=13.5, gap=6)
notes(s, "Production doesn't chunk — it sends whole chapters. This slide documents the legacy strategy "
         "that the benchmark selected, and that a future dense-retrieval restore would use.")

# ── SLIDE 8 — Chunking benchmark (chart) ───────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Chunking Benchmark — the One Measured Eval", kicker="Evidence")
img(s, "01_chunking_benchmark.png", Inches(2.5), Inches(1.3), Inches(8.3))
callout(s, "Read it honestly:", "Fixed-Size has the highest raw Precision@5 (0.96) but catastrophic coherence (0.16) — it splits mid-sentence. Semantic wins overall (0.849) on coherence. 'Relevance' = keyword presence, not gold labels.",
        top=Inches(6.4), h=Inches(0.8), label_color=GREEN)
notes(s, "30 queries (2 per topic). Weights 50% relevance / 30% coherence / 20% an LLM-quality proxy "
         "that never calls Groq — it scores context quality (relevance density + imperative-verb "
         "actionability + Jaccard diversity). It is NOT an answer-quality measurement.")

# ── SLIDE 9 — Embedding + vector store ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Embeddings & Vector Store — by Type", kicker="Legacy pipeline")
tf = box(s, Inches(0.6), Inches(1.25), Inches(12), Inches(0.35)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Embeddings"; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
table(s, [
    ["Property", "Value"],
    ["Model", "all-MiniLM-L6-v2"],
    ["Dim", "384"],
    ["Location", "local CPU (no GPU/API)"],
    ["Batch", "64 (chunk) / 32 (add+query)"],
], Inches(0.6), Inches(1.65), Inches(6.0),
   col_widths=[Inches(2.1), Inches(3.9)], fsize=12.5)
tf = box(s, Inches(7.0), Inches(1.25), Inches(6), Inches(0.35)); p = tf.paragraphs[0]
r = p.add_run(); r.text = "Vector store"; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
table(s, [
    ["Property", "Value"],
    ["Store", "ChromaDB PersistentClient"],
    ["Collection", "safety_guidelines"],
    ["Distance", "cosine (hnsw:space)"],
    ["HNSW params", "defaults (M≈16, efC≈100)"],
    ["Vectors", "167"],
], Inches(7.0), Inches(1.65), Inches(5.9),
   col_widths=[Inches(2.4), Inches(3.5)], fsize=12.5)
callout(s, "Why Chroma here:", "embedded in-process, zero-ops, native cosine — fits a free-tier model. At 167 vectors HNSW is effectively exact; FAISS/Qdrant/Pinecone/pgvector add ops or a dependency for no recall gain.",
        top=Inches(5.0), h=Inches(0.85))
callout(s, "Why MiniLM:", "80 MB, CPU-friendly, 384-dim is enough for one English technical doc; trade lower recall vs 768/1536-dim for speed + zero cost.",
        top=Inches(6.0), h=Inches(0.75), color=LIGHT, label_color=DEEP)

# ── SLIDE 10 — Retrieval + reranking ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Retrieval & Reranking — by Type")
table(s, [
    ["Aspect", "Production (deployed)", "Legacy (benchmarked)"],
    ["Retrieval type", "Exact key lookup", "Dense-only (cosine ANN)"],
    ["Mechanism", "_chapters.get(type), O(1)", "col.query(query_embeddings=…)"],
    ["top_k", "n/a (whole chapter)", "5"],
    ["Similarity threshold", "none", "none (always returns 5)"],
    ["Metadata filter", "case-insensitive key", "none (disaster_type stored, unused)"],
    ["Hybrid (BM25+dense)", "no", "no"],
    ["MMR / diversity", "no", "no"],
    ["Reranking", "NONE", "NONE (no cross-encoder)"],
], Inches(0.5), Inches(1.45), Inches(12.3),
   col_widths=[Inches(3.0), Inches(4.6), Inches(4.7)], fsize=12.5)
callout(s, "The gap:", "no reranker in either path; legacy retrieval ignores the disaster_type metadata it stores, and applies no score cutoff — both are on the improvement list (§ later).",
        top=Inches(6.4), h=Inches(0.7), color=PINK, label_color=RED)

# ── SLIDE 11 — Context assembly + truncation (chart) ───────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Context Assembly — and a Real Limitation", kicker="Honest")
img(s, "03_chapter_truncation.png", Inches(0.45), Inches(1.35), Inches(7.4))
bullets(s, [
    ("Production context = chapter_text[:6000].", 0, DEEP, True),
    ("Every one of the 8 chapters is > 6,000 chars,", 0, RED, True),
    ("so all 8 are truncated — tail-cut drops the", 1, SLATE, False),
    ("Emergency Contacts + part of Evacuation.", 1, SLATE, False),
    ("Token budget ≈ 2,060/call; context is 73%.", 0, AMBER, True),
    ("System + user prompt enforce: 6 items,", 0, SLATE, False),
    ("5 categories, char caps, 'do not invent'.", 1, SLATE, False),
    ("No source/citation is returned to the client.", 0, GREY, False),
], left=Inches(8.05), top=Inches(1.5), width=Inches(5.0), size=13.5, gap=6)
notes(s, "Be upfront: the [:6000] tail-cut is a real defect — a priority-ordered section extraction "
         "would keep Contacts/Evacuation. Sort by CATEGORY_ORDER is deterministic in the service; the "
         "LLM only assigns categories.")

# ── SLIDE 12 — Generation + token budget (chart) ───────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Generation — Groq llama-3.1-8b-instant", kicker="LLM")
img(s, "04_token_budget.png", Inches(0.5), Inches(1.4), Inches(7.4))
table(s, [
    ["Param", "Value"],
    ["Model", "llama-3.1-8b-instant"],
    ["Temperature", "0.3"],
    ["response_format", "json_object"],
    ["max_tokens", "unset (default)"],
    ["Streaming", "off (blocking)"],
    ["Client", "startup singleton"],
], Inches(8.1), Inches(1.55), Inches(4.9),
   col_widths=[Inches(2.2), Inches(2.7)], fsize=12.5)
callout(s, "Why Groq:", "LPU latency ~200–500 ms · free tier 6,000 req/day · API-key-only deploy. Trade-off: weaker instruction-following than GPT-4o — JSON mode + an explicit schema example compensate; bad responses fall back to DB.",
        top=Inches(6.4), h=Inches(0.8), label_color=GREEN)

# ── SLIDE 13 — Recommendation system: approach ─────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Recommendation System — Its Type")
bullets(s, [
    ("Type: KNOWLEDGE-BASED (LLM generation from retrieved context).", 0, DEEP, True),
    ("NOT collaborative filtering · NOT content-based · NOT embedding-similarity recs.", 0, RED, True),
    ("There is no user–item matrix and no item catalogue.", 1, SLATE, False),
    ("Signal in = (disaster_type, severity, region) only. Same tuple → same 6 items for every user.", 0, SLATE, False),
    ("Severity conditions the LLM via free-text prompt — not via retrieval filtering.", 0, AMBER, True),
    ("Scoring/ranking: NONE learned — Groq assigns categories; the service sorts by CATEGORY_ORDER.", 0, SLATE, False),
    ("The only quality gate: schema validation (exactly 6, valid category) → else fallback.", 0, GREEN, True),
    ("Cold-start: not applicable — all signal is the disaster type; history never feeds generation.", 0, SLATE, False),
], top=Inches(1.5), size=16.5, gap=12)
callout(s, "The honest label:", "this is RAG-as-structured-output dressed as a recommender — the '6 items / 5 categories' is a product contract imposed on retrieval+generation.",
        top=Inches(6.5), h=Inches(0.7))

# ── SLIDE 14 — Recsys ↔ RAG + forecast dedup (chart) ───────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Recsys ↔ RAG — One Stack, Not Two")
bullets(s, [
    ("The recommender IS the RAG stack —", 0, DEEP, True),
    ("no separate recsys layer.", 1, SLATE, False),
    ("recommendation_service: rag.recommender", 0, SLATE, False),
    ("→ catch GroqUnavailableError → DB fallback.", 1, SLATE, False),
    ("predictor_service: inline per prediction;", 0, SLATE, False),
    ("forecast dedups by severity (chart →).", 1, SLATE, False),
    ("router: thin pass-through; adds a", 0, SLATE, False),
    ("personalisation_notice (Alert × Subscription).", 1, SLATE, False),
    ("Notice does NOT change items — UI hint only.", 0, AMBER, True),
    ("DB fallback: 12 seeded rows; returns < 6", 0, SLATE, False),
    ("rather than fabricate.", 1, SLATE, False),
], left=Inches(0.6), top=Inches(1.5), width=Inches(6.4), size=15, gap=8)
img(s, "05_forecast_dedup.png", Inches(7.2), Inches(1.7), Inches(5.7))
callout(s, "Swap-safety:", "restore dense retrieval and recommendation_service needs zero changes — the swap lives entirely inside rag.recommender.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)

# ── SLIDE 15 — Reliability & failure modes ─────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Reliability & Failure Modes")
table(s, [
    ["Failure", "Detection", "Response"],
    ["GROQ_API_KEY empty", "client is None", "GroqUnavailableError → DB"],
    ["Network / rate-limit", "except around create()", "GroqUnavailableError → DB"],
    ["Malformed JSON", "json.JSONDecodeError", "→ DB fallback"],
    ["Item count ≠ 6", "len check", "→ DB fallback"],
    ["Invalid category", "Pydantic model_validate", "→ DB fallback"],
    ["Unknown disaster_type", "empty chapter lookup", "→ DB fallback"],
    ["DB fallback empty", "len < 6", "return < 6 (no fabrication)"],
    ["Hallucinated procedure", "NONE", "prompt instruction only"],
], Inches(0.5), Inches(1.45), Inches(12.3),
   col_widths=[Inches(3.5), Inches(4.3), Inches(4.5)], fsize=12.5)
callout(s, "Design stance:", "the client never sees a 500 caused by Groq — every RAG failure degrades to the static table. The one uncovered risk is hallucination: there is no output verification.",
        top=Inches(6.4), h=Inches(0.75), color=PINK, label_color=RED)

# ── SLIDE 16 — Evaluation coverage (chart) ─────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Evaluation — What We Can & Cannot Claim", kicker="Honesty", color=RED)
img(s, "06_eval_coverage.png", Inches(2.35), Inches(1.35), Inches(8.6))
callout(s, "The claim we CAN make:", "Semantic chunking beat 3 alternatives on a 30-query structural+keyword benchmark. We have measured NO retrieval recall, NO faithfulness, NO hallucination rate, and NO recsys metric.",
        top=Inches(6.4), h=Inches(0.8), color=PINK, label_color=RED)
notes(s, "Deliver this as a strength: I know exactly which numbers are load-bearing and which are absent. "
         "The next slide and the recommendations turn each red row into a concrete experiment.")

# ── SLIDE 17 — What improved (levers) ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What Improved Quality / Latency / Cost")
table(s, [
    ["Lever", "Effect", "Domain", "Evidence"],
    ["Chapter-based vs ChromaDB+PyTorch", "boots in 512 MB (was OOM)", "infra/cost", "measured"],
    ["…same change", "context 400→6,000 chars", "retrieval qual.", "trade-off (expected)"],
    ["Semantic chunking (vs Fixed/Recursive)", "coherence 0.99 vs 0.06–0.16", "retrieval qual.", "measured (benchmark)"],
    ["Forecast severity-dedup", "30 → ≤4 Groq calls/batch", "latency/cost", "expected ~7.5×"],
    ["response_format=json_object", "parse failures ≈ 0", "generation", "functional"],
    ["query_embeddings= (not query_texts=)", "avoids 79 MB ONNX pull", "latency/cost", "measured (legacy)"],
    ["Startup singletons (chapters+client)", "zero per-request I/O", "latency", "by design"],
    ["DB fallback on any RAG failure", "zero Groq-caused 500s", "reliability", "tested"],
], Inches(0.5), Inches(1.4), Inches(12.3),
   col_widths=[Inches(4.6), Inches(3.1), Inches(2.0), Inches(2.6)], fsize=11.5)
callout(s, "The verdict:", "the biggest lever was an availability decision (deploy at all), not a quality metric. The only MEASURED quality win is chunking coherence; the rest are functional or expected.",
        top=Inches(6.45), h=Inches(0.75), label_color=GREEN)

# ── SLIDE 18 — Recommendations ─────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Recommendations — Highest Impact First")
bullets(s, [
    ("P1 · Add a faithfulness eval (RAGAS or chapter-similarity check) — verify items are grounded, not invented. The single biggest trust gap.", 0, RED, True),
    ("P1 · Restore dense retrieval via ONNX MiniLM (no PyTorch): context 6,000 → ~800 chars, ~75% fewer Groq tokens, more targeted grounding.", 0, GREEN, True),
    ("P1 · Build a labeled retrieval set → recall@k, MRR, NDCG; tune similarity_threshold (0.35–0.55).", 0, SLATE, False),
    ("P2 · Cross-encoder reranker (ms-marco-MiniLM, ~70 MB CPU) over top-k — typically +5–15% NDCG@3.", 0, SLATE, False),
    ("P2 · Chunk overlap ablation (0/1/2 sentences) + recsys quality eval (LLM-judge over 32 type×severity buckets).", 0, SLATE, False),
    ("P3 · Severity-conditioned section filter (During/Evacuation for High/Critical); priority-ordered truncation to stop dropping Emergency Contacts.", 0, SLATE, False),
    ("P3 · Set max_tokens=512; add region-specific corpus for true geo-conditioned retrieval.", 0, GREY, False),
], top=Inches(1.45), size=15, gap=12)
notes(s, "Ordered by expected payoff. P1 faithfulness + dense-retrieval restore are the two that move "
         "both trust and cost; the labeled retrieval set is the prerequisite for claiming any retrieval gain.")

# ── SLIDE 19 — Takeaways ───────────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("Production RAG is chapter-based Groq: O(1) key lookup → whole-chapter context → llama-3.1-8b-instant, no runtime embeddings or vector store.", 0, DEEP, True),
    ("The defining decision was deployability: PyTorch's ~2 GB OOMs a 512 MB tier, so we traded chunk-precision for a 60 KB JSON that runs.", 0, SLATE, False),
    ("The recommender is knowledge-based RAG-as-structured-output — a 6-item / 5-category contract, not collaborative or content-based.", 0, SLATE, False),
    ("Every RAG failure degrades to a static DB table — zero Groq-caused 500s — but there is no hallucination guardrail beyond the prompt.", 0, SLATE, False),
    ("Only the chunking benchmark is measured. Faithfulness, retrieval recall, and recsys metrics are honest, named gaps with concrete next experiments.", 0, RED, True),
], top=Inches(1.7), size=17, gap=18)
notes(s, "End on the same note as the ML talk: precise component types, one real engineering trade-off, "
         "and explicit honesty about the evaluation surface.")

# ── SLIDE 20 — Thank you / Q&A ─────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(2.55), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(2.7), Inches(11.3), Inches(0.6)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
tf = box(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(3.2)); p = tf.paragraphs[0]
qa = [
    ("Anticipated questions:", DEEP, True, 16),
    ("Is it 'real' RAG?  →  Yes — retrieve-then-generate; production retrieval is chapter-level key lookup, not vector search.", SLATE, False, 14),
    ("Why drop the vector DB?  →  sentence-transformers+PyTorch ≈ 2 GB OOMs Render's 512 MB; chapters.json is 60 KB.", SLATE, False, 14),
    ("How do you stop hallucination?  →  Today only the prompt + schema validation; faithfulness eval is P1 (not yet measured).", SLATE, False, 14),
    ("Best single improvement?  →  ONNX MiniLM dense retrieval — ~75% fewer tokens AND tighter grounding.", SLATE, False, 14),
    ("Why Semantic chunking?  →  Coherence 0.994 vs 0.06–0.16; it respects the PDF's section boundaries.", SLATE, False, 14),
]
for i, (txt, col, bold, sz) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT
    p.space_after = Pt(7)
    r = p.add_run(); r.text = ("•  " if i else "") + txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0) if i else GREEN
    r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(6.95), Inches(11.3), Inches(0.4)); p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Full write-up: docs/RAG_AND_RECSYS_DEEP_DIVE.md   ·   KB: 51-page safety PDF   ·   safeearth.tech"
r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Keep the deployability + faithfulness answers ready — the two most likely probes. Owning the "
         "evaluation gap is the win.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "SafeEarth_RAG_Presentation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides, charts embedded)")
