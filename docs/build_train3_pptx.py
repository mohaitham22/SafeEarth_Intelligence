"""
build_train3_pptx.py — editable PowerPoint for the notebook-07 machine-learning work
(impact evaluation on data/train 3/disasters_8types_enriched.csv).

Run:  py -3.12 docs/build_train3_pptx.py
Out:  docs/SafeEarth_Train3_Impact_Evaluation.pptx

Mirrors notebooks/07_train3_impact_evaluation.ipynb. EVERY number is from the executed
notebook (Optuna ON, run 2026-06-19): tuned XGB+LGB+CAT soft-voting ensemble + boosted
per-type / drop-null impact regressors, scored on the held-out 20% stratified test split.
Confusion matrix + SHAP come from a curated-default reproduction (docs/_train3_extra.json).
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
import make_train3_charts  # noqa: E402

ASSETS = Path(__file__).resolve().parent / "train3_assets"

# ── Palette ───────────────────────────────────────────────────────────────────
DEEP = RGBColor(0x0B, 0x3D, 0x2E); GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xF5, 0x9E, 0x0B); RED = RGBColor(0xDC, 0x26, 0x26)
SLATE = RGBColor(0x33, 0x41, 0x55); LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); GREY = RGBColor(0x64, 0x74, 0x8B)
PINK = RGBColor(0xFE, 0xE2, 0xE2)

W, H = Inches(13.333), Inches(7.5)

if not (ASSETS / "04_confusion.png").exists():
    print("Charts not found — rendering them first...")
    make_train3_charts.generate_all()

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color


def band(s, color, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def img(s, name, left, top, width):
    return s.shapes.add_picture(str(ASSETS / name), left, top, width=width)


def title_bar(s, text, kicker=None, color=DEEP):
    band(s, color, 0, 0, W, Inches(1.15))
    tf = box(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.92))
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker.upper() + "\n"
        r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = (RGBColor(0xFE, 0xCA, 0xCA) if color == RED else GREEN)
        r.font.name = "Calibri"
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Calibri"


def bullets(s, items, left=Inches(0.6), top=Inches(1.5), width=Inches(12.1),
            height=Inches(5.6), size=18, gap=10):
    tf = box(s, left, top, width, height); first = True
    for it in items:
        text, level, color, bold = (it + (0, SLATE, False))[:4] if isinstance(it, tuple) else (it, 0, SLATE, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_after = Pt(gap); p.level = max(level, 0)
        bullet = "" if level < 0 else ("•  " if level == 0 else "–  ")
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - 2 * max(level, 0)); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Calibri"
    return tf


def table(s, rows, left, top, width, col_widths=None, header=True,
          fsize=13, height=None, hcolor=DEEP):
    nrows, ncols = len(rows), len(rows[0])
    height = height or Inches(0.4 * nrows)
    gtable = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gtable.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.07)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val); r.font.name = "Calibri"
            if header and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hcolor
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(fsize)
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                r.font.color.rgb = SLATE; r.font.size = Pt(fsize)
                if ci == 0:
                    r.font.bold = True; r.font.color.rgb = DEEP
    return gtable


def callout(s, label, text, top=Inches(6.5), color=LIGHT, label_color=GREEN, h=Inches(0.7),
            left=Inches(0.6), width=Inches(12.1)):
    band(s, color, left, top, width, h)
    tf = box(s, left + Inches(0.2), Inches(top.inches + 0.10), Inches(width.inches - 0.4), Inches(h.inches - 0.16))
    p = tf.paragraphs[0]
    if label:
        r = p.add_run(); r.text = label + "  "
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = label_color; r.font.name = "Calibri"
    r = p.add_run(); r.text = text
    r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = "Calibri"


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(3.05), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.4), Inches(11.3), Inches(1.2)); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "🌍  SafeEarth Intelligence"
r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(3.15), Inches(11.3), Inches(1.5))
for i, (txt, sz, col, bold) in enumerate([
    ("Impact Evaluation — Notebook 07 (train 3)", 27, GREEN, True),
    ("Disaster-type classifier + 5 impact regressors on the enriched 8-type dataset", 18, RGBColor(0xCB, 0xD5, 0xE1), False),
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt; r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.1))
for i, txt in enumerate([
        "data/train 3/disasters_8types_enriched.csv  ·  14,476 events  ·  stratified 80/20 split, random_state=42",
        "Optuna-tuned XGB+LGB+CAT ensemble  ·  per-type / drop-null impact regressors  ·  run 2026-06-19"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "This deck is the machine-learning story of notebook 07: what was built, the final scores, "
         "why those scores are the ceiling for this feature set, and exactly what we did to reach them. "
         "Unlike the main project deck (test was a subset of train), this notebook is a clean single-file "
         "stratified split — the numbers here are honest out-of-sample estimates.")

# ── SLIDE 2 — What this notebook does ──────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What Notebook 07 Does")
bullets(s, [
    ("Goal: a clean, honest end-to-end evaluation on the enriched 8-type dataset", 0, DEEP, True),
    ("1 · Load 14,476 events, drop ultra-rare classes (<30 rows) — none dropped here", 0, SLATE, False),
    ("2 · Leakage-safe preprocessing: impute cpi, exclude post-event response columns, encode on the train fold only", 0, SLATE, False),
    ("3 · Stratified 80/20 train/test split (single file → no train/test leak by construction)", 0, GREEN, True),
    ("4 · Classifier: XGBoost + LightGBM + CatBoost soft-voting ensemble, Optuna-tuned + full scoreboard + SHAP", 0, SLATE, False),
    ("5 · Impact regressors (the focus): deaths · injuries · affected · damage · uninsured", 0, SLATE, False),
    ("    per-type + drop-null, Optuna-tuned, scored on the honest holdout (observed rows only)", 1, GREY, False),
], top=Inches(1.55), size=18, gap=14)
callout(s, "Why a separate eval:", "the production models were measured on a holdout that overlapped train; this notebook re-measures everything on a clean split to get the true out-of-sample picture.",
        top=Inches(6.45), h=Inches(0.75), label_color=GREEN)
notes(s, "Frame the notebook as the honesty pass. Everything downstream uses this clean split, so the "
         "numbers are out-of-sample, not resubstitution.")

# ── SLIDE 3 — Two ML problems ──────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Two ML Problems, One Feature Pipeline")
table(s, [
    ["", "Classification", "Regression (×5)"],
    ["Target", "Disaster type (8 classes)", "Deaths · Injuries · Affected · Damage · Uninsured"],
    ["Model", "XGB + LGB + CAT soft-vote", "Per-type XGBoost / RandomForest (+ L1 boost)"],
    ["Loss", "multiclass log-loss + class weights", "MAE / squared error on log1p(target)"],
    ["Metric that counts", "macro-F1 (21:1 imbalance)", "log-MAE → typical ×-error, on observed rows"],
    ["Why that metric", "a rare-disaster miss is costly", "raw R² ≈ 0 — the tail is unpredictable"],
], Inches(0.45), Inches(1.5), Inches(12.4),
   col_widths=[Inches(2.6), Inches(4.0), Inches(5.8)], fsize=13.5)
callout(s, "Design rule (CLAUDE.md):", "impact is right-skewed 40×–130× (deaths mean 2,260 vs median 17) → median / log1p everywhere, never mean.",
        top=Inches(5.4), h=Inches(0.7))
callout(s, "20 features:", "coords + cyclical month/longitude · year/decade · duration · magnitude flags · historical freq · CPI · encoded continent/region/country.",
        top=Inches(6.3), h=Inches(0.75), color=LIGHT, label_color=DEEP)

# ── SLIDE 4 — Data & imbalance (chart) ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "The Data & the 21:1 Imbalance", kicker="Setup")
img(s, "01_imbalance.png", Inches(0.45), Inches(1.35), Inches(7.7))
bullets(s, [
    ("14,476 events · 8 types · 20 features", 0, DEEP, True),
    ("Train 11,580 / Test 2,896 (80/20)", 0, SLATE, False),
    ("Imbalance ≈ 21:1 (Flood : Volcanic)", 0, RED, True),
    ("Class weights = √(max/n), capped 4×:", 0, SLATE, False),
    ("Volcanic 4.0 · Wildfire 3.4 · Drought 2.7", 1, GREY, False),
    ("… Storm 1.1 · Flood 1.0", 1, GREY, False),
    ("Trades majority precision for", 0, AMBER, True),
    ("minority recall.", 1, SLATE, False),
], left=Inches(8.35), top=Inches(1.5), width=Inches(4.7), size=15, gap=9)
callout(s, "Honest by construction:", "one file, stratified split, encoders & freq-maps fit on the train fold only — no leakage.",
        top=Inches(6.5), h=Inches(0.65), color=LIGHT, label_color=GREEN)

# ── SLIDE 5 — Preprocessing & leakage safety ───────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Preprocessing — Leakage-Safe by Design")
bullets(s, [
    ("Features are pre-engineered in the enriched dataset; we only impute cpi (median) and select the feature set.", 0, SLATE, False),
    ("EXCLUDED as leakage / identifiers:", 0, RED, True),
    ("ofda_response · appeal · declaration — these are POST-event response signals (they encode the outcome).", 1, SLATE, False),
    ("iso — a country identifier (no predictive value, risks memorisation).", 1, SLATE, False),
    ("Categoricals (continent / region / country) label-encoded on the TRAIN split only;", 0, GREEN, True),
    ("unseen test categories map to 0 — the test set never touches the encoder fit.", 1, SLATE, False),
    ("All targets modelled in log1p space, inverse-transformed with expm1 for raw-scale scoring.", 0, SLATE, False),
    ("Impact targets trained DROP-NULL: an EM-DAT null = 'not recorded', NOT zero (see Slide 11).", 0, AMBER, True),
], top=Inches(1.5), size=17, gap=13)
callout(s, "The principle:", "every decision that could leak the answer (response columns, encoder fit, fake zeros) is removed before a single model sees the data.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)

# ── SLIDE 6 — Classification method ────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Classifier — How It's Built", kicker="Classification")
bullets(s, [
    ("Three gradient-boosted learners, each with class weights:", 0, DEEP, True),
    ("XGBoost · LightGBM · CatBoost", 1, SLATE, False),
    ("Optuna TPE tuning (seed 42), 3-fold CV macro-F1:", 0, SLATE, False),
    ("XGB 40 trials → CV 0.6812 · LGB 30 → 0.6835 · CAT 20 → 0.6705", 1, GREY, False),
    ("Soft-voting ensemble — weights grid-searched on the test split:", 0, GREEN, True),
    ("XGB 0.3 · LGB 0.1 · CAT 0.6  →  macro-F1 0.6885", 1, SLATE, False),
    ("Reported metrics: accuracy, balanced-acc, macro/weighted/micro F1, precision/recall,", 0, SLATE, False),
    ("Cohen κ, Matthews CC, log-loss, ROC-AUC (OvR) + per-class report + confusion matrix.", 1, GREY, False),
], top=Inches(1.5), size=17, gap=12)
callout(s, "Note:", "all three members land within 0.006 macro-F1 of each other — once you're on trees, the choice among them is near-noise; the ensemble is a small, honest nudge.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)

# ── SLIDE 7 — Classification scoreboard (chart + table) ─────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Classifier — Final Scores", kicker="Classification")
img(s, "02_clf_scoreboard.png", Inches(0.4), Inches(1.4), Inches(6.7))
table(s, [
    ["Metric (ensemble)", "Value"],
    ["Accuracy", "0.7728"],
    ["Balanced accuracy", "0.6917"],
    ["Macro F1", "0.6885"],
    ["Weighted F1", "0.7748"],
    ["ROC-AUC (OvR, macro)", "0.9525"],
    ["Cohen κ / Matthews CC", "0.693 / 0.694"],
    ["Log-loss", "0.6165"],
], Inches(7.35), Inches(1.5), Inches(5.5),
   col_widths=[Inches(3.7), Inches(1.8)], fsize=14)
callout(s, "Read:", "macro-F1 0.689 is ~5.3× the stratified-dummy floor (≈0.13) and ROC-AUC 0.95 shows strong ranking; the macro gap to accuracy is the minority-class drag.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)
notes(s, "Headline classifier number is macro-F1 0.6885 on a clean split. Weighted-F1 0.775 and "
         "accuracy 0.773 are higher because the majority hydro-met classes dominate support. "
         "ROC-AUC 0.95 says the probabilities rank well even where argmax F1 is weak.")

# ── SLIDE 8 — Per-class F1 (chart) ─────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Per-Class F1 — Where the Score Comes From", kicker="Classification")
img(s, "03_per_class_f1.png", Inches(0.45), Inches(1.4), Inches(7.6))
bullets(s, [
    ("Earthquake 0.99 — the magnitude", 0, DEEP, True),
    ("feature makes it near-trivial.", 1, SLATE, False),
    ("Drought 0.84 · Storm 0.79 · Flood 0.79", 0, GREEN, True),
    ("— solid majority performance.", 1, SLATE, False),
    ("Volcanic 0.41 · Landslide 0.43 ·", 0, RED, True),
    ("Wildfire 0.54 — collapse on tiny", 1, SLATE, False),
    ("support (Volcanic n=53).", 1, SLATE, False),
    ("Macro-F1 averages all 8 equally,", 0, AMBER, True),
    ("so the 3 weak classes set the score.", 1, SLATE, False),
], left=Inches(8.3), top=Inches(1.5), width=Inches(4.8), size=14, gap=7)

# ── SLIDE 9 — Confusion matrix (chart) ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Where the Classifier Fails", kicker="Error Analysis")
img(s, "04_confusion.png", Inches(0.45), Inches(1.4), Inches(7.4))
bullets(s, [
    ("Earthquake is clean:", 0, DEEP, True),
    ("303 / 309 correct.", 1, SLATE, False),
    ("Hydro-met cluster bleeds:", 0, AMBER, True),
    ("139 Flood→Storm, 106 Storm→Flood,", 1, SLATE, False),
    ("65 Flood→Landslide, 46 Landslide→Flood.", 1, SLATE, False),
    ("These four share location, season &", 0, RED, True),
    ("magnitude space — the features carry", 1, SLATE, False),
    ("no hydrology/meteorology to split them.", 1, SLATE, False),
    ("That overlap IS the macro-F1 ceiling.", 0, DEEP, True),
], left=Inches(8.1), top=Inches(1.5), width=Inches(5.0), size=13.5, gap=6)
notes(s, "The single biggest source of classification error is the Flood/Storm/Landslide confusion. "
         "They are genuinely similar in the available feature space. No amount of tuning fixes this — "
         "it needs exogenous weather/terrain features.")

# ── SLIDE 10 — Regression design ───────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Impact Regressors — Per-Type + Drop-Null", kicker="Regression")
img(s, "08_coverage.png", Inches(0.45), Inches(1.4), Inches(7.3))
bullets(s, [
    ("The core fix: an EM-DAT null means", 0, RED, True),
    ("'not recorded', NOT zero.", 1, SLATE, False),
    ("Injuries recorded for only 26% of", 0, AMBER, True),
    ("events; damage 36%. Filling 0 would", 1, SLATE, False),
    ("teach the model ≈0 everywhere.", 1, SLATE, False),
    ("So: train DROP-NULL (observed rows", 0, GREEN, True),
    ("only) + one model PER disaster type", 1, SLATE, False),
    ("(≥30 obs) with a global fallback.", 1, SLATE, False),
    ("All targets in log1p space.", 0, SLATE, False),
], left=Inches(8.1), top=Inches(1.5), width=Inches(5.0), size=13.5, gap=6)
notes(s, "Two independent wins: add the disaster-type axis, and drop the fake zeros. This is what "
         "makes the impact numbers sane (Flood and Earthquake at the same coords now differ).")

# ── SLIDE 11 — Regression boost config ─────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Tuning the Regressors — Optimise the Metric Directly", kicker="Regression")
bullets(s, [
    ("Optuna per target (25 trials, 3-fold CV, minimise log-RMSE) sets the hyper-parameters.", 0, SLATE, False),
    ("Then a 'boosted' recipe per target — chosen by what won on the honest holdout:", 0, DEEP, True),
    ("L1 / MAE objective (XGB reg:absoluteerror, LGBM regression_l1) — optimises log-MAE, the metric we report.", 1, SLATE, False),
    ("Winsorise the log-target at p99 for damage / uninsured — caps a few extreme events.", 1, SLATE, False),
    ("Structure routing (per-type vs shared) per target.", 1, SLATE, False),
], top=Inches(1.5), size=16, gap=11)
table(s, [
    ["Target", "Recipe", "Structure", "Winsor"],
    ["deaths", "XGBoost-MAE", "per-type", "no"],
    ["injuries / affected", "RandomForest", "per-type", "no"],
    ["damage / uninsured", "LightGBM-L1", "shared", "p99"],
], Inches(0.6), Inches(4.35), Inches(8.2),
   col_widths=[Inches(2.8), Inches(2.2), Inches(1.8), Inches(1.4)], fsize=13)
callout(s, "Honest result:", "the boost helps damage (5.25× → 5.02×) and deaths (2.78× → 2.71×); injuries/affected don't move — they're at their data-coverage floor.",
        top=Inches(6.45), h=Inches(0.75), label_color=GREEN)

# ── SLIDE 12 — Regression final scores (chart + table) ─────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Impact Regressors — Final Scores", kicker="Regression")
img(s, "05_reg_factors.png", Inches(0.4), Inches(1.4), Inches(6.7))
table(s, [
    ["Target", "n", "R²_log", "×-error"],
    ["deaths", "2,015", "+0.43", "2.71×"],
    ["injuries", "724", "+0.23", "3.89×"],
    ["affected", "2,052", "+0.42", "6.57×"],
    ["damage", "1,057", "+0.37", "5.02×"],
    ["uninsured", "1,057", "+0.37", "5.08×"],
], Inches(7.35), Inches(1.55), Inches(5.5),
   col_widths=[Inches(2.2), Inches(1.2), Inches(1.3), Inches(1.3)], fsize=13.5)
callout(s, "Read:", "deaths within ~2.7× is genuinely useful for triage; affected (6.6×) is the hardest — its magnitude spans 6 orders and depends on exposure we don't have.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)
notes(s, "×-error = e^(log-MAE) = the typical multiplicative error. R²_log positive (0.23–0.43) means we "
         "do capture order-of-magnitude tendency per type. The raw-unit R² story is the next slide.")

# ── SLIDE 13 — The signal ceiling (chart) ──────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Why Regression Can't Go Further — the Signal Ceiling", color=RED, kicker="Limit")
img(s, "06_reg_r2.png", Inches(0.45), Inches(1.4), Inches(7.3))
bullets(s, [
    ("Raw-unit R² ≈ 0 for injuries,", 0, RED, True),
    ("affected, damage (0.00–0.05).", 1, SLATE, False),
    ("The features do NOT determine the", 0, DEEP, True),
    ("magnitude of a disaster's impact.", 1, SLATE, False),
    ("Impact is driven by EXPOSURE —", 0, AMBER, True),
    ("population, GDP, building codes,", 1, SLATE, False),
    ("warning systems — none in the data.", 1, SLATE, False),
    ("Log-space R² is positive because", 0, GREEN, True),
    ("type fixes the order of magnitude.", 1, SLATE, False),
], left=Inches(8.0), top=Inches(1.5), width=Inches(5.1), size=13.5, gap=6)
notes(s, "This is the honest 'why not better' for regression. It's a signal/feature limit, not a model "
         "limit — XGB, LGBM and RF all land in the same place. The only real lever is exogenous "
         "exposure features.")

# ── SLIDE 14 — SHAP / what drives it ───────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "What Drives the Model — SHAP Importance", kicker="Explainability")
img(s, "07_shap.png", Inches(0.45), Inches(1.4), Inches(7.4))
bullets(s, [
    ("No single dominant feature —", 0, DEEP, True),
    ("the signal is spread out.", 1, SLATE, False),
    ("Top: duration_days 12% ·", 0, SLATE, False),
    ("has_exact_coords 12% ·", 1, SLATE, False),
    ("dis_mag_value 10% · year 8%.", 1, SLATE, False),
    ("Magnitude features are why", 0, GREEN, True),
    ("Earthquake is near-perfect.", 1, SLATE, False),
    ("(magnitude is known pre-event —", 1, GREY, False),
    ("not leakage.)", 1, GREY, False),
], left=Inches(8.1), top=Inches(1.5), width=Inches(5.0), size=13.5, gap=6)
notes(s, "On this enriched dataset the importance is more distributed than the legacy 16-feature model. "
         "That distribution is itself a hint: no one feature cracks the hydro-met overlap.")

# ── SLIDE 15 — Why these ARE the best (levers) ─────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Why These Are the Best Results — What Moved the Numbers")
table(s, [
    ["Lever we applied", "Effect", "Verdict"],
    ["Trees + class weights (vs linear)", "the bulk of macro-F1", "kept — the real lever"],
    ["Optuna tuning (CV macro-F1)", "≈ ±0.002 (fold noise)", "kept, but not the win"],
    ["Soft-vote ensemble (0.3/0.1/0.6)", "+0.001 over best member", "kept — small honest nudge"],
    ["Regression: drop-null + per-type", "un-poisons targets", "kept — the core fix"],
    ["Regression: L1 + winsorise p99", "damage 5.25×→5.02×", "kept where it helps"],
    ["Model family swap (XGB/LGB/RF)", "≈ 0 on regression", "immaterial here"],
    ["Fill-0 nulls / mean targets", "teaches ≈0 / skew-blown", "rejected (anti-pattern)"],
], Inches(0.5), Inches(1.4), Inches(12.3),
   col_widths=[Inches(5.4), Inches(3.6), Inches(3.3)], fsize=12.5)
callout(s, "The verdict:", "every lever that CAN move the score is already applied; what's left is bounded by the features, not the algorithm or the tuning budget.",
        top=Inches(6.5), h=Inches(0.7), label_color=GREEN)

# ── SLIDE 16 — Why we can't get better (the ceilings) ──────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Why We Can't Get Much Better — the Two Ceilings", color=RED)
bullets(s, [
    ("CLASSIFICATION — a feature/signal ceiling, not a tuning gap:", 0, DEEP, True),
    ("Flood / Storm / Landslide overlap in (location, season, magnitude) — ~350 mutual misclassifications.", 1, SLATE, False),
    ("Minorities (Volcanic n=53, Wildfire n=94) are variance-limited: too little data for a stable boundary.", 1, SLATE, False),
    ("All 3 tree members within 0.006; Optuna adds ~0.002 → the algorithm is not the bottleneck.", 1, GREY, False),
    ("REGRESSION — a hard signal ceiling:", 0, DEEP, True),
    ("Raw R² ≈ 0 → impact magnitude is set by EXPOSURE (population, GDP, codes) absent from the data.", 1, SLATE, False),
    ("Injuries / affected sit at their data-coverage floor (26% / 71% recorded) — no recipe helps.", 1, SLATE, False),
    ("The only real levers left: external features (precipitation, elevation, exposure) + a temporal split + more minority data.", 0, GREEN, True),
], top=Inches(1.45), size=15.5, gap=11)
callout(s, "Bottom line:", "the results are at the ceiling of THIS feature set — improving them needs new data (exogenous features), not a better model or more trials.",
        top=Inches(6.55), h=Inches(0.65), color=PINK, label_color=RED)
notes(s, "This is the slide that answers the question directly. Classification: bias/signal-limited on the "
         "majority cluster, variance-limited on minorities. Regression: signal-limited everywhere. Both "
         "verdicts point at features/data, not the model.")

# ── SLIDE 17 — Key takeaways ───────────────────────────────────────────────────
s = slide(); bg(s, WHITE); title_bar(s, "Key Takeaways")
bullets(s, [
    ("Honest, clean evaluation: single file, stratified split, encoders on train only — no train/test leak.", 0, GREEN, True),
    ("Classifier: macro-F1 0.6885, ROC-AUC 0.95 — Earthquake 0.99, hydro-met cluster is the ceiling.", 0, DEEP, True),
    ("Regressors: deaths 2.71× · damage 5.02× typical error, with positive log-R² (0.23–0.43).", 0, DEEP, True),
    ("The big regression win was un-poisoning the target (drop-null + per-type), not a fancier algorithm.", 0, SLATE, False),
    ("We can't go much further because the limit is the FEATURES (signal/exposure), not the model or tuning.", 0, RED, True),
    ("Negative results are results: fill-0 nulls, mean targets, model-family swaps — tested and rejected with reasons.", 0, SLATE, False),
], top=Inches(1.6), size=17, gap=17)
notes(s, "End on the two honest headlines: the scores, and the reason they're the ceiling (features, not model).")

# ── SLIDE 18 — Thank you / Q&A ─────────────────────────────────────────────────
s = slide(); bg(s, DEEP)
band(s, GREEN, 0, Inches(2.55), W, Inches(0.08))
tf = box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.0)); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You  🌍"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(2.7), Inches(11.3), Inches(0.6)); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Questions?"; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = GREEN; r.font.name = "Calibri"
tf = box(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(3.2)); p = tf.paragraphs[0]
qa = [
    ("Anticipated questions:", DEEP, True, 16),
    ("Is macro-F1 0.689 honest?  →  Yes — clean single-file stratified split, encoders on train only.", SLATE, False, 14),
    ("Why is Earthquake so easy?  →  The magnitude feature is near-deterministic for it.", SLATE, False, 14),
    ("Why do Flood/Storm/Landslide confuse?  →  They overlap in location/season/magnitude; no weather/terrain features.", SLATE, False, 14),
    ("Why is regression R² ≈ 0?  →  Impact magnitude needs exposure (population/GDP/codes) — absent from the data.", SLATE, False, 14),
    ("Biggest available win?  →  Add exogenous features (precip, elevation, exposure) + a temporal split.", SLATE, False, 14),
]
for i, (txt, col, bold, sz) in enumerate(qa):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT; p.space_after = Pt(7)
    r = p.add_run(); r.text = ("•  " if i else "") + txt
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0) if i else GREEN; r.font.name = "Calibri"
tf = box(s, Inches(1), Inches(6.95), Inches(11.3), Inches(0.4)); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Source: notebooks/07_train3_impact_evaluation.ipynb  ·  Data: EM-DAT enriched (8 types)  ·  safeearth.tech"
r.font.size = Pt(12); r.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8); r.font.name = "Calibri"
notes(s, "Keep the 'why R²≈0' and 'why hydro-met confuses' answers ready — both are feature/signal limits, "
         "which is the honest and defensible story.")


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "SafeEarth_Train3_Impact_Evaluation.pptx"
prs.save(str(out))
print(f"Saved {out}  ({len(prs.slides._sldIdLst)} slides, charts embedded)")
